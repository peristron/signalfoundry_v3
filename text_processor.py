import csv
import html
import io
import re
import string
from collections import Counter
from datetime import datetime
from itertools import pairwise
from typing import Any, Dict, Iterable, List, Optional, Tuple

import pandas as pd

try:
    from wordcloud import STOPWORDS as WORDCLOUD_STOPWORDS
except ImportError:
    WORDCLOUD_STOPWORDS = {
        "a", "an", "and", "are", "as", "be", "been", "being", "can", "could",
        "did", "do", "does", "for", "had", "has", "have", "he", "her", "his",
        "i", "if", "in", "is", "it", "its", "me", "my", "not", "of", "on",
        "or", "our", "she", "so", "that", "the", "their", "them", "there",
        "they", "this", "to", "was", "we", "were", "what", "when", "where",
        "which", "who", "will", "with", "would", "you", "your",
    }

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from dateutil import parser as date_parser
except ImportError:
    date_parser = None


HTML_TAG_RE = re.compile(r"<[^>]+>")
CHAT_ARTIFACT_RE = re.compile(
    r":\w+:"
    r"|\b(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|today|yesterday) at \d{1,2}:\d{2}\b"
    r"|\b\d+\s+repl(?:y|ies)\b"
    r"|\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}"
    r"|\[[^\]]+\]",
    flags=re.IGNORECASE,
)
URL_EMAIL_RE = re.compile(
    r"(?:https?://|www\.)[a-zA-Z0-9-]+(?:\.[a-zA-Z0-9-]+)+[^\s]*"
    r"|(?:[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})",
    flags=re.IGNORECASE,
)
NER_CAPS_RE = re.compile(
    r"\b(?:"
    r"[A-Z]{2,}"
    r"|[A-Z][a-zA-Z0-9-]*\d[a-zA-Z0-9-]*"
    r"|[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*"
    r")\b"
)


def default_prepositions() -> set:
    return {
        "about", "above", "across", "after", "against", "along", "among",
        "around", "at", "before", "behind", "below", "beneath", "beside",
        "between", "beyond", "but", "by", "concerning", "despite", "down",
        "during", "except", "for", "from", "in", "inside", "into", "like",
        "near", "of", "off", "on", "onto", "out", "outside", "over", "past",
        "regarding", "since", "through", "throughout", "to", "toward", "under",
        "underneath", "until", "up", "upon", "with", "within", "without",
        "something", "someone", "somebody", "anything", "anyone", "anybody",
        "everything", "everyone", "everybody", "nothing", "none", "maybe",
        "perhaps", "really", "basically", "generally", "kind", "sort", "stuff",
        "thing", "things",
    }


def build_default_stopwords(remove_prepositions: bool = True) -> set:
    stopwords = set(WORDCLOUD_STOPWORDS)
    if remove_prepositions:
        stopwords.update(default_prepositions())
    return stopwords


def build_punct_translation(keep_hyphens: bool, keep_apostrophes: bool) -> dict:
    punct = string.punctuation + "“”‘’–—…"
    if keep_hyphens:
        for char in "-–—":
            punct = punct.replace(char, "")
    if keep_apostrophes:
        for char in "'’":
            punct = punct.replace(char, "")
    return str.maketrans("", "", punct)


def build_phrase_pattern(phrases: List[str]) -> Optional[re.Pattern]:
    if not phrases:
        return None
    escaped = [re.escape(p) for p in phrases if p]
    if not escaped:
        return None
    return re.compile(rf"\b(?:{'|'.join(escaped)})\b", flags=re.IGNORECASE)


def parse_user_stopwords(raw: str) -> Tuple[List[str], List[str]]:
    raw = (raw or "").replace("\n", ",").replace(".", ",")
    phrases, singles = [], []
    for item in [x.strip() for x in raw.split(",") if x.strip()]:
        if " " in item:
            phrases.append(item.lower())
        else:
            singles.append(item.lower())
    return phrases, singles


def make_unique_header(raw_names: List[Optional[str]]) -> List[str]:
    seen: Dict[str, int] = {}
    result: List[str] = []
    for i, name in enumerate(raw_names):
        clean_name = str(name).strip() if name is not None else ""
        if not clean_name:
            clean_name = f"col_{i}"
        if clean_name in seen:
            seen[clean_name] += 1
            clean_name = f"{clean_name}__{seen[clean_name]}"
        else:
            seen[clean_name] = 1
        result.append(clean_name)
    return result


def detect_text_encoding(file_bytes: bytes, encoding_choice: str = "auto") -> str:
    if encoding_choice == "latin-1":
        return "latin-1"
    if file_bytes.startswith((b"\xff\xfe", b"\xfe\xff")):
        return "utf-16"
    if file_bytes.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    return "utf-8"


def clean_date_str(raw: Any) -> Optional[str]:
    if raw is None:
        return None
    s = str(raw).strip()
    if not s or s.lower() in {"nan", "none", "null"}:
        return None
    if date_parser:
        try:
            return date_parser.parse(s, fuzzy=True).strftime("%Y-%m-%d")
        except (ValueError, OverflowError):
            pass
    match = re.search(r"\d{4}-\d{2}-\d{2}", s)
    if match:
        return match.group(0)
    match_us = re.search(r"(\d{1,2})/(\d{1,2})/(\d{4})", s)
    if match_us:
        m, d, y = match_us.groups()
        return f"{y}-{int(m):02d}-{int(d):02d}"
    return None


def clean_category(raw: Any) -> Optional[str]:
    if raw is None:
        return None
    s = re.sub(r"\s+", " ", str(raw)).strip()
    if not s or s.lower() in {"nan", "none", "null"}:
        return None
    return s[:120]


def apply_text_cleaning(
    text: str,
    remove_chat: bool = True,
    remove_html: bool = True,
    unescape: bool = True,
    remove_urls: bool = True,
    phrase_pattern: Optional[re.Pattern] = None,
) -> str:
    if not isinstance(text, str):
        return ""
    text = text.replace("\ufeff", " ").replace("\ufffd", " ")
    text = "".join(ch if (ch >= " " or ch in "\n\r\t") else " " for ch in text)
    if remove_chat:
        text = CHAT_ARTIFACT_RE.sub(" ", text)
    if remove_html:
        text = HTML_TAG_RE.sub(" ", text)
    if unescape:
        try:
            text = html.unescape(text)
        except Exception:
            pass
    text = re.sub(r"(\w)-\s+(\w)", r"\1\2", text)
    text = re.sub(r"([A-Z]{2,})([a-z]{2,})", r"\1 \2", text)
    if remove_urls:
        text = URL_EMAIL_RE.sub(" ", text)
    text = text.lower()
    if phrase_pattern:
        text = phrase_pattern.sub(" ", text)
    return text.strip()


def clean_and_tokenize(
    text: str,
    remove_chat: bool,
    remove_html: bool,
    unescape: bool,
    remove_urls: bool,
    trans_map: dict,
    stopwords: set,
    phrase_pattern: Optional[re.Pattern],
    min_len: int,
    drop_int: bool,
) -> List[str]:
    cleaned = apply_text_cleaning(
        text,
        remove_chat=remove_chat,
        remove_html=remove_html,
        unescape=unescape,
        remove_urls=remove_urls,
        phrase_pattern=phrase_pattern,
    )
    tokens = []
    strip_chars = string.punctuation + "“”‘’–—"
    for token in cleaned.split():
        token = token.translate(trans_map).strip(strip_chars)
        if not token:
            continue
        if drop_int and token.isdigit():
            continue
        if len(token) < min_len:
            continue
        if token in stopwords:
            continue
        tokens.append(token)
    return tokens


def extract_entities_regex(text: str, stopwords: set) -> List[str]:
    if not isinstance(text, str):
        return []
    entities = []
    for candidate in NER_CAPS_RE.findall(text):
        if candidate.lower() in stopwords:
            continue
        if len(candidate) < 3:
            continue
        entities.append(candidate[:120])
    return entities


def build_evidence_doc(
    doc_id: int,
    raw_text: str,
    tokens: List[str],
    date_key: Optional[str],
    category_key: Optional[str],
    max_chars: int = 700,
    max_tokens: int = 160,
) -> Optional[Dict[str, Any]]:
    if not raw_text or not tokens:
        return None
    excerpt = html.unescape(str(raw_text))
    excerpt = re.sub(r"\s+", " ", excerpt).strip()
    if not excerpt:
        return None
    if len(excerpt) > max_chars:
        excerpt = excerpt[:max_chars].rsplit(" ", 1)[0] + "..."
    return {
        "id": doc_id,
        "excerpt": excerpt,
        "tokens": list(dict.fromkeys(tokens[:max_tokens])),
        "date": date_key,
        "category": category_key,
    }


def summarize_dashboard(
    counts: Counter,
    bigrams: Counter,
    category_counts: Dict[str, Counter],
    temporal_counts: Dict[str, Counter],
    entity_counts: Counter,
    total_rows: int,
    top_n: int = 10,
) -> Dict[str, Any]:
    top_words = [{"term": term, "count": count} for term, count in counts.most_common(top_n)]
    top_bigrams = [
        {"term": f"{pair[0]} {pair[1]}", "count": count}
        for pair, count in bigrams.most_common(top_n)
    ]
    top_entities = [
        {"entity": entity, "count": count}
        for entity, count in entity_counts.most_common(top_n)
    ]
    categories = [
        {"category": category, "rows_or_terms": sum(counter.values())}
        for category, counter in sorted(
            category_counts.items(),
            key=lambda item: sum(item[1].values()),
            reverse=True,
        )[:top_n]
    ]
    return {
        "generated_at": datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ"),
        "total_rows": total_rows,
        "unique_terms": len(counts),
        "total_tokens": sum(counts.values()),
        "top_words": top_words,
        "top_bigrams": top_bigrams,
        "top_entities": top_entities,
        "category_count": len(category_counts),
        "date_count": len(temporal_counts),
        "top_categories": categories,
    }


def read_rows_raw_lines(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[str]:
    encoding = detect_text_encoding(file_bytes, encoding_choice)
    bio = io.BytesIO(file_bytes)
    with io.TextIOWrapper(bio, encoding=encoding, errors="replace", newline=None) as wrapper:
        for line in wrapper:
            yield line.rstrip("\r\n")


def read_rows_vtt(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[str]:
    for line in read_rows_raw_lines(file_bytes, encoding_choice):
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit():
            continue
        if ":" in line:
            speaker, utterance = line.split(":", 1)
            if speaker and len(speaker) < 30 and " " in speaker:
                yield utterance.strip()
                continue
        yield line


def read_rows_pdf(file_bytes: bytes) -> Iterable[str]:
    import pypdf

    bio = io.BytesIO(file_bytes)
    try:
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                yield text
    except Exception:
        yield ""


def read_rows_pptx(file_bytes: bytes) -> Iterable[str]:
    import pptx

    bio = io.BytesIO(file_bytes)
    try:
        prs = pptx.Presentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text:
                    yield shape.text
    except Exception:
        yield ""


def get_csv_preview(file_bytes: bytes, encoding_choice: str = "auto") -> pd.DataFrame:
    encoding = detect_text_encoding(file_bytes, encoding_choice)
    bio = io.BytesIO(file_bytes)
    try:
        return pd.read_csv(bio, header=0, nrows=5, encoding=encoding, on_bad_lines="skip")
    except Exception:
        return pd.DataFrame()


def iter_csv_selected_columns(
    file_bytes: bytes,
    encoding_choice: str,
    delimiter: str,
    has_header: bool,
    selected_columns: List[str],
) -> Iterable[str]:
    encoding = detect_text_encoding(file_bytes, encoding_choice)
    bio = io.BytesIO(file_bytes)
    with io.TextIOWrapper(bio, encoding=encoding, errors="replace", newline="") as wrapper:
        reader = csv.reader(wrapper, delimiter=delimiter)
        first = next(reader, None)
        if first is None:
            return

        if has_header:
            header = make_unique_header(list(first))
            name_to_idx = {name: idx for idx, name in enumerate(header)}
            idxs = [name_to_idx[name] for name in selected_columns if name in name_to_idx]
        else:
            name_to_idx = {f"col_{idx}": idx for idx in range(len(first))}
            idxs = [name_to_idx[name] for name in selected_columns if name in name_to_idx]
            vals = [first[idx] if idx < len(first) else "" for idx in idxs]
            if any(vals):
                yield " ".join(str(v) for v in vals if v)

        for row in reader:
            vals = [row[idx] if idx < len(row) else "" for idx in idxs]
            if any(vals):
                yield " ".join(str(v) for v in vals if v)


def perform_topic_modeling(
    file_counts: List[Counter],
    n_topics: int = 4,
    top_n_words: int = 6,
    model_type: str = "LDA",
) -> Optional[List[Dict]]:
    from sklearn.decomposition import LatentDirichletAllocation, NMF
    from sklearn.feature_extraction import DictVectorizer

    valid_counts = [counter for counter in file_counts if counter and len(counter) > 0]
    if not valid_counts:
        return None

    vectorizer = DictVectorizer(sparse=True)
    dtm = vectorizer.fit_transform(valid_counts)
    n_samples, n_features = dtm.shape
    if n_samples == 0 or n_features == 0:
        return None

    safe_topics = min(n_topics, min(n_samples, n_features)) if model_type == "NMF" else min(n_topics, n_samples)
    if safe_topics < 1:
        return None

    if model_type == "NMF":
        model = NMF(n_components=safe_topics, random_state=42, init="nndsvd")
    else:
        model = LatentDirichletAllocation(
            n_components=safe_topics,
            random_state=42,
            learning_method="batch",
            max_iter=10,
        )
    model.fit(dtm)

    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(model.components_):
        top_indices = topic.argsort()[:-top_n_words - 1:-1]
        topics.append({
            "id": topic_idx + 1,
            "words": [feature_names[idx] for idx in top_indices],
        })
    return topics
