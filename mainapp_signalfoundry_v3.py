import io
import os
import re
import html
import gc
import time
import csv
import json
import math
import string
import zipfile
import tempfile
import logging
import secrets
from dataclasses import dataclass, field
from functools import lru_cache
from urllib.parse import urlparse
from collections import Counter, defaultdict
from typing import Dict, List, Tuple, Iterable, Optional, Callable, Any, Union, Set
from datetime import datetime

import numpy as np
import pandas as pd
import streamlit as st

try:
    import altair as alt
except Exception:
    alt = None
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from io import BytesIO
from wordcloud import WordCloud, STOPWORDS
from matplotlib import font_manager
# Requires Python 3.10+. Streamlit Community Cloud currently supports this,
# but requirements/README should keep the Python version expectation explicit.
from itertools import pairwise
import openai

# --- graph imports
import networkx as nx
import networkx.algorithms.community as nx_comm
from streamlit_agraph import agraph, Node, Edge, Config

# -3rd party imports checks
try:
    import requests
    from bs4 import BeautifulSoup
except ImportError:
    requests = None
    BeautifulSoup = None

try:
    import qrcode
except ImportError:
    qrcode = None

try:
    from scipy.stats import beta as beta_dist
except ImportError:
    beta_dist = None

try:
    from sklearn.decomposition import LatentDirichletAllocation, NMF
    from sklearn.feature_extraction import DictVectorizer
except ImportError:
    LatentDirichletAllocation = None
    NMF = None
    DictVectorizer = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from dateutil import parser as date_parser
except ImportError:
    date_parser = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import nltk
    from nltk.sentiment.vader import SentimentIntensityAnalyzer
    from nltk.stem import WordNetLemmatizer
except ImportError:
    nltk = None
    SentimentIntensityAnalyzer = None
    WordNetLemmatizer = None


# ⚙️ constants and config
# ========================================

MAX_TOPIC_DOCS = 50_000
MAX_EVIDENCE_DOCS = 10_000
MAX_EVIDENCE_CHARS = 700
MAX_SPEAKER_NAME_LENGTH = 30
SENTIMENT_ANALYSIS_TOP_N = 5000
URL_SCRAPE_RATE_LIMIT_SECONDS = 1.0
PROGRESS_UPDATE_MIN_INTERVAL = 100
NPMI_MIN_FREQ = 3
MAX_FILE_SIZE_MB = 1024

# regex patterns
HTML_TAG_RE = re.compile(r"<[^>]+>")
CHAT_ARTIFACT_RE = re.compile(
    r":\w+:"
    r"|\b(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|today|yesterday) at \d{1,2}:\d{2}\b"
    r"|\b\d+\s+repl(?:y|ies)\b"
    r"|\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}"
    r"|\[[^\]]+\]",
    flags=re.IGNORECASE
)
URL_EMAIL_RE = re.compile(
    r'(?:https?://|www\.)[a-zA-Z0-9-]+(?:\.[a-zA-Z0-9-]+)+[^\s]*'
    r'|(?:[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})',
    flags=re.IGNORECASE
)
# NER Pattern: matches 3 types of entities:
# acronyms (2+ Uppercase): DARPA, EU, NATO
# complex IDs (Cap start + digits/hyphens): COVID-19, Mi-6, G-7
# standard proper nouns (title case phrases): John Doe, Project Gutenberg
NER_CAPS_RE = re.compile(
    r'\b(?:'
    r'[A-Z]{2,}'                         # acronyms (DARPA)
    r'|[A-Z][a-zA-Z0-9-]*\d[a-zA-Z0-9-]*' # complex IDs (COVID-19)
    r'|[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*'    # standard names (John Doe)
    r')\b'
)

# logger setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("IntelEngine")

# custom exceptions
class ReaderError(Exception):
    pass

class ValidationError(Exception):
    pass

# 📦 dataclassed
# ===========================================

@dataclass
class CleaningConfig:
    remove_chat: bool = True
    remove_html: bool = True
    remove_urls: bool = True
    unescape: bool = True
    phrase_pattern: Optional[re.Pattern] = None

@dataclass
class ProcessingConfig:
    min_word_len: int = 2
    drop_integers: bool = True
    compute_bigrams: bool = True
    use_lemmatization: bool = False
    translate_map: Dict[int, Optional[int]] = field(default_factory=dict)
    stopwords: Set[str] = field(default_factory=set)
    excluded_speakers: Set[str] = field(default_factory=set)
    partial_speaker_match: bool = False


# 🛡️ security and validation utils
# ==========================================

def get_auth_password() -> str:
    pwd = st.secrets.get("auth_password")
    if not pwd:
        st.error("🚨 Configuration Error: 'auth_password' not set in .streamlit/secrets.toml.")
        st.stop()
    return pwd

def validate_url(url: str) -> bool:
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ('http', 'https'):
            return False
        if parsed.hostname in ('localhost', '127.0.0.1', '0.0.0.0', '::1'):
            return False
        return True
    except Exception:
        return False

def validate_sketch_data(data: Dict) -> bool:
    REQUIRED_KEYS = {"total_rows", "counts", "bigrams", "topic_docs"}
    if not isinstance(data, dict): return False
    if not REQUIRED_KEYS.issubset(data.keys()): return False
    return True


# 🧠 core logic, scanner
# =========================================

class StreamScanner:
    def __init__(self, doc_batch_size=5):
        self.global_counts = Counter()
        self.global_bigrams = Counter()
        self.total_rows_processed = 0
        self.topic_docs: List[Counter] = []
        
        # newer analytical structures (v2.0+)
        self.temporal_counts = defaultdict(Counter) # { '2023-10-01': {'word': 5} }
        self.category_counts = defaultdict(Counter) # { 'CategoryA': {'word': 10} }
        self.doc_freqs = Counter() # DF for TF-IDF
        self.entity_counts = Counter() # NER lite storage
        self.evidence_docs: List[Dict[str, Any]] = []
        self.evidence_limit_reached = False
        self.dashboard_summary: Dict[str, Any] = {}
        
        self.DOC_BATCH_SIZE = doc_batch_size
        self.limit_reached = False

    def set_batch_size(self, size: int):
        self.DOC_BATCH_SIZE = size

    def update_global_stats(self, counts: Counter, bigrams: Counter, rows: int):
        self.global_counts.update(counts)
        self.global_bigrams.update(bigrams)
        self.total_rows_processed += rows
    
    # only updates doc_freqs if still sampling (prevents unnecessary work)    
    def add_topic_sample(self, doc_counts: Counter):
        if not doc_counts: return
        
        if not self.limit_reached and len(self.topic_docs) < MAX_TOPIC_DOCS:
            self.doc_freqs.update(doc_counts.keys())
            self.topic_docs.append(doc_counts)
            if len(self.topic_docs) >= MAX_TOPIC_DOCS:
                self.limit_reached = True
        # if limit already reached, do nothing — saves memory/CPU
       
    def update_metadata_stats(self, date_key: Optional[str], cat_key: Optional[str], tokens: List[str]):
        if date_key:
            self.temporal_counts[date_key].update(tokens)
        if cat_key:
            self.category_counts[cat_key].update(tokens)

    def update_entities(self, entities: List[str]):
        if entities:
            self.entity_counts.update(entities)

    def add_evidence_doc(
        self,
        raw_text: str,
        tokens: List[str],
        date_key: Optional[str],
        cat_key: Optional[str],
    ):
        if self.evidence_limit_reached or not raw_text or not tokens:
            return
        if len(self.evidence_docs) >= MAX_EVIDENCE_DOCS:
            self.evidence_limit_reached = True
            return

        excerpt = html.unescape(str(raw_text))
        excerpt = re.sub(r"\s+", " ", excerpt).strip()
        if not excerpt:
            return
        if len(excerpt) > MAX_EVIDENCE_CHARS:
            excerpt = excerpt[:MAX_EVIDENCE_CHARS].rsplit(" ", 1)[0] + "..."

        self.evidence_docs.append({
            "id": len(self.evidence_docs) + 1,
            "excerpt": excerpt,
            "tokens": list(dict.fromkeys(tokens[:160])),
            "date": date_key,
            "category": cat_key,
        })

    def to_json(self) -> str:
        # simplifying complex structures for JSON serialization
        serializable_bigrams = {f"{k[0]}|{k[1]}": v for k, v in self.global_bigrams.items()}
        data = {
            "total_rows": self.total_rows_processed,
            "counts": dict(self.global_counts),
            "bigrams": serializable_bigrams,
            "topic_docs": [dict(c) for c in self.topic_docs],
            "limit_reached": self.limit_reached,
            # Persistence for new features
            "temporal_counts": {k: dict(v) for k, v in self.temporal_counts.items()},
            "category_counts": {k: dict(v) for k, v in self.category_counts.items()},
            "entity_counts": dict(self.entity_counts),
            "doc_freqs": dict(self.doc_freqs),
            "evidence_docs": self.evidence_docs,
            "evidence_limit_reached": self.evidence_limit_reached,
            "dashboard_summary": self.dashboard_summary,
        }
        return json.dumps(data)

    def load_from_json(self, json_str: str) -> bool:
        try:
            data = json.loads(json_str)
            if not validate_sketch_data(data): return False
            
            self.total_rows_processed = data.get("total_rows", 0)
            self.global_counts = Counter(data.get("counts", {}))
            
            raw_bigrams = data.get("bigrams", {})
            self.global_bigrams = Counter()
            for k, v in raw_bigrams.items():
                if "|" in k:
                    parts = k.split("|", 1)
                    self.global_bigrams[(parts[0], parts[1])] = v
            
            self.topic_docs = [Counter(d) for d in data.get("topic_docs", [])]
            self.limit_reached = data.get("limit_reached", False)
            
            # new
            self.entity_counts = Counter(data.get("entity_counts", {}))
            self.doc_freqs = Counter(data.get("doc_freqs", {}))
            self.evidence_docs = data.get("evidence_docs", [])[:MAX_EVIDENCE_DOCS]
            self.evidence_limit_reached = data.get("evidence_limit_reached", False)
            self.dashboard_summary = data.get("dashboard_summary", {})
            
            raw_temp = data.get("temporal_counts", {})
            self.temporal_counts = defaultdict(Counter)
            for k, v in raw_temp.items():
                self.temporal_counts[k] = Counter(v)

            raw_cat = data.get("category_counts", {})
            self.category_counts = defaultdict(Counter)
            for k, v in raw_cat.items():
                self.category_counts[k] = Counter(v)
                
            return True
        except Exception as e:
            logger.error(f"JSON Load Error: {e}")
            return False

# 📈 maturity modeling logic (multi-Persona)
# ==========================================

class MaturityAssessor:
    """
    Evaluates maturity based on linguistic markers using selectable 'Personas'.
    Supports switching between Business, EdTech, and Policy contexts.

    Now supports:
    - Single-word tokens (unigrams)
    - Multi-word phrases (bigrams), via the scanner's global_bigrams
    """
    def __init__(self):
        # library of models
        # NOTE: each level can now optionally define:
        #   - "terms": set[str]          -> single tokens
        #   - "phrases": list[str]       -> multi-word phrases (space-separated)
        self.models = {
            "🏫 EdTech & LMS Ops": {
                "desc": "Evaluates LMS utilization from 'Digital Repository' (L1) to 'Connected Ecosystem' (L5).",
                "levels": {
                    1: {
                        "name": "Digital Repository (Static)",
                        "color": "#d62728",
                        "terms": {
                            "upload", "download", "pdf", "file", "syllabus", "login", "password",
                            "access", "content", "link", "ppt", "doc", "email", "submit", "paper", "static"
                        },
                        "phrases": [
                            "course shell", "file repository", "content dump"
                        ]
                    },
                    2: {
                        "name": "Managed Courseware (Tools)",
                        "color": "#ff7f0e",
                        "terms": {
                            "quiz", "gradebook", "discussion", "rubric", "module", "assignment",
                            "calendar", "announcement", "feedback", "group", "template", "checklist", "forum"
                        },
                        "phrases": [
                            "online quiz", "assignment submission", "discussion forum"
                        ]
                    },
                    3: {
                        "name": "Integrated (Connected)",
                        "color": "#f7b731",
                        "terms": {
                            "lti", "integration", "api", "plugin", "interoperability", "tool",
                            "external", "sso", "vendor", "connect", "ecosystem", "zoom", "teams", "turnitin", "scorm"
                        },
                        "phrases": [
                            "learning tools", "third party", "external tool", "deep integration"
                        ]
                    },
                    4: {
                        "name": "Data-Informed (Adaptive)",
                        "color": "#2ca02c",
                        "terms": {
                            "analytics", "engagement", "retention", "risk", "dashboard", "report",
                            "outcome", "competency", "mastery", "release", "adaptive", "personalized", "pathway", "agent"
                        },
                        "phrases": [
                            "learning analytics", "course analytics", "early alert", "release conditions"
                        ]
                    },
                    5: {
                        "name": "Optimized (Strategic)",
                        "color": "#9467bd",
                        "terms": {
                            "governance", "strategy", "accessibility", "udl", "equity", "inclusion",
                            "continuous", "scale", "innovation", "agency", "holistic", "success", "lifelong", "transform"
                        },
                        "phrases": [
                            "governance framework", "continuous improvement", "student success", "institutional strategy"
                        ]
                    }
                }
            },
            "🏢 General Business Ops": {
                "desc": "Standard CMMI model: From 'Ad-Hoc' chaos to 'Optimized' strategy.",
                "levels": {
                    1: {
                        "name": "Ad-Hoc / Reactive",
                        "color": "#d62728",
                        "terms": {
                            "urgent", "fix", "panic", "broken", "late", "fail", "incident", "manual", "fire", "chaos"
                        },
                        "phrases": [
                            "fire drill", "last minute", "workaround"
                        ]
                    },
                    2: {
                        "name": "Managed / Project",
                        "color": "#ff7f0e",
                        "terms": {
                            "plan", "track", "project", "deadline", "schedule", "assign", "meeting", "status", "budget"
                        },
                        "phrases": [
                            "project plan", "status report", "project charter"
                        ]
                    },
                    3: {
                        "name": "Defined / Standardized",
                        "color": "#f7b731",
                        "terms": {
                            "standard", "process", "policy", "document", "compliance", "audit", "workflow", "consistent"
                        },
                        "phrases": [
                            "standard operating", "standardized process", "policy framework"
                        ]
                    },
                    4: {
                        "name": "Measured / Quantitative",
                        "color": "#2ca02c",
                        "terms": {
                            "metric", "kpi", "measure", "data", "analysis", "trend", "dashboard", "roi", "forecast"
                        },
                        "phrases": [
                            "key performance indicator", "data driven", "variance analysis"
                        ]
                    },
                    5: {
                        "name": "Optimizing / Strategic",
                        "color": "#9467bd",
                        "terms": {
                            "innovate", "strategy", "vision", "culture", "synergy", "scale", "optimize", "best-in-class"
                        },
                        "phrases": [
                            "strategic roadmap", "continuous improvement", "organizational transformation"
                        ]
                    }
                }
            },
            "⚖️ Policy & Governance": {
                "desc": "Evaluates policy maturity from 'Reactive/Enforcement' (L1) to 'Systemic/Holistic' (L5).",
                "levels": {
                    1: {
                        "name": "Enforcement / Reactive",
                        "color": "#d62728",
                        "terms": {
                            "violation", "sanction", "ban", "prohibit", "force", "threat", "danger",
                            "emergency", "crisis", "incident", "restriction", "penalty", "risk", "security"
                        },
                        "phrases": [
                            "zero tolerance", "strict enforcement"
                        ]
                    },
                    2: {
                        "name": "Procedural / Draft",
                        "color": "#ff7f0e",
                        "terms": {
                            "draft", "proposal", "clause", "article", "amendment", "review",
                            "committee", "meeting", "agenda", "timeline", "signature", "ratify", "consensus"
                        },
                        "phrases": [
                            "draft policy", "working group", "policy proposal"
                        ]
                    },
                    3: {
                        "name": "Operational / Implemented",
                        "color": "#f7b731",
                        "terms": {
                            "framework", "mechanism", "guideline", "standard", "monitor",
                            "verify", "compliance", "report", "mandate", "coordination", "treaty"
                        },
                        "phrases": [
                            "compliance framework", "implementation plan", "governance mechanism"
                        ]
                    },
                    4: {
                        "name": "Evidence-Based / Analysis",
                        "color": "#2ca02c",
                        "terms": {
                            "assessment", "evaluation", "impact", "data", "research", "finding",
                            "indicator", "measure", "trend", "forecast", "efficacy", "evidence"
                        },
                        "phrases": [
                            "impact assessment", "evidence based", "data driven"
                        ]
                    },
                    5: {
                        "name": "Systemic / Sustainable",
                        "color": "#9467bd",
                        "terms": {
                            "sustainable", "resilient", "holistic", "global", "ecosystem",
                            "peace", "development", "cooperation", "future", "inclusive", "norms", "universal"
                        },
                        "phrases": [
                            "holistic approach", "systemic change", "sustainable development"
                        ]
                    }
                }
            },
            "🎓 TAM Maturity Model (12-Domain)": {
    "desc": "Official 12-domain TAM Admin Maturity Model. Evaluates from Foundational to Leading Edge across Platform Admin, Curriculum, Student Engagement, Analytics, Assessment, Instructor Efficiency, Change Management, Knowledge Management, Accessibility, User Support, Innovation, and Collaboration.",
    "type": "domain_based",
    "domains": {
        "01_platform_admin": {
            "name": "Platform & Technical Administration",
            "short": "Platform Admin",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "upload", "download", "csv", "manual", "default", "password",
                        "login", "enrol", "enroll", "patching", "basic", "static",
                        "broad", "minimal", "limited"
                    },
                    "phrases": [
                        "course shell", "default settings", "csv uploads",
                        "manual updates", "default configurations",
                        "minimal customization", "basic password",
                        "manual user", "limited api", "limited use"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "api", "automation", "sso", "lti", "scorm", "xapi",
                        "scripting", "compliance", "dashboard", "provisioning",
                        "workflow", "documentation", "blueprint", "naming",
                        "monitoring", "wcag", "gdpr", "impersonation"
                    },
                    "phrases": [
                        "role based", "access control", "data hub",
                        "naming conventions", "least privilege",
                        "automated user", "course provisioning",
                        "standard operating", "regular review",
                        "documented processes", "sis integration"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "governance", "middleware", "forecasting", "transformation",
                        "scalable", "monitoring", "capacity", "branded",
                        "lifecycle", "archiving", "penetration", "siem"
                    },
                    "phrases": [
                        "continuous improvement", "capacity planning",
                        "custom scripts", "integration roadmap",
                        "real-time monitoring", "power bi", "tableau",
                        "automation strategy", "performance monitoring",
                        "full automation", "cross-platform automation",
                        "system health"
                    ]
                }
            }
        },
        "02_curriculum": {
            "name": "Curriculum Development & Delivery",
            "short": "Curriculum",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "syllabus", "outline", "upload", "organize",
                        "template", "schedule", "guideline", "accessible"
                    },
                    "phrases": [
                        "learning objectives", "course materials",
                        "curriculum goals", "basic accessibility",
                        "course design", "content effectively"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "instructional", "multimedia", "interactive",
                        "rubric", "inclusive", "formative", "differentiated",
                        "personalize"
                    },
                    "phrases": [
                        "instructional design", "active learning",
                        "diverse learning", "release conditions",
                        "learning styles", "course content"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "accreditation", "competency", "programmatic",
                        "innovation", "pedagogical"
                    },
                    "phrases": [
                        "program-level learning", "cross-functional teams",
                        "competency-based learning", "continuous improvement",
                        "curriculum evaluation", "institutional goals",
                        "curriculum improvements"
                    ]
                }
            }
        },
        "03_student_engagement": {
            "name": "Student Engagement & Success",
            "short": "Student Engagement",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "navigate", "announcement", "discussion", "tutorial",
                        "onboarding", "calendar", "login", "access"
                    },
                    "phrases": [
                        "course shell", "technical support",
                        "student learning", "access issues",
                        "course access", "onboarding materials"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "analytics", "disengagement", "intervention",
                        "quiz", "checklist", "advisor", "workshop"
                    },
                    "phrases": [
                        "intelligent agents", "early alerts",
                        "class progress", "engagement dashboards",
                        "active learning", "digital engagement",
                        "engagement reports"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "predictive", "retention", "progression",
                        "holistic", "inclusive", "scalable"
                    },
                    "phrases": [
                        "predictive analytics", "student success",
                        "data hub", "cross-functional initiatives",
                        "institutional goals", "real-time engagement",
                        "students at risk"
                    ]
                }
            }
        },
        "04_data_analytics": {
            "name": "Data & Learning Analytics",
            "short": "Data Analytics",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "report", "export", "gradebook", "basic",
                        "manual", "spreadsheet"
                    },
                    "phrases": [
                        "class progress", "basic reports",
                        "grade exports", "manual tracking",
                        "default reports"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "analytics", "dashboard", "visualization",
                        "benchmark", "kpi", "metric", "trend"
                    },
                    "phrases": [
                        "data-informed", "engagement dashboards",
                        "learning analytics", "custom reports",
                        "data hub", "decision making"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "predictive", "modeling", "warehouse",
                        "longitudinal", "algorithm"
                    },
                    "phrases": [
                        "predictive analytics", "power bi", "tableau",
                        "data warehouse", "institutional research",
                        "cross-functional initiatives",
                        "data governance"
                    ]
                }
            }
        },
        "05_assessment": {
            "name": "Assessment & Evaluation",
            "short": "Assessment",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "quiz", "assignment", "gradebook", "summative",
                        "multiple-choice", "grading", "manual"
                    },
                    "phrases": [
                        "auto-graded quizzes", "default settings",
                        "assignment uploads", "basic assignments",
                        "getting started", "step-by-step tutorials"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "rubric", "formative", "peer", "reflective",
                        "randomization", "integrity", "feedback"
                    },
                    "phrases": [
                        "peer review", "question banks",
                        "learning outcomes", "peer assessment",
                        "assessment design", "rubric libraries",
                        "knowledge checks", "video assignments"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "adaptive", "simulation", "personalized",
                        "branching", "turnitin", "h5p", "innovation"
                    },
                    "phrases": [
                        "item analysis", "real-time feedback",
                        "adaptive learning", "innovation labs",
                        "pilot projects", "assessment innovation",
                        "delegated marking", "co-marking"
                    ]
                }
            }
        },
        "06_instructor_efficiency": {
            "name": "Instructor Efficiency",
            "short": "Instructor Efficiency",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "onboarding", "reactive", "tutorial", "guide",
                        "slides", "compliance", "upload", "training",
                        "support", "help", "manual", "orientation",
                        "faculty", "instructor"
                    },
                    "phrases": [
                        "step-by-step guides", "video tutorials",
                        "getting started", "drop-in help",
                        "technical support", "support dependency",
                        "faculty training", "instructor support",
                        "basic training", "how-to guides"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "proactive", "contextual", "champion", "proficiency",
                        "troubleshoot", "multimedia", "padlet", "jamboard",
                        "captioned", "coaching", "workshop", "consultation",
                        "adoption", "readiness", "enablement"
                    },
                    "phrases": [
                        "best practices", "faculty champions",
                        "pedagogical integration", "analytics dashboards",
                        "peer learning", "faculty showcases",
                        "user groups", "one-on-one coaching",
                        "faculty development", "contextual support",
                        "training sessions", "teaching support"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "collaborative", "strategic", "innovation",
                        "mentorship", "gamified", "flipped", "adaptive",
                        "personalized", "agency", "transformation",
                        "scalable", "empowerment"
                    },
                    "phrases": [
                        "change agents", "custom tool integrations",
                        "data-driven", "beta testing",
                        "ai tools", "edtech innovation",
                        "data exports", "continuous improvement",
                        "faculty learning communities",
                        "instructional design partnership",
                        "teaching innovation"
                    ]
                }
            }
        },
        "07_change_management": {
            "name": "Change Management",
            "short": "Change Mgmt",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "update", "outage", "rollout", "disruption",
                        "communicate", "documentation"
                    },
                    "phrases": [
                        "system updates", "new feature",
                        "just-in-time support", "minimize disruption",
                        "pros and cons"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "adoption", "sandbox", "webinar", "feedback",
                        "champion", "training"
                    },
                    "phrases": [
                        "training sessions", "sandbox environments",
                        "rollout strategies", "instructor feedback",
                        "best practices", "faculty champions"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "transformation", "strategy", "innovation",
                        "pilot", "culture"
                    },
                    "phrases": [
                        "change management plans", "digital transformation",
                        "continuous improvement", "communities of practice",
                        "cross-functional collaboration",
                        "pilot programs"
                    ]
                }
            }
        },
        "08_knowledge_management": {
            "name": "Knowledge & Resource Management",
            "short": "Knowledge Mgmt",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "documentation", "faq", "guide", "reactive",
                        "manual", "organize", "upload"
                    },
                    "phrases": [
                        "knowledge base", "help site",
                        "how-to guides", "support channels",
                        "course shells", "case-by-case"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "template", "library", "curate", "repository",
                        "lor", "reuse", "redundancy", "analytics"
                    },
                    "phrases": [
                        "shared content", "content libraries",
                        "course copy", "instructional designers",
                        "institutional standards", "content usage"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "governance", "lifecycle", "versioning",
                        "archiving", "metadata", "tagging",
                        "discoverability", "stewardship"
                    },
                    "phrases": [
                        "communities of practice", "knowledge strategy",
                        "content lifecycle", "metadata standards",
                        "cross-departmental collaboration",
                        "oer repositories"
                    ]
                }
            }
        },
        "09_accessibility": {
            "name": "Accessibility & Compliance",
            "short": "Accessibility",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "accessibility", "ada", "wcag", "captioning",
                        "keyboard", "alt", "readspeaker", "ally"
                    },
                    "phrases": [
                        "screen reader", "alternative text",
                        "accessibility standards", "colour contrast",
                        "accessible content", "legal requirements"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "audit", "remediation", "inclusive", "assistive",
                        "consultation", "template", "checklist"
                    },
                    "phrases": [
                        "inclusive course design", "assistive technologies",
                        "course accessibility", "accessible course materials",
                        "content remediation", "captioning legacy"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "champion", "mentor", "procurement", "innovation",
                        "immersive"
                    },
                    "phrases": [
                        "accessibility champions", "accessibility compliance",
                        "procurement policies", "ai-driven captioning",
                        "inclusive pedagogy", "accessibility audits"
                    ]
                }
            }
        },
        "10_user_support": {
            "name": "User Support & Training",
            "short": "User Support",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "login", "navigation", "tutorial", "faq",
                        "guide", "gradebook", "assignment"
                    },
                    "phrases": [
                        "quick-start guides", "video tutorials",
                        "login support", "basic training",
                        "core tools"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "tiered", "workshop", "rubric", "multimedia",
                        "contextualized"
                    },
                    "phrases": [
                        "tiered training", "course design",
                        "assessment strategies", "student engagement",
                        "release conditions", "contextualized support"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "coaching", "consultation", "champion",
                        "personalize", "competency"
                    },
                    "phrases": [
                        "one-on-one coaching", "instructional design consultations",
                        "faculty learning communities",
                        "intelligent agents", "competency frameworks",
                        "learning analytics"
                    ]
                }
            }
        },
        "11_innovation": {
            "name": "Innovation & Emerging Technologies",
            "short": "Innovation",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "sandbox", "demo", "showcase", "curiosity",
                        "explore", "experiment", "gamification"
                    },
                    "phrases": [
                        "sandbox environments", "demo sessions",
                        "tech showcases", "immersive media",
                        "ai-powered tools"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "pilot", "adaptive", "simulation", "framework",
                        "evaluate", "effectiveness"
                    },
                    "phrases": [
                        "hands-on training", "pilot programs",
                        "adaptive learning", "ai feedback",
                        "virtual simulations", "technology adoption"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "research", "innovation", "equity", "transformation",
                        "scaled"
                    },
                    "phrases": [
                        "innovation labs", "communities of practice",
                        "digital pedagogy", "institutional goals",
                        "change agents", "faculty-led research"
                    ]
                }
            }
        },
        "12_collaboration": {
            "name": "Collaboration & Communication",
            "short": "Collaboration",
            "tiers": {
                1: {
                    "name": "Foundational",
                    "color": "#d62728",
                    "terms": {
                        "support", "troubleshooting", "reactive",
                        "coordinate", "respond"
                    },
                    "phrases": [
                        "support requests", "system updates",
                        "service platform", "issue resolution",
                        "one-way communication", "tool focused"
                    ]
                },
                2: {
                    "name": "Advanced",
                    "color": "#ff7f0e",
                    "terms": {
                        "consultative", "engage", "curriculum",
                        "faculty", "feedback", "align"
                    },
                    "phrases": [
                        "pedagogical partner", "two-way communication",
                        "curriculum planning", "faculty development",
                        "cross-functional meetings",
                        "academic priorities"
                    ]
                },
                3: {
                    "name": "Leading Edge",
                    "color": "#2ca02c",
                    "terms": {
                        "strategic", "governance", "transformation",
                        "vision", "proactive"
                    },
                    "phrases": [
                        "institutional strategy", "governance committees",
                        "strategic initiatives", "digital transformation",
                        "shared vision", "data-informed",
                        "emerging technologies"
                    ]
                }
            }
        }
    }
}
        }

    def get_model_names(self) -> List[str]:
        return list(self.models.keys())

    def get_model_desc(self, name: str) -> str:
        return self.models.get(name, {}).get("desc", "")

    def assess(self, counts: Counter, bigrams: Counter, model_name: str) -> Dict:
        """Routes to flat or domain-based assessment based on model type."""
        if model_name not in self.models:
            return None
        model = self.models[model_name]
        if model.get("type") == "domain_based":
            return self._assess_domain_based(counts, bigrams, model_name)
        else:
            return self._assess_flat(counts, bigrams, model_name)
    def _assess_flat(self, counts: Counter, bigrams: Counter, model_name: str) -> Dict:
        """Original 5-level flat assessment. Unchanged logic."""
        if model_name not in self.models:
            return None
        levels = self.models[model_name]["levels"]
        scores = {1: 0.0, 2: 0.0, 3: 0.0, 4: 0.0, 5: 0.0}
        total_hits = 0.0
        for lvl, data in levels.items():
            terms = data.get("terms", set())
            for term in terms:
                qty = counts.get(term, 0)
                if qty > 0:
                    scores[lvl] += qty
                    total_hits += qty
        phrase_counts: Dict[str, int] = {}
        for (w1, w2), freq in bigrams.items():
            phrase = f"{w1} {w2}"
            phrase_counts[phrase] = phrase_counts.get(phrase, 0) + freq
        # Phrase hits carry more context than single-word hits, so they receive
        # a modest boost without letting one phrase dominate the whole score.
        PHRASE_WEIGHT = 1.5
        for lvl, data in levels.items():
            phrases = data.get("phrases", []) or []
            for phrase in phrases:
                freq = phrase_counts.get(phrase, 0)
                if freq > 0:
                    weighted = freq * PHRASE_WEIGHT
                    scores[lvl] += weighted
                    total_hits += weighted
        if total_hits == 0:
            return None
        distribution = {k: v / total_hits for k, v in scores.items()}
        weighted_sum = sum(lvl * pct for lvl, pct in distribution.items())
        dominant_level = max(distribution, key=distribution.get)
        return {
            "type": "flat",
            "overall_score": round(weighted_sum, 2),
            "distribution": distribution,
            "dominant_stage": levels[dominant_level],
            "total_signals_found": int(total_hits),
            "levels_ref": levels
        }
    def _assess_domain_based(self, counts: Counter, bigrams: Counter, model_name: str) -> Dict:
        """
        12-domain, 3-tier assessment. Returns per-domain scores + composite.
        Score range: 1.0 (all Foundational) to 3.0 (all Leading Edge).
        """
        if model_name not in self.models:
            return None
        domains = self.models[model_name]["domains"]
        # Phrase hits carry more context than single-word hits, so they receive
        # a modest boost without letting one phrase dominate a domain score.
        PHRASE_WEIGHT = 1.5
        # Build bigram lookup once
        phrase_counts: Dict[str, int] = {}
        for (w1, w2), freq in bigrams.items():
            phrase = f"{w1} {w2}"
            phrase_counts[phrase] = phrase_counts.get(phrase, 0) + freq
        domain_results = {}
        total_signals_all = 0
        for domain_key, domain_data in domains.items():
            tiers = domain_data["tiers"]
            tier_scores = {1: 0.0, 2: 0.0, 3: 0.0}
            domain_hits = 0.0
            driver_terms = {}
            for tier_num, tier_data in tiers.items():
                tier_term_hits = []
                # Unigram hits
                for term in tier_data.get("terms", set()):
                    qty = counts.get(term, 0)
                    if qty > 0:
                        tier_scores[tier_num] += qty
                        domain_hits += qty
                        tier_term_hits.append((term, qty))
                # Phrase hits
                for phrase in tier_data.get("phrases", []):
                    freq = phrase_counts.get(phrase, 0)
                    if freq > 0:
                        weighted = freq * PHRASE_WEIGHT
                        tier_scores[tier_num] += weighted
                        domain_hits += weighted
                        tier_term_hits.append((phrase, freq))
                # Sort by count descending, keep top 5
                driver_terms[tier_num] = sorted(
                    tier_term_hits, key=lambda x: x[1], reverse=True
                )[:5]
            if domain_hits == 0:
                domain_results[domain_key] = {
                    "name": domain_data["name"],
                    "short": domain_data["short"],
                    "score": 0.0,
                    "tier_label": "No Data",
                    "distribution": {1: 0.0, 2: 0.0, 3: 0.0},
                    "signals": 0,
                    "drivers": driver_terms
                }
                continue
            total_signals_all += domain_hits
            # Normalize distribution
            distribution = {k: v / domain_hits for k, v in tier_scores.items()}
            # Weighted score (1.0 - 3.0)
            weighted_score = sum(tier * pct for tier, pct in distribution.items())
            # Determine tier label
            if weighted_score >= 2.5:
                tier_label = "Leading Edge"
            elif weighted_score >= 1.5:
                tier_label = "Advanced"
            else:
                tier_label = "Foundational"
            domain_results[domain_key] = {
                "name": domain_data["name"],
                "short": domain_data["short"],
                "score": round(weighted_score, 2),
                "tier_label": tier_label,
                "distribution": distribution,
                "signals": int(domain_hits),
                "drivers": driver_terms
            }
        # Composite score (average of domains that have data)
        scored_domains = [
            d for d in domain_results.values() if d["signals"] > 0
        ]
        if not scored_domains:
            return None
        composite_score = sum(d["score"] for d in scored_domains) / len(scored_domains)
        return {
            "type": "domain_based",
            "overall_score": round(composite_score, 2),
            "max_score": 3.0,
            "domains_assessed": len(scored_domains),
            "domains_total": 12,
            "total_signals_found": int(total_signals_all),
            "domain_results": domain_results,
            "domains_ref": domains
        }

    def render_radar_chart(self, result: Dict):
        """Generates a Radar/Spider chart."""
        if not result:
            return None

        levels_ref = result['levels_ref']
        categories = [levels_ref[i]['name'] for i in range(1, 6)]
        values = [result['distribution'][i] for i in range(1, 6)]

        values += values[:1]
        angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
        angles += angles[:1]

        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
        ax.fill(angles, values, color='#1f77b4', alpha=0.25)
        ax.plot(angles, values, color='#1f77b4', linewidth=2)
        ax.set_theta_offset(np.pi / 2)
        ax.set_theta_direction(-1)
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories, fontsize=9)
        ax.set_yticklabels([])
        ax.spines["polar"].set_visible(False)
        ax.grid(color='#444444', linestyle='--', alpha=0.5)
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)

        return fig
    def render_domain_radar_chart(
        self,
        result: Dict,
        font_family: str = "sans-serif",
        label_size: int = 9,
        label_color: str = "#222222",
        wrap_width: int = 14,
        show_tier_labels: bool = True,
        tier_label_angle: int = 35,
    ):
        """12-spoke spider chart for domain-based assessment."""
        if not result or result.get("type") != "domain_based":
            return None

        import textwrap

        domain_results = result["domain_results"]
        labels = []
        values = []
        colors = []

        for key in sorted(domain_results.keys()):
            dr = domain_results[key]
            wrapped_label = "\n".join(
                textwrap.wrap(str(dr["short"]), width=max(6, wrap_width))
            )
            labels.append(wrapped_label)
            values.append(dr["score"])

            if dr["score"] >= 2.5:
                colors.append("#2ca02c")
            elif dr["score"] >= 1.5:
                colors.append("#ff7f0e")
            else:
                colors.append("#d62728")

        n = len(labels)
        if n == 0:
            return None

        values_closed = values + values[:1]
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        angles_closed = angles + angles[:1]

        fig, ax = plt.subplots(figsize=(8.2, 8.2), subplot_kw=dict(polar=True))
        ax.fill(angles_closed, values_closed, color="#1f77b4", alpha=0.15)
        ax.plot(angles_closed, values_closed, color="#1f77b4", linewidth=2)

        for angle, val, color in zip(angles, values, colors):
            ax.scatter(angle, val, color=color, s=80, zorder=5)

        ax.set_theta_offset(np.pi / 2)
        ax.set_theta_direction(-1)
        ax.set_xticks(angles)
        ax.set_xticklabels(
            labels,
            fontsize=label_size,
            fontfamily=font_family,
            color=label_color,
        )

        # Move domain labels outward to reduce collisions with ring labels.
        ax.tick_params(axis="x", pad=18)

        ax.set_ylim(0, 3.0)
        ax.set_yticks([1.0, 2.0, 3.0])

        if show_tier_labels:
            ax.set_yticklabels(
                ["Foundational", "Advanced", "Leading\nEdge"],
                fontsize=max(7, label_size - 2),
                fontfamily=font_family,
                color=label_color,
            )
        else:
            ax.set_yticklabels(
                ["1.0", "2.0", "3.0"],
                fontsize=max(7, label_size - 2),
                fontfamily=font_family,
                color=label_color,
            )

        # Move maturity ring labels away from the crowded top-right label area.
        ax.set_rlabel_position(tier_label_angle)

        ax.spines["polar"].set_visible(False)
        ax.grid(color="#444444", linestyle="--", alpha=0.5)
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        fig.tight_layout(pad=2.5)

        return fig
        
        
    def render_domain_breakdown_chart(self, result: Dict):
        """Horizontal stacked bar chart showing tier distribution per domain."""
        if not result or result.get("type") != "domain_based":
            return None
        domain_results = result["domain_results"]
        sorted_keys = sorted(domain_results.keys())
        labels = []
        found_vals = []
        adv_vals = []
        lead_vals = []
        for key in sorted_keys:
            dr = domain_results[key]
            labels.append(dr["short"])
            dist = dr["distribution"]
            found_vals.append(dist.get(1, 0) * 100)
            adv_vals.append(dist.get(2, 0) * 100)
            lead_vals.append(dist.get(3, 0) * 100)
        y_pos = np.arange(len(labels))
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.barh(y_pos, found_vals, color='#d62728', label='Foundational', height=0.6)
        ax.barh(y_pos, adv_vals, left=found_vals, color='#ff7f0e', label='Advanced', height=0.6)
        lefts = [f + a for f, a in zip(found_vals, adv_vals)]
        ax.barh(y_pos, lead_vals, left=lefts, color='#2ca02c', label='Leading Edge', height=0.6)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=9)
        ax.set_xlabel("Signal Distribution (%)")
        ax.set_xlim(0, 100)
        ax.legend(loc='lower right', fontsize=8)
        ax.invert_yaxis()
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        ax.grid(axis='x', alpha=0.3)
        plt.tight_layout()
        return fig
    def render_trend_chart(self, snapshots: list):
        """Line chart showing composite score over time across snapshots."""
        if not snapshots or len(snapshots) < 2:
            return None
        dates = [s["assessment_date"] for s in snapshots]
        scores = [s["composite_score"] for s in snapshots]
        fig, ax = plt.subplots(figsize=(10, 4))
        ax.plot(dates, scores, marker='o', linewidth=2, color='#1f77b4', markersize=8)
        # Color-code background bands
        ax.axhspan(1.0, 1.5, alpha=0.08, color='#d62728', label='Foundational')
        ax.axhspan(1.5, 2.5, alpha=0.08, color='#ff7f0e', label='Advanced')
        ax.axhspan(2.5, 3.0, alpha=0.08, color='#2ca02c', label='Leading Edge')
        ax.set_ylim(0.8, 3.2)
        ax.set_ylabel("Composite Maturity Score")
        ax.set_xlabel("Assessment Date")
        ax.legend(loc='lower right', fontsize=8)
        plt.xticks(rotation=45, ha='right')
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        ax.grid(axis='y', alpha=0.3)
        plt.tight_layout()
        return fig
    def render_domain_trend_chart(self, snapshots: list):
        """Multi-line chart showing per-domain scores over time."""
        if not snapshots or len(snapshots) < 2:
            return None
        dates = [s["assessment_date"] for s in snapshots]
        # Collect all domain keys across all snapshots
        all_domains = set()
        for s in snapshots:
            all_domains.update(s.get("domain_results", {}).keys())
        fig, ax = plt.subplots(figsize=(12, 6))
        domain_colors = [
            "#e6194b", "#3cb44b", "#ffe119", "#4363d8", "#f58231", "#911eb4",
            "#42d4f4", "#f032e6", "#bfef45", "#fabed4", "#469990", "#dcbeff"
        ]
        for i, dk in enumerate(sorted(all_domains)):
            domain_scores = []
            for s in snapshots:
                dr = s.get("domain_results", {}).get(dk, {})
                domain_scores.append(dr.get("score", 0.0))
            short_name = snapshots[0].get("domain_results", {}).get(
                dk, {}
            ).get("short", dk)
            color = domain_colors[i % len(domain_colors)]
            ax.plot(dates, domain_scores, marker='o', linewidth=1.5,
                    label=short_name, color=color, markersize=5)
        ax.set_ylim(0.8, 3.2)
        ax.set_ylabel("Domain Score")
        ax.set_xlabel("Assessment Date")
        ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=7)
        plt.xticks(rotation=45, ha='right')
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        ax.grid(axis='y', alpha=0.3)
        plt.tight_layout()
        return fig



TAM_DOMAIN_HELP = {
    "01_platform_admin": {
        "plain": "How the LMS environment is configured, secured, integrated, monitored, and maintained.",
        "covers": "Roles and permissions, SSO, SIS/LTI/API integrations, course provisioning, configuration standards, monitoring, automation, lifecycle practices, and governance of the technical environment.",
        "leading": "A leading-edge environment uses automated provisioning, clear governance, integrated systems, proactive monitoring, scalable admin patterns, and documented lifecycle practices."
    },
    "02_curriculum": {
        "plain": "How courses and learning experiences are designed, structured, delivered, reviewed, and improved.",
        "covers": "Course templates, learning objectives, instructional design, content structure, release conditions, multimedia, accessibility within course design, competency alignment, and program-level curriculum improvement.",
        "leading": "Leading-edge curriculum work is programmatic, evidence-informed, accessible, competency-aware, and continuously improved through cross-functional collaboration."
    },
    "03_student_engagement": {
        "plain": "How the platform supports learner participation, persistence, retention, and success.",
        "covers": "Announcements, discussions, checklists, intelligent agents, engagement dashboards, early alerts, student-success workflows, advisor touchpoints, retention signals, and interventions for students at risk.",
        "leading": "Leading-edge engagement uses predictive signals, scalable interventions, real-time insight, and coordinated student-success practices."
    },
    "04_data_analytics": {
        "plain": "How data is collected, interpreted, shared, and used to improve teaching, learning, and operations.",
        "covers": "Reports, dashboards, Data Hub, learning analytics, KPIs, trend analysis, institutional research, data governance, longitudinal analysis, and visualization tools such as Power BI or Tableau.",
        "leading": "Leading-edge analytics connects platform data to institutional decision-making, predictive models, governed data practices, and cross-functional improvement cycles."
    },
    "05_assessment": {
        "plain": "How assessment activities are designed, delivered, evaluated, and improved.",
        "covers": "Quizzes, assignments, rubrics, gradebook practices, formative assessment, peer review, question banks, academic integrity, feedback workflows, item analysis, adaptive assessment, and assessment innovation.",
        "leading": "Leading-edge assessment uses high-quality feedback loops, adaptive or authentic assessment design, analytics-driven evaluation, and scalable assessment practices."
    },
    "06_instructor_efficiency": {
        "plain": "How effectively instructors can use LMS tools without unnecessary friction or manual effort.",
        "covers": "Instructor onboarding, training, support dependency, templates, reusable workflows, troubleshooting, faculty champions, instructional design partnership, automation, and time-saving practices.",
        "leading": "Leading-edge instructor efficiency is supported by scalable enablement, coaching, reusable models, automation, faculty learning communities, and innovation partnerships."
    },
    "07_change_management": {
        "plain": "How the organization plans, communicates, supports, and sustains platform-related change.",
        "covers": "Rollouts, system updates, training plans, sandbox environments, feedback loops, champion networks, adoption planning, change management plans, pilot programs, and digital transformation work.",
        "leading": "Leading-edge change management is proactive, strategic, cross-functional, evidence-informed, and tied to continuous improvement rather than one-time rollout support."
    },
    "08_knowledge_management": {
        "plain": "How institutional knowledge, documentation, templates, resources, and reusable assets are managed.",
        "covers": "Knowledge bases, FAQs, help sites, how-to guides, shared content libraries, learning object repositories, metadata, tagging, lifecycle management, versioning, archiving, and stewardship.",
        "leading": "Leading-edge knowledge management uses governed content lifecycles, discoverable resources, reusable standards, metadata practices, and communities of practice."
    },
    "09_accessibility": {
        "plain": "How accessibility, inclusion, compliance, and universal design are built into platform and course practices.",
        "covers": "WCAG, ADA/AODA-style compliance, alternative text, screen reader compatibility, captioning, templates, remediation, audits, assistive technology, inclusive pedagogy, procurement, and accessibility champions.",
        "leading": "Leading-edge accessibility is proactive, embedded in design and procurement, supported by champions, and continuously audited and improved."
    },
    "10_user_support": {
        "plain": "How learners, instructors, and staff receive help, training, onboarding, and ongoing support.",
        "covers": "Quick-start guides, login help, tutorials, FAQs, workshops, tiered support, contextualized training, coaching, consultations, learning communities, and role-specific enablement.",
        "leading": "Leading-edge user support combines self-service, role-based training, coaching, communities of practice, analytics-informed support, and continuous enablement."
    },
    "11_innovation": {
        "plain": "How the organization experiments with, evaluates, and scales emerging teaching and learning technologies.",
        "covers": "Sandboxes, demos, showcases, pilot programs, adaptive learning, simulations, AI tools, immersive media, innovation labs, faculty-led research, digital pedagogy, and emerging technology governance.",
        "leading": "Leading-edge innovation is tied to institutional goals, evaluated for impact and equity, and scaled through pilots, research, communities, and change agents."
    },
    "12_collaboration": {
        "plain": "How technology, academic, support, and leadership teams communicate and work together around the LMS.",
        "covers": "Support requests, service coordination, two-way communication, pedagogical partnership, faculty development, cross-functional meetings, governance committees, shared vision, and strategic initiatives.",
        "leading": "Leading-edge collaboration is strategic, proactive, cross-functional, governance-supported, and aligned to institutional priorities and digital transformation goals."
    },
}


def get_tam_domain_help(domain_key: str) -> Dict[str, str]:
    """Returns plain-language help text for a TAM / LMS maturity domain."""
    return TAM_DOMAIN_HELP.get(domain_key, {
        "plain": "No description is available for this domain yet.",
        "covers": "No coverage details are available yet.",
        "leading": "No leading-edge example is available yet.",
    })


def render_tam_domain_glossary(domain_results: Dict[str, Dict]):
    """Renders a compact glossary for the 12 maturity domains."""
    if not domain_results:
        return

    rows = []
    for domain_key in sorted(domain_results.keys()):
        dr = domain_results[domain_key]
        help_text = get_tam_domain_help(domain_key)
        rows.append({
            "Domain": dr.get("short", dr.get("name", domain_key)),
            "Plain-language meaning": help_text["plain"],
            "What it covers": help_text["covers"],
            "Leading-edge signals": help_text["leading"],
        })

    with st.expander("📘 What do these 12 domains mean?", expanded=False):
        st.caption(
            "Use this glossary to interpret the radar and domain cards. "
            "These explanations do not change the scoring logic; they explain what each domain is intended to represent."
        )
        st.dataframe(
            pd.DataFrame(rows),
            use_container_width=True,
            hide_index=True,
            column_config={
                "Domain": st.column_config.TextColumn(
                    "Domain",
                    help="The short label used in the radar and breakdown charts."
                ),
                "Plain-language meaning": st.column_config.TextColumn(
                    "Plain-language meaning",
                    help="A quick explanation of the domain."
                ),
                "What it covers": st.column_config.TextColumn(
                    "What it covers",
                    help="Examples of practices, workflows, or concepts included in this domain."
                ),
                "Leading-edge signals": st.column_config.TextColumn(
                    "Leading-edge signals",
                    help="What stronger maturity can look like in this area."
                ),
            },
        )


# session state init
if 'sketch' not in st.session_state: st.session_state['sketch'] = StreamScanner()
if 'total_cost' not in st.session_state: st.session_state['total_cost'] = 0.0
if 'total_tokens' not in st.session_state: st.session_state['total_tokens'] = 0
if 'authenticated' not in st.session_state: st.session_state['authenticated'] = False
if 'auth_error' not in st.session_state: st.session_state['auth_error'] = False
if 'ai_response' not in st.session_state: st.session_state['ai_response'] = ""
if 'last_sketch_hash' not in st.session_state: st.session_state['last_sketch_hash'] = None

def reset_sketch():
    st.session_state['sketch'] = StreamScanner()
    st.session_state['ai_response'] = ""
    st.session_state['last_sketch_hash'] = None
    gc.collect()

def perform_login():
    try:
        correct_password = get_auth_password()
        if secrets.compare_digest(st.session_state.password_input, correct_password):
            st.session_state['authenticated'] = True
            st.session_state['auth_error'] = False
            st.session_state['password_input'] = ""
        else:
            st.session_state['auth_error'] = True
    except Exception:
        st.session_state['auth_error'] = True

def logout():
    st.session_state['authenticated'] = False
    st.session_state['ai_response'] = ""


# 🛠️ helpers, setup
# ============================================

@st.cache_resource(show_spinner="Init NLTK...")
def setup_nlp_resources():
    if nltk is None: return None, None
    try: 
        nltk.data.find('sentiment/vader_lexicon.zip')
        nltk.data.find('corpora/wordnet.zip')
        nltk.data.find('corpora/omw-1.4.zip')
    except LookupError: 
        try:
            nltk.download('vader_lexicon')
            nltk.download('wordnet')
            nltk.download('omw-1.4')
        except:
            pass
    
    sia = SentimentIntensityAnalyzer()
    lemmatizer = WordNetLemmatizer()
    return sia, lemmatizer

@st.cache_data(show_spinner=False)
def list_system_fonts() -> Dict[str, str]:
    mapping = {}
    for fe in font_manager.fontManager.ttflist:
        if fe.name not in mapping: mapping[fe.name] = fe.fname
    return dict(sorted(mapping.items(), key=lambda x: x[0].lower()))

def build_punct_translation(keep_hyphens: bool, keep_apostrophes: bool) -> dict:
    # 1. Standard ASCII punctuation
    punct = string.punctuation
    
    # 2. Add Unicode "Smart" quotes & dashes
    punct += "“”‘’–—…" 

    if keep_hyphens: 
        for char in "-–—": punct = punct.replace(char, "")
    if keep_apostrophes: 
        for char in "'’": punct = punct.replace(char, "")
    
    return str.maketrans("", "", punct)

def parse_user_stopwords(raw: str) -> Tuple[List[str], List[str]]:
    raw = raw.replace("\n", ",").replace(".", ",")
    phrases, singles = [], []
    for item in [x.strip() for x in raw.split(",") if x.strip()]:
        if " " in item: phrases.append(item.lower())
        else: singles.append(item.lower())
    return phrases, singles

def default_prepositions() -> set:
    return {
        'about', 'above', 'across', 'after', 'against', 'along', 'among',
        'around', 'at', 'before', 'behind', 'below', 'beneath', 'beside',
        'between', 'beyond', 'but', 'by', 'concerning', 'despite', 'down',
        'during', 'except', 'for', 'from', 'in', 'inside', 'into', 'like',
        'near', 'of', 'off', 'on', 'onto', 'out', 'outside', 'over', 'past',
        'regarding', 'since', 'through', 'throughout', 'to', 'toward', 'under',
        'underneath', 'until', 'up', 'upon', 'with', 'within', 'without',
        'something', 'someone', 'somebody', 'anything', 'anyone', 'anybody',
        'everything', 'everyone', 'everybody', 'nothing', 'none', 'maybe',
        'perhaps', 'really', 'basically', 'generally', 'kind', 'sort', 'stuff',
        'thing', 'things'
    }

def build_phrase_pattern(phrases: List[str]) -> Optional[re.Pattern]:
    if not phrases:
        return None
    escaped = [re.escape(p) for p in phrases if p]
    if not escaped:
        return None
    return re.compile(rf"\b(?:{'|'.join(escaped)})\b", flags=re.IGNORECASE)

def normalize_speaker_name(raw: str) -> str:
    if not isinstance(raw, str):
        return ""
    text = raw.replace("\ufeff", " ").replace("\ufffd", " ")
    text = "".join(ch if ch >= " " else " " for ch in text)
    text = html.unescape(text).lower()
    text = re.sub(r"\s+", " ", text)
    return text.strip(" \t\r\n:-–—")


def parse_speaker_exclusions(raw: str) -> Set[str]:
    if not raw:
        return set()
    entries = raw.replace("\n", ",").split(",")
    return {
        normalized
        for item in entries
        if (normalized := normalize_speaker_name(item))
    }


def extract_speaker_label(text: str) -> Tuple[Optional[str], str]:
    """
    Extracts transcript-style speaker labels such as:
    - Speaker 18 (Blaze): text
    - Omar Akhtar: text

    Returns (speaker, utterance). If no safe label is found, speaker is None.
    """
    if not isinstance(text, str) or ":" not in text:
        return None, text

    speaker, utterance = text.split(":", 1)
    speaker = speaker.strip()
    if not speaker or len(speaker) > MAX_SPEAKER_NAME_LENGTH:
        return None, text
    if "-->" in speaker or speaker.isdigit():
        return None, text
    if not re.search(r"[A-Za-z]", speaker):
        return None, text

    return speaker, utterance.strip()


def speaker_is_excluded(
    speaker: Optional[str],
    excluded_speakers: Set[str],
    partial_match: bool = False
) -> bool:
    if not speaker or not excluded_speakers:
        return False

    normalized = normalize_speaker_name(speaker)
    if not normalized:
        return False
    if normalized in excluded_speakers:
        return True
    if partial_match:
        return any(entry in normalized or normalized in entry for entry in excluded_speakers)
    return False


def collect_speaker_labels_from_text(
    text: str,
    max_lines: int = 20000
) -> Counter:
    """
    Counts explicit transcript speaker labels without attempting identity detection.

    This only looks for labels already present before a colon, such as
    "Speaker 18 (Blaze):" or "Participant A:".
    """
    counts = Counter()
    if not isinstance(text, str):
        return counts

    for idx, line in enumerate(text.splitlines()):
        if idx >= max_lines:
            break
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit():
            continue
        speaker, _utterance = extract_speaker_label(line)
        if speaker:
            counts[speaker.strip()] += 1
    return counts


def collect_speaker_labels_from_file(
    file_bytes: bytes,
    filename: str = "",
    max_lines: int = 20000
) -> Counter:
    """
    Safely previews transcript-style speaker labels from uploaded text/VTT files.
    Does not inspect PDFs, Office files, CSVs, or JSON to avoid surprising work.
    """
    lower = (filename or "").lower()
    if not (lower.endswith((".txt", ".vtt")) or is_probably_vtt(file_bytes)):
        return Counter()

    try:
        text = file_bytes.decode(detect_text_encoding(file_bytes), errors="replace")
    except Exception:
        return Counter()
    return collect_speaker_labels_from_text(text, max_lines=max_lines)


@lru_cache(maxsize=1)
def collect_maturity_vocabulary() -> Set[str]:
    """
    Collects maturity-model vocabulary so maturity-critical words are not
    accidentally removed by generic stopword / filler-word cleanup.

    This intentionally includes:
    - explicit one-word terms
    - individual words from maturity phrases

    Example: "student success" protects both "student" and "success".

    Cached because maturity vocabularies are static during an app run and this
    function may be called every time scan settings are rebuilt.
    """
    protected_terms: Set[str] = set()

    try:
        assessor = MaturityAssessor()
        for model in assessor.models.values():
            if model.get("type") == "domain_based":
                domain_iter = model.get("domains", {}).values()
                for domain_data in domain_iter:
                    for tier_data in domain_data.get("tiers", {}).values():
                        for term in tier_data.get("terms", set()):
                            protected_terms.add(str(term).lower())
                        for phrase in tier_data.get("phrases", []):
                            protected_terms.update(
                                part.lower()
                                for part in str(phrase).split()
                                if part.strip()
                            )
            else:
                for level_data in model.get("levels", {}).values():
                    for term in level_data.get("terms", set()):
                        protected_terms.add(str(term).lower())
                    for phrase in level_data.get("phrases", []):
                        protected_terms.update(
                            part.lower()
                            for part in str(phrase).split()
                            if part.strip()
                        )
    except Exception:
        # Safety-first fallback: never break app startup because help/protection
        # vocabulary could not be collected.
        return set()

    return protected_terms


def protect_maturity_vocabulary(stopwords: Set[str]) -> Set[str]:
    """
    Removes maturity-critical terms from the active stopword set.

    This preserves maturity scoring while still allowing generic filler words
    and user-entered junk terms to clean the broader analysis.
    """
    maturity_terms = collect_maturity_vocabulary()
    if not maturity_terms:
        return stopwords
    return set(stopwords) - maturity_terms

def estimate_row_count_from_bytes(file_bytes: bytes) -> int:
    if not file_bytes: return 0
    return file_bytes.count(b'\n') + 1

def make_unique_header(raw_names: List[Optional[str]]) -> List[str]:
    seen: Dict[str, int] = {}
    result: List[str] = []
    for i, nm in enumerate(raw_names):
        name = (str(nm).strip() if nm is not None else "")
        if not name: name = f"col_{i}"
        if name in seen:
            seen[name] += 1
            unique = f"{name}__{seen[name]}"
        else:
            seen[name] = 1
            unique = name
        result.append(unique)
    return result

def extract_entities_regex(text: str, stopwords: Set[str]) -> List[str]:
    # lightweight NER without heavy models
    candidates = NER_CAPS_RE.findall(text)
    valid = []
    for c in candidates:
        # filter out if it's just a common stopword capitalized at start of sentence
        if c.lower() in stopwords: continue
        if len(c) < 3: continue
        valid.append(c)
    return valid

# --- virtual files and web/url
class VirtualFile:
    def __init__(self, name: str, text_content: str):
        self.name = name
        self._bytes = text_content.encode('utf-8')
    
    def getvalue(self) -> bytes:
        return self._bytes
    
    def getbuffer(self) -> memoryview:
        return memoryview(self._bytes)

def fetch_url_content(url: str) -> Optional[str]:
    if not requests or not BeautifulSoup: return None
    if not validate_url(url): return None
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        for script in soup(["script", "style", "nav", "footer"]):
            script.decompose()
        return soup.get_text(separator=' ', strip=True)
    except Exception: return None


# 📄 file readers (tuple yielding)
# ==========================================

# all readers yield (text_content, date_str, category_str)

def detect_text_encoding(file_bytes: bytes, encoding_choice: str = "auto") -> str:
    """
    Detect common transcript/text encodings without adding dependencies.

    Several meeting transcript exports are UTF-16 even when the file name ends
    in .txt. Reading those as UTF-8 leaves embedded NUL characters that render
    as square boxes in the word cloud and pollute downstream counts.
    """
    if encoding_choice == "latin-1":
        return "latin-1"
    if file_bytes.startswith((b"\xff\xfe", b"\xfe\xff")):
        return "utf-16"
    if file_bytes.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    return "utf-8"

def decode_text_sample(file_bytes: bytes, encoding_choice: str = "auto", max_bytes: int = 8192) -> str:
    enc = detect_text_encoding(file_bytes, encoding_choice)
    return file_bytes[:max_bytes].decode(enc, errors="replace")

def is_probably_vtt(file_bytes: bytes, encoding_choice: str = "auto") -> bool:
    try:
        sample = decode_text_sample(file_bytes, encoding_choice)
    except Exception:
        return False
    sample = sample.lstrip("\ufeff\r\n\t ")
    return sample.startswith("WEBVTT") or "-->" in sample[:1000]

def read_rows_raw_lines(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[Tuple[str, None, None]]:
    def _iter(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield (line.rstrip("\r\n"), None, None)
    try:
        yield from _iter(detect_text_encoding(file_bytes, encoding_choice))
    except UnicodeDecodeError:
        yield ("", None, None)

def read_rows_vtt(
    file_bytes: bytes,
    encoding_choice: str = "auto",
    excluded_speakers: Optional[Set[str]] = None,
    partial_speaker_match: bool = False
) -> Iterable[Tuple[str, None, None]]:
    # robust VTT reader that yields tuples
    def _iter_lines(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield line.rstrip("\r\n")
    
    iterator = _iter_lines(detect_text_encoding(file_bytes, encoding_choice))
    
    for line in iterator:
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit(): continue
        speaker, utterance = extract_speaker_label(line)
        if speaker:
            excluded = speaker_is_excluded(speaker, excluded_speakers or set(), partial_speaker_match)
            if excluded:
                continue
            if " " not in speaker:
                yield (line, None, None)
                continue
            yield (utterance, None, None)
            continue
        yield (line, None, None)

def read_rows_pdf(file_bytes: bytes) -> Iterable[Tuple[str, None, None]]:
    if pypdf is None: 
        st.error("pypdf missing")
        return
    bio = io.BytesIO(file_bytes)
    try:
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text = page.extract_text()
            if text: yield (text, None, None)
    except Exception: yield ("", None, None)

def read_rows_pptx(file_bytes: bytes) -> Iterable[Tuple[str, None, None]]:
    if pptx is None: return
    bio = io.BytesIO(file_bytes)
    try:
        prs = pptx.Presentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text: yield (shape.text, None, None)
    except Exception:
        yield ("", None, None)

def read_rows_json(file_bytes: bytes, selected_key: str = None) -> Iterable[Tuple[str, None, None]]:
    bio = io.BytesIO(file_bytes)
    try:
        wrapper = io.TextIOWrapper(bio, encoding="utf-8", errors="replace")
        for line in wrapper:
            if not line.strip(): continue
            try:
                obj = json.loads(line)
                txt = ""
                if selected_key and isinstance(obj, dict): txt = str(obj.get(selected_key, ""))
                elif isinstance(obj, str): txt = obj
                else: txt = str(obj)
                yield (txt, None, None)
            except json.JSONDecodeError:
                pass
    except Exception:
        pass

def read_rows_csv_structured(
    file_bytes: bytes, 
    encoding_choice: str, 
    delimiter: str, 
    has_header: bool, 
    text_cols: List[str],
    date_col: Optional[str],
    cat_col: Optional[str],
    join_with: str
) -> Iterable[Tuple[str, Optional[str], Optional[str]]]:
    
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    
    with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
        rdr = csv.reader(wrapper, delimiter=delimiter)
        first = next(rdr, None)
        if first is None: return

        if has_header:
            header = make_unique_header(list(first))
            name_to_idx = {n: i for i, n in enumerate(header)}
            
            text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
            date_idx = name_to_idx.get(date_col) if date_col else None
            cat_idx = name_to_idx.get(cat_col) if cat_col else None
        else:
            # if no header, user likely selected "col_0", "col_1" etc.
            name_to_idx = {f"col_{i}": i for i in range(len(first))}
            text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
            date_idx = name_to_idx.get(date_col) if date_col else None
            cat_idx = name_to_idx.get(cat_col) if cat_col else None
            
            # yield 1st row data
            txt_parts = [first[i] if i < len(first) else "" for i in text_idxs]
            d_val = first[date_idx] if (date_idx is not None and date_idx < len(first)) else None
            c_val = first[cat_idx] if (cat_idx is not None and cat_idx < len(first)) else None
            yield (join_with.join(txt_parts), d_val, c_val)

        for row in rdr:
            txt_parts = [row[i] if i < len(row) else "" for i in text_idxs]
            d_val = row[date_idx] if (date_idx is not None and date_idx < len(row)) else None
            c_val = row[cat_idx] if (cat_idx is not None and cat_idx < len(row)) else None
            yield (join_with.join(txt_parts), d_val, c_val)

def iter_excel_structured(
    file_bytes: bytes, 
    sheet_name: str, 
    has_header: bool, 
    text_cols: List[str],
    date_col: Optional[str],
    cat_col: Optional[str],
    join_with: str
) -> Iterable[Tuple[str, Optional[str], Optional[str]]]:
    if openpyxl is None: return
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    rows_iter = ws.iter_rows(values_only=True)
    
    first = next(rows_iter, None)
    if first is None: 
        wb.close()
        return

    # header logic
    if has_header:
        header = make_unique_header(list(first))
        name_to_idx = {n: i for i, n in enumerate(header)}
        text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
        date_idx = name_to_idx.get(date_col) if date_col else None
        cat_idx = name_to_idx.get(cat_col) if cat_col else None
    else:
        name_to_idx = {f"col_{i}": i for i in range(len(first))}
        text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
        date_idx = name_to_idx.get(date_col) if date_col else None
        cat_idx = name_to_idx.get(cat_col) if cat_col else None
        
        # Yield first row
        txt_parts = [str(first[i]) if (i < len(first) and first[i] is not None) else "" for i in text_idxs]
        d_val = first[date_idx] if (date_idx is not None and date_idx < len(first)) else None
        c_val = first[cat_idx] if (cat_idx is not None and cat_idx < len(first)) else None
        yield (join_with.join(txt_parts), str(d_val) if d_val else None, str(c_val) if c_val else None)

    for row in rows_iter:
        txt_parts = [str(row[i]) if (i < len(row) and row[i] is not None) else "" for i in text_idxs]
        d_val = row[date_idx] if (date_idx is not None and date_idx < len(row)) else None
        c_val = row[cat_idx] if (cat_idx is not None and cat_idx < len(row)) else None
        yield (join_with.join(txt_parts), str(d_val) if d_val else None, str(c_val) if c_val else None)
    
    wb.close()

def detect_csv_headers(file_bytes: bytes, delimiter: str = ",") -> List[str]:
    try:
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding="utf-8", errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return make_unique_header(row) if row else []
    except: return []

def detect_csv_num_cols(file_bytes: bytes, delimiter: str = ",") -> int:
    try:
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding="utf-8", errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return len(row) if row else 0
    except: return 0

def get_excel_sheetnames(file_bytes: bytes) -> List[str]:
    if openpyxl is None: return []
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    sheets = list(wb.sheetnames)
    wb.close()
    return sheets

def get_excel_preview(file_bytes: bytes, sheet_name: str, has_header: bool, rows: int = 5) -> pd.DataFrame:
    if openpyxl is None: return pd.DataFrame()
    bio = io.BytesIO(file_bytes)
    try:
        df = pd.read_excel(bio, sheet_name=sheet_name, header=0 if has_header else None, nrows=rows, engine='openpyxl')
        if not has_header: df.columns = [f"col_{i}" for i in range(len(df.columns))]
        return df
    except:
        return pd.DataFrame()

def excel_estimate_rows(file_bytes: bytes, sheet_name: str, has_header: bool) -> int:
    if openpyxl is None: return 0
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    total = ws.max_row or 0
    wb.close()
    if has_header and total > 0: total -= 1
    return max(total, 0)


# ⚙️ processing logic
# ==========================================
def clean_date_str(raw: Any) -> Optional[str]:
    """Robust extraction using dateutil if available, falling back to regex."""
    if not raw: return None
    s = str(raw).strip()
    
    # tries smart parsing (covers "Jan 5, 2024", "2024/01/01", etc)
    if date_parser:
        try:
            # fuzzy=True allows extracting dates buried in strings like "Date: Jan 1"
            dt = date_parser.parse(s, fuzzy=True)
            return dt.strftime("%Y-%m-%d")
        except (ValueError, OverflowError):
            pass # falls through to regex if dateutil fails

    # the fallback Regex (simple ISO & US formats)
    match = re.search(r'\d{4}-\d{2}-\d{2}', s)
    if match: return match.group(0)
    match_us = re.search(r'(\d{1,2})/(\d{1,2})/(\d{4})', s)
    if match_us:
        m, d, y = match_us.groups()
        return f"{y}-{int(m):02d}-{int(d):02d}"
    
    return None

def apply_text_cleaning(text: str, config: CleaningConfig) -> str:
    if not isinstance(text, str): return ""

    # Strip non-printing artifacts from badly decoded transcript/text files.
    # This prevents control glyphs from becoming word-cloud tokens.
    text = text.replace("\ufeff", " ").replace("\ufffd", " ")
    text = "".join(
        ch if (ch >= " " or ch in "\n\r\t") else " "
        for ch in text
    )
    
    # standard cleaning FIRST (to convert tags/artifacts to spaces)
    # ensures artifacts like <br> become spaces so the de-hyphenator can see them
    if config.remove_chat: text = CHAT_ARTIFACT_RE.sub(" ", text)
    if config.remove_html: text = HTML_TAG_RE.sub(" ", text)
    if config.unescape:
        try: text = html.unescape(text)
        except: pass

    # the de-hyphenation repair
    # merges words split by hyphen+whitespace (e.g., "equiv-\nalent" -> "equivalent")
    text = re.sub(r'(\w)-\s+(\w)', r'\1\2', text)
    
    # sticky Acronyms Fix (e.g., "WHOrecommended" -> "WHO recommended")
    # logic: 2+ CAPS followed by 2+ lowercase letters
    text = re.sub(r'([A-Z]{2,})([a-z]{2,})', r'\1 \2', text)

    # final cleanups
    if config.remove_urls: text = URL_EMAIL_RE.sub(" ", text)
    
    # normalizing
    text = text.lower()
    
    if config.phrase_pattern: text = config.phrase_pattern.sub(" ", text)
    return text.strip()

def process_chunk_iter(
    rows_iter: Iterable[Tuple[str, Optional[str], Optional[str]]],
    clean_conf: CleaningConfig,
    proc_conf: ProcessingConfig,
    scanner: StreamScanner,
    lemmatizer: Optional[WordNetLemmatizer],
    progress_cb: Optional[Callable[[int], None]] = None
):
    _min_len = proc_conf.min_word_len
    _drop_int = proc_conf.drop_integers
    _trans = proc_conf.translate_map
    _stopwords = proc_conf.stopwords
    _lemma = proc_conf.use_lemmatization and (lemmatizer is not None)
    _excluded_speakers = proc_conf.excluded_speakers
    _partial_speaker_match = proc_conf.partial_speaker_match
    
    # defines set of "edge junk" to strip (quotes, brackets, dashes)
    # to catch "word" or (word) or -word-
    _strip_chars = string.punctuation + "“”‘’–—" 
    
    local_global_counts = Counter()
    local_global_bigrams = Counter() if proc_conf.compute_bigrams else Counter()
    
    batch_accum = Counter()
    batch_rows = 0
    row_count = 0
    
    # pre-caching lemmatizer methods for speed loop
    lemmatize = lemmatizer.lemmatize if _lemma else None

    for (raw_text, date_val, cat_val) in rows_iter:
        row_count += 1

        if _excluded_speakers:
            speaker, _utterance = extract_speaker_label(raw_text)
            if speaker_is_excluded(speaker, _excluded_speakers, _partial_speaker_match):
                continue
        
        # entities (Before lowercase)
        if raw_text:
            entities = extract_entities_regex(raw_text, _stopwords)
            scanner.update_entities(entities)

        # cleaning
        text = apply_text_cleaning(raw_text, clean_conf)
        
        # tokenization & filter
        filtered_tokens_line: List[str] = []
        for t in text.split():
            # 1. Internal Translation (removes internal chars like don't -> dont)
            t = t.translate(_trans)
            
            # 2. Edge Stripping (The Fix for "between and operations")
            t = t.strip(_strip_chars)
            
            if not t: continue
            if _drop_int and t.isdigit(): continue
            if len(t) < _min_len: continue
            
            # lemmatize?
            if _lemma:
                t = lemmatize(t, pos='v')
                t = lemmatize(t, pos='n')
            
            if t in _stopwords: continue
            filtered_tokens_line.append(t)
        
        if filtered_tokens_line:
            # Stats Update
            line_counts = Counter(filtered_tokens_line)
            local_global_counts.update(filtered_tokens_line)
            
            if proc_conf.compute_bigrams and len(filtered_tokens_line) > 1:
                local_global_bigrams.update(pairwise(filtered_tokens_line))
            
            # metadata update
            clean_date = clean_date_str(date_val)
            clean_cat = str(cat_val).strip() if cat_val else None
            scanner.update_metadata_stats(clean_date, clean_cat, filtered_tokens_line)
            scanner.add_evidence_doc(raw_text, filtered_tokens_line, clean_date, clean_cat)

            # topic modeling batching
            batch_accum.update(line_counts)
            batch_rows += 1
            if batch_rows >= scanner.DOC_BATCH_SIZE:
                scanner.add_topic_sample(batch_accum)
                batch_accum = Counter()
                batch_rows = 0

        if progress_cb and (row_count % 2000 == 0): progress_cb(row_count)

    # flushing last batch
    if batch_accum and batch_rows > 0:
        scanner.add_topic_sample(batch_accum)

    scanner.update_global_stats(local_global_counts, local_global_bigrams, row_count)
    if progress_cb: progress_cb(row_count)
    gc.collect()

def perform_refinery_job(file_obj, chunk_size, clean_conf: CleaningConfig):
    with tempfile.TemporaryDirectory() as temp_dir:
        original_name = os.path.splitext(file_obj.name)[0]
        status_container = st.status(f"⚙️ Refining {file_obj.name}...", expanded=True)
        part_num = 1
        created_files = []
        
        try:
            file_obj.seek(0)
            df_iterator = pd.read_csv(file_obj, chunksize=chunk_size, on_bad_lines='skip', dtype=str)
            
            for chunk in df_iterator:
                for col in chunk.columns:
                    chunk[col] = chunk[col].fillna("")
                    chunk[col] = chunk[col].apply(lambda x: apply_text_cleaning(x, clean_conf))
                
                new_filename = f"{original_name}_cleaned_part_{part_num}.csv"
                temp_path = os.path.join(temp_dir, new_filename)
                chunk.to_csv(temp_path, index=False)
                created_files.append(temp_path)
                status_container.write(f"✅ Processed chunk {part_num} ({len(chunk)} rows)")
                part_num += 1
            
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                for file_path in created_files:
                    zip_file.write(file_path, arcname=os.path.basename(file_path))
            
            zip_buffer.seek(0)
            status_container.update(label="🎉 Refinery Job Complete!", state="complete", expanded=False)
            return zip_buffer
            
        except Exception as e:
            status_container.update(label="❌ Error", state="error")
            st.error(f"Refinery Error: {str(e)}")
            return None


# 📊 UI, analytics renderers
# ==========================================

def calculate_text_stats(counts: Counter, total_rows: int) -> Dict:
    total_tokens = sum(counts.values())
    unique_tokens = len(counts)
    avg_len = sum(len(word) * count for word, count in counts.items()) / total_tokens if total_tokens else 0
    return {
        "Total Rows": total_rows, "Total Tokens": total_tokens,
        "Unique Vocabulary": unique_tokens, "Avg Word Length": round(avg_len, 2),
        "Lexical Diversity": round(unique_tokens / total_tokens, 4) if total_tokens else 0
    }

def calculate_npmi(bigram_counts: Counter, unigram_counts: Counter, total_words: int, min_freq: int = 3) -> pd.DataFrame:
    results = []
    if not bigram_counts: return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
    epsilon = 1e-10 
    for (w1, w2), freq in bigram_counts.items():
        if freq < min_freq: continue
        count_w1 = unigram_counts.get(w1, 0)
        count_w2 = unigram_counts.get(w2, 0)
        if count_w1 == 0 or count_w2 == 0: continue
        prob_bigram = freq / total_words
        try: 
            pmi = math.log(prob_bigram / ((count_w1 / total_words) * (count_w2 / total_words)))
        except ValueError: continue
        
        log_prob_bigram = math.log(prob_bigram)
        if abs(log_prob_bigram) < epsilon: npmi = 1.0
        else: npmi = pmi / -log_prob_bigram
        results.append({"Bigram": f"{w1} {w2}", "Count": freq, "NPMI": round(npmi, 3)})
        
    df = pd.DataFrame(results)
    if df.empty: return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
    return df.sort_values("NPMI", ascending=False)

def calculate_tfidf(scanner: StreamScanner, top_n=50) -> pd.DataFrame:
    # IDF = log(total docs / doc freq)
    # TF (Global) = total count / total words (simplified for sketch)
    total_docs = len(scanner.topic_docs)
    if total_docs == 0: return pd.DataFrame()
    
    results = []
    # Analyze only terms that appear in at least 2 docs to avoid noise
    candidates = [t for t, c in scanner.doc_freqs.items() if c > 1]
    
    for term in candidates:
        tf = scanner.global_counts[term]
        df = scanner.doc_freqs[term]
        idf = math.log(total_docs / (1 + df))
        score = tf * idf
        results.append({"Term": term, "TF (Count)": tf, "DF (Docs)": df, "Keyphrase Score": round(score, 2)})
        
    df = pd.DataFrame(results)
    if df.empty: return df
    return df.sort_values("Keyphrase Score", ascending=False).head(top_n)

def perform_topic_modeling(synthetic_docs: List[Counter], n_topics: int, model_type: str) -> Optional[List[Dict]]:
    if not DictVectorizer or len(synthetic_docs) < 1: return None
    vectorizer = DictVectorizer(sparse=True)
    dtm = vectorizer.fit_transform(synthetic_docs)
    n_samples, n_features = dtm.shape
    if n_samples == 0 or n_features == 0: return None
    
    safe_n_topics = min(n_topics, min(n_samples, n_features)) if model_type == "NMF" else min(n_topics, n_samples)
    if safe_n_topics < 1: return None

    model = None
    try:
        if model_type == "LDA": model = LatentDirichletAllocation(n_components=safe_n_topics, random_state=42, max_iter=10)
        elif model_type == "NMF": model = NMF(n_components=safe_n_topics, random_state=42, init='nndsvd')
        model.fit(dtm)
    except ValueError: return None
    
    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(model.components_):
        top_indices = topic.argsort()[:-7:-1]
        top_words = [feature_names[i] for i in top_indices]
        strength = sum(topic[i] for i in top_indices)
        topics.append({"id": topic_idx + 1, "words": top_words, "strength": strength})
    return topics

def perform_bayesian_sentiment_analysis(counts: Counter, sentiments: Dict[str, float], pos_thresh: float, neg_thresh: float) -> Optional[Dict]:
    if not beta_dist: return None
    pos_count = sum(counts[w] for w, s in sentiments.items() if s >= pos_thresh)
    neg_count = sum(counts[w] for w, s in sentiments.items() if s <= neg_thresh)
    total_informative = pos_count + neg_count
    if total_informative < 1: return None

    alpha_post = 1 + pos_count
    beta_post = 1 + neg_count
    mean_prob = alpha_post / (alpha_post + beta_post)
    lower_ci, upper_ci = beta_dist.ppf([0.025, 0.975], alpha_post, beta_post)
    x = np.linspace(0, 1, 300)
    y = beta_dist.pdf(x, alpha_post, beta_post)
    return {
        "pos_count": pos_count, "neg_count": neg_count, "total": total_informative,
        "mean_prob": mean_prob, "ci_low": lower_ci, "ci_high": upper_ci,
        "x_axis": x, "pdf_y": y
    }

def build_theme_evidence_cards(scanner: StreamScanner, counts: Counter, top_n: int = 6) -> pd.DataFrame:
    total_words = max(sum(counts.values()), 1)
    npmi_df = calculate_npmi(scanner.global_bigrams, counts, total_words)
    tfidf_df = calculate_tfidf(scanner, 100)
    tfidf_scores = {
        row["Term"]: row["Keyphrase Score"]
        for _, row in tfidf_df.iterrows()
    } if not tfidf_df.empty else {}

    rows = []
    used_terms = set()

    if not npmi_df.empty:
        for _, row in npmi_df.head(top_n).iterrows():
            phrase = str(row["Bigram"])
            terms = phrase.split()
            used_terms.update(terms)
            related = Counter()
            for bg, freq in scanner.global_bigrams.items():
                if any(term in bg for term in terms):
                    for term in bg:
                        if term not in terms and term in counts:
                            related[term] += freq
            rows.append({
                "Theme Evidence": phrase,
                "Evidence Type": "Sticky phrase",
                "Support": int(row["Count"]),
                "Distinctiveness": float(row["NPMI"]),
                "Related Terms": ", ".join([t for t, _ in related.most_common(5)]),
                "Read As": "Terms that appear together more often than chance; useful as a candidate theme or concept.",
            })

    for term, score in sorted(tfidf_scores.items(), key=lambda item: item[1], reverse=True):
        if len(rows) >= top_n:
            break
        if term in used_terms or term not in counts:
            continue
        related = Counter()
        for bg, freq in scanner.global_bigrams.items():
            if term in bg:
                for neighbor in bg:
                    if neighbor != term and neighbor in counts:
                        related[neighbor] += freq
        rows.append({
            "Theme Evidence": term,
            "Evidence Type": "Distinctive term",
            "Support": counts[term],
            "Distinctiveness": round(score, 3),
            "Related Terms": ", ".join([t for t, _ in related.most_common(5)]),
            "Read As": "A corpus-specific word that may point to a distinctive theme, issue, or domain signal.",
        })

    return pd.DataFrame(rows)


def build_signal_quadrant_df(scanner: StreamScanner, counts: Counter, top_n: int = 150) -> pd.DataFrame:
    tfidf_df = calculate_tfidf(scanner, top_n)
    tfidf_scores = {
        row["Term"]: row["Keyphrase Score"]
        for _, row in tfidf_df.iterrows()
    } if not tfidf_df.empty else {}

    rows = []
    for term, count in counts.most_common(top_n):
        rows.append({
            "Term": term,
            "Frequency": count,
            "Distinctiveness": float(tfidf_scores.get(term, 0.0)),
        })
    df = pd.DataFrame(rows)
    if df.empty:
        return df

    freq_cutoff = df["Frequency"].median()
    distinct_cutoff = df["Distinctiveness"].median()

    def classify(row):
        high_freq = row["Frequency"] >= freq_cutoff
        high_distinct = row["Distinctiveness"] >= distinct_cutoff
        if high_freq and high_distinct:
            return "Core signal"
        if high_freq and not high_distinct:
            return "Common backdrop"
        if not high_freq and high_distinct:
            return "Niche signal"
        return "Low evidence"

    df["Quadrant"] = df.apply(classify, axis=1)
    return df.sort_values(["Quadrant", "Frequency"], ascending=[True, False])


def parse_expected_terms(raw: str) -> List[str]:
    if not raw:
        return []
    entries = raw.replace("\n", ",").split(",")
    return [entry.strip().lower() for entry in entries if entry.strip()]


def build_expected_terms_df(raw_terms: str, counts: Counter, bigrams: Counter) -> pd.DataFrame:
    rows = []
    for term in parse_expected_terms(raw_terms):
        parts = term.split()
        if len(parts) == 1:
            observed = counts.get(term, 0)
        elif len(parts) == 2:
            observed = bigrams.get((parts[0], parts[1]), 0) + bigrams.get((parts[1], parts[0]), 0)
        else:
            observed = 0

        if observed == 0:
            status = "Missing"
        elif observed < 3:
            status = "Weak"
        else:
            status = "Present"

        rows.append({
            "Expected Signal": term,
            "Observed Count": observed,
            "Status": status,
        })
    return pd.DataFrame(rows)


def compare_counter_terms(left_counts: Counter, right_counts: Counter, stopwords: Set[str], min_word_len: int, top_n: int = 50) -> pd.DataFrame:
    left_total = max(sum(left_counts.values()), 1)
    right_total = max(sum(right_counts.values()), 1)
    vocab = set(left_counts) | set(right_counts)
    rows = []
    for term in vocab:
        if len(str(term)) < min_word_len or term in stopwords:
            continue
        left_rate = left_counts.get(term, 0) / left_total
        right_rate = right_counts.get(term, 0) / right_total
        diff = left_rate - right_rate
        if left_counts.get(term, 0) + right_counts.get(term, 0) < 3:
            continue
        rows.append({
            "Term": term,
            "Left Count": left_counts.get(term, 0),
            "Right Count": right_counts.get(term, 0),
            "Left Rate": round(left_rate, 6),
            "Right Rate": round(right_rate, 6),
            "Difference": round(diff, 6),
            "Leans Toward": "Left" if diff > 0 else "Right",
        })
    df = pd.DataFrame(rows)
    if df.empty:
        return df
    df["Abs Difference"] = df["Difference"].abs()
    return df.sort_values("Abs Difference", ascending=False).drop(columns=["Abs Difference"]).head(top_n)


def build_temporal_drift_df(temporal_counts: Dict[str, Counter], stopwords: Set[str], min_word_len: int, top_n: int = 50) -> pd.DataFrame:
    dates = sorted(temporal_counts.keys())
    if len(dates) < 2:
        return pd.DataFrame()
    midpoint = max(1, len(dates) // 2)
    early = Counter()
    late = Counter()
    for d in dates[:midpoint]:
        early.update(temporal_counts[d])
    for d in dates[midpoint:]:
        late.update(temporal_counts[d])
    df = compare_counter_terms(late, early, stopwords, min_word_len, top_n)
    if df.empty:
        return df
    return df.rename(columns={
        "Left Count": "Late Count",
        "Right Count": "Early Count",
        "Left Rate": "Late Rate",
        "Right Rate": "Early Rate",
    }).replace({"Leans Toward": {"Left": "Rising later", "Right": "Fading later"}})


SIGNAL_TAXONOMY = {
    "Pain / Friction": {
        "pain", "problem", "issue", "difficult", "hard", "frustrat", "confusing",
        "slow", "manual", "burden", "struggle", "complaint", "broken", "stuck",
        "workaround", "delay", "bottleneck", "overwhelm",
    },
    "Need / Request": {
        "need", "needs", "want", "request", "require", "must", "should", "ask",
        "support", "training", "help", "guidance", "resource", "clarity",
        "documentation", "template", "example",
    },
    "Blocker / Constraint": {
        "block", "blocked", "barrier", "constraint", "limited", "lack", "missing",
        "capacity", "budget", "time", "access", "permission", "dependency",
        "approval", "technical", "policy", "risk",
    },
    "Aspiration / Opportunity": {
        "opportunity", "improve", "better", "future", "strategy", "strategic",
        "scale", "innovation", "transform", "optimize", "automate", "pilot",
        "enable", "modernize", "growth", "mature",
    },
    "Risk / Concern": {
        "risk", "concern", "security", "privacy", "compliance", "audit",
        "liability", "failure", "fragile", "exposure", "uncertain", "threat",
        "vulnerable", "sensitive",
    },
    "Decision / Tradeoff": {
        "decision", "decide", "choice", "tradeoff", "priority", "prioritize",
        "option", "recommend", "roadmap", "scope", "governance", "ownership",
        "accountable", "alignment",
    },
    "Contradiction / Tension": {
        "but", "however", "although", "yet", "except", "despite", "instead",
        "nevertheless", "while", "whereas", "tension", "conflict", "gap",
        "misalignment",
    },
}

TENSION_PATTERNS = [
    re.compile(r"\b(?:but|however|although|yet|despite|whereas|while)\b", re.IGNORECASE),
    re.compile(r"\b(?:want|need|should|goal|strategy)\b.{0,80}\b(?:but|however|blocked|lack|limited|constraint|barrier)\b", re.IGNORECASE),
    re.compile(r"\b(?:gap|misalignment|disconnect|tradeoff|workaround|bottleneck)\b", re.IGNORECASE),
]


def classify_signal_dimension(text: str, tokens: Iterable[str]) -> Tuple[str, int]:
    token_set = set(tokens)
    lowered = text.lower()
    scores = {}

    for label, lexicon in SIGNAL_TAXONOMY.items():
        score = 0
        for marker in lexicon:
            if marker in token_set or marker in lowered:
                score += 1
        scores[label] = score

    best_label, best_score = max(scores.items(), key=lambda item: item[1])
    if best_score == 0:
        return "Theme / Topic", 0
    return best_label, best_score


SEMANTIC_SIGNAL_FAMILIES = {
    "Infrastructure / System Dependence": {
        "markers": {
            "system", "machine", "apparatus", "infrastructure", "platform", "tool",
            "server", "cloud", "model", "automation", "automated", "workflow",
            "central", "power", "station", "network", "dependency", "depends",
            "depend", "mending", "repair", "service", "provider", "feeds",
            "clothes", "houses", "connects", "communication-system",
        },
        "interpretation": (
            "This points to reliance on an enabling system or infrastructure. "
            "The important question is whether the system is simply supporting the work, "
            "or quietly shaping what people can notice, do, or decide."
        ),
        "question": "What breaks, changes, or becomes impossible if this system stops working?",
    },
    "Embodiment / Lived Experience": {
        "markers": {
            "body", "bodies", "physical", "touch", "touched", "hand", "hands",
            "walk", "walking", "muscles", "breath", "breathe", "face",
            "voice", "near", "far", "felt", "feel", "pain", "experience",
            "direct", "silence", "hearing", "see", "seen", "saw", "steady",
            "stumbled", "pressed", "button",
        },
        "interpretation": (
            "This suggests a lived or embodied dimension: how people physically experience "
            "the situation, not just what they say about it. These signals often reveal what "
            "is lost, avoided, constrained, or newly felt."
        ),
        "question": "What lived experience is being revealed beneath the surface language?",
    },
    "Isolation / Disconnection": {
        "markers": {
            "alone", "isolated", "isolation", "separate", "separated", "remote",
            "distance", "distant", "room", "cell", "silence", "communicate",
            "communicated", "communication", "message", "visit", "visits",
            "connection", "disconnect", "lonely", "apart", "never", "correspondents",
        },
        "interpretation": (
            "This points to separation, mediated contact, or weakened connection. The signal "
            "may matter even when the text does not literally say 'isolation'."
        ),
        "question": "Who or what is separated here, and what form of connection is missing or mediated?",
    },
    "Authority / Legitimacy": {
        "markers": {
            "authority", "authorized", "governance", "policy", "rule", "rules",
            "committee", "approval", "book", "law", "duties", "ownership",
            "accountable", "ritual", "worship", "prayer", "divine", "blessed",
            "praise", "official", "leader", "management", "government",
            "constitution", "constitutional", "magistrate", "supreme", "court",
            "judicial", "executive", "legislative", "sovereignty", "legitimate",
        },
        "interpretation": (
            "This suggests a legitimacy or authority structure: rules, committees, rituals, "
            "official texts, or accepted narratives that tell people what is normal or permissible."
        ),
        "question": "What authority is being accepted, questioned, or protected by this language?",
    },
    "Decision / Tradeoff": {
        "markers": {
            "decision", "decide", "choice", "tradeoff", "trade-off", "balance",
            "compromise", "alternative", "option", "taxation", "taxes", "revenue",
            "imports", "exports", "commerce", "regulation", "power", "powers",
            "liberty", "security", "representation", "majority", "minority",
        },
        "interpretation": (
            "This points to a decision or tradeoff structure: competing goods, powers, "
            "constraints, or institutional choices that the text is asking readers to weigh."
        ),
        "question": "What tradeoff or institutional choice is being framed here?",
    },
    "Risk / Concern": {
        "markers": {
            "risk", "concern", "danger", "dangerous", "threat", "war", "armies",
            "standing", "faction", "tyranny", "disorder", "instability",
            "violence", "invasion", "security", "defense", "defence", "usurpation",
            "corruption", "jealousy", "ambition", "oppression", "disease",
            "infection", "infected", "plague", "bacillus", "pathogen", "hazard",
        },
        "interpretation": (
            "This is risk or concern language. The text is marking what could go wrong, "
            "what must be guarded against, or what danger justifies a proposed structure."
        ),
        "question": "What danger or failure is this language trying to prevent?",
    },
    "Disease / Hazard": {
        "markers": {
            "disease", "infection", "infected", "infectious", "plague", "bacillus",
            "bacteria", "bacterium", "virus", "typhi", "murium", "trypanosoma",
            "trypanozoon", "parasite", "pathogen", "contagion", "epidemic",
            "mortality", "cases", "symptoms", "rats", "mice", "rodents",
        },
        "interpretation": (
            "This points to a hazard or biological agent: disease, infection, organisms, "
            "vectors, or exposed populations. The important question is what risk is being "
            "tracked and how it spreads."
        ),
        "question": "What hazard, vector, or biological agent is being tracked here?",
    },
    "Evidence / Experiment": {
        "markers": {
            "experiment", "experiments", "experimental", "test", "tests", "tested",
            "observed", "observation", "sample", "samples", "specimen", "specimens",
            "survey", "biological", "guinea", "pigs", "laboratory", "inoculated",
            "culture", "cultures", "control", "controls", "results", "evidence",
            "data", "method", "methods", "table", "fig", "figure",
        },
        "interpretation": (
            "This looks like evidence or experiment language: observations, test subjects, "
            "samples, methods, or findings used to support a technical claim."
        ),
        "question": "What evidence, observation, or test is supporting the claim?",
    },
    "Intervention / Control Method": {
        "markers": {
            "control", "controlled", "extermination", "exterminate", "destruction",
            "destroy", "destroyed", "killed", "poison", "poisoning", "trap", "traps",
            "fumigation", "bisulphide", "carbon", "sanitation", "inspection",
            "prevention", "prevent", "eradication", "campaign", "measures",
            "treatment", "intervention", "removal", "disinfection",
        },
        "interpretation": (
            "This points to an intervention or control method: a practical action meant to "
            "reduce risk, interrupt spread, or change conditions."
        ),
        "question": "What intervention or control method is being proposed, tested, or relied on?",
    },
    "Public Health / Institutional Response": {
        "markers": {
            "public", "health", "service", "commission", "hospital", "marinehospital",
            "marine-hospital", "department", "bureau", "board", "officer", "surgeon",
            "asst", "assistant", "inspection", "quarantine", "sanitary", "sanitation",
            "municipal", "federal", "state", "report", "regulation", "authority",
        },
        "interpretation": (
            "This points to an organized public-health or institutional response: agencies, "
            "services, commissions, inspections, rules, or coordinated prevention work."
        ),
        "question": "Which institution or public-health response is being described, and what role does it play?",
    },
    "Risk / Failure Mode": {
        "markers": {
            "risk", "failure", "fail", "failed", "fragile", "threat", "danger",
            "collapse", "stops", "stopped", "broken", "defect", "repair", "mending",
            "darkness", "foul", "dying", "death", "complaints", "remedies",
            "vulnerable", "outage", "crisis", "warning", "alarm",
        },
        "interpretation": (
            "This looks like a failure-mode signal: the text is exposing what can break, "
            "what people depend on, or what becomes dangerous when normal supports degrade."
        ),
        "question": "What failure mode is being signaled, and who would feel it first?",
    },
    "Standardization / Loss of Difference": {
        "markers": {
            "alike", "same", "identical", "uniform", "standard", "standardized",
            "template", "common", "conform", "homogeneous", "similar", "everywhere",
            "ordinary", "generic", "repeated", "routine", "exactly", "flattened",
            "interchangeable",
        },
        "interpretation": (
            "This suggests standardization or loss of meaningful difference. In many corpora, "
            "this can point to efficiency, conformity, commoditization, or reduced local context."
        ),
        "question": "What difference or local context is being flattened here?",
    },
    "Institutional Structure / Social Design": {
        "markers": {
            "institution", "institutional", "public", "nurseries", "parents",
            "children", "duties", "assigned", "committee", "school", "organization",
            "team", "role", "roles", "department", "process", "stakeholder",
            "governance", "structure", "engagements", "gatherings",
            "federal", "national", "state", "states", "government", "governments",
            "representatives", "senate", "house", "congress", "union",
            "confederacy", "republic", "legislature", "legislatures",
            "department", "court", "jury", "convention",
        },
        "interpretation": (
            "This points to a social or institutional arrangement: how roles, duties, access, "
            "or relationships are organized by a system rather than by individual choice alone."
        ),
        "question": "What social arrangement is this language normalizing or exposing?",
    },
    "Aspiration / Ideology": {
        "markers": {
            "advanced", "progress", "future", "improve", "better", "innovation",
            "modern", "transformation", "ideal", "belief", "believed", "praise",
            "thanks", "blessed", "vision", "strategy", "promise", "hope",
        },
        "interpretation": (
            "This may be aspirational or ideological language. It can reveal what the text "
            "presents as progress, virtue, inevitability, or shared belief."
        ),
        "question": "What idea of progress or value is being promoted, assumed, or challenged?",
    },
    "Motif / Image Pattern": {
        "markers": {
            "colour", "color", "pearl", "cloud", "hills", "stars", "light",
            "dark", "darkness", "sun", "moon", "sky", "blue", "white",
            "surface", "earth", "hemisphere", "northern", "southern",
            "vision", "visions", "image", "symbol", "metaphor",
        },
        "interpretation": (
            "This appears to be recurring image or motif language. It may matter as atmosphere, "
            "symbolism, or framing, but it should not be treated as a practical need or conclusion "
            "unless the evidence ties it to a stronger pattern."
        ),
        "question": "Is this image pattern carrying interpretive weight, or is it mainly descriptive texture?",
    },
}

LOW_INFORMATION_SIGNAL_TERMS = {
    "case", "either", "thing", "things", "said", "will", "still", "time", "made",
    "make", "must", "really", "perhaps", "course", "however", "therefore",
    "seemed", "whose", "whereunder", "diurnal", "passed",
}

GENERIC_SIGNAL_PHRASES = {
    "either case",
    "in either case",
    "of course",
    "said vashti",
}

REQUIRED_SEMANTIC_MARKERS = {
    "Embodiment / Lived Experience": {
        "body", "bodies", "physical", "touch", "touched", "hand", "hands",
        "walk", "walking", "muscles", "breath", "breathe", "face", "voice",
        "felt", "feel", "pain", "experience", "direct", "silence", "hearing",
        "steady", "stumbled", "pressed", "button",
    },
    "Standardization / Loss of Difference": {
        "alike", "same", "identical", "uniform", "standard", "standardized",
        "homogeneous", "everywhere", "exactly", "interchangeable",
    },
    "Motif / Image Pattern": {
        "colour", "color", "pearl", "cloud", "hills", "stars", "light",
        "dark", "darkness", "sun", "moon", "sky", "blue", "white",
        "vision", "visions", "image", "symbol", "metaphor",
    },
}

SIGNAL_TYPE_PRIORITY = {
    "Infrastructure / System Dependence": 1.35,
    "Risk / Failure Mode": 1.3,
    "Institutional Structure / Social Design": 1.25,
    "Standardization / Loss of Difference": 1.2,
    "Authority / Legitimacy": 1.15,
    "Decision / Tradeoff": 1.15,
    "Risk / Concern": 1.1,
    "Disease / Hazard": 1.25,
    "Evidence / Experiment": 1.1,
    "Intervention / Control Method": 1.2,
    "Public Health / Institutional Response": 1.2,
    "Isolation / Disconnection": 1.1,
    "Embodiment / Lived Experience": 1.0,
    "Aspiration / Ideology": 0.95,
    "Contradiction / Tension": 0.75,
    "Motif / Image Pattern": 0.65,
    "Source / Boilerplate": 0.15,
    "Low-Specificity Signal": 0.25,
    "Absence / Weak Signal": 0.2,
}

SIGNAL_ROLE_PRIORITY = {
    "Core Insight": 4,
    "Supporting Signal": 3,
    "Supporting Motif": 2,
    "Context / Reference": 1,
    "Low-Specificity": 0,
}

CONTEXT_REFERENCE_TERMS = {
    "hemisphere", "northern", "southern", "eastern", "western", "north",
    "south", "east", "west", "brisbane", "peking", "shrewsbury", "sumatra",
    "island", "school", "epoch", "century", "chapter", "page", "volume",
    "city", "country", "region", "place", "location", "french", "revolution",
    "mongolian", "mongols", "australian", "brazil", "courland", "wessex",
    "rhode", "carolina", "virginia", "maryland", "delaware", "pennsylvania",
    "jersey", "connecticut", "georgia", "york", "mclean", "mcleans",
}

BOILERPLATE_SIGNAL_TERMS = {
    "edition", "editions", "copyright", "license", "licence", "gutenberg",
    "transcriber", "publisher", "published", "publication", "printer",
    "printing", "press", "packet", "journal", "newspaper", "magazine",
    "advertiser", "gazette", "edition", "mclean", "mcleans",
    "appendix", "preface", "foreword", "introduction", "contents", "index",
    "chapter", "section", "volume", "vol", "page", "pages", "continued",
    "download", "ebook", "archive", "source", "scan", "scanned", "ocr",
    "proofread", "proofreading", "editor", "editors", "errata", "title",
    "heading", "header", "footer", "rights", "reserved",
    "monday", "tuesday", "wednesday", "thursday", "friday", "saturday",
    "sunday", "january", "february", "march", "april", "june", "july",
    "august", "september", "october", "november", "december",
}

BOILERPLATE_SIGNAL_PHRASES = {
    "subject continued",
    "all rights",
    "rights reserved",
    "project gutenberg",
    "daily advertiser",
    "independent journal",
    "york packet",
    "mcleans edition",
    "state york",
    "tuesday february",
}

BOILERPLATE_SIGNAL_PATTERNS = [
    re.compile(
        r"\b(?:edition|editions|copyright|license|licence|gutenberg|transcriber|"
        r"publisher|published|publication|printer|printing|packet|journal|"
        r"newspaper|magazine|advertiser|gazette|mcleans?|appendix|preface|foreword|contents|index|"
        r"chapter|section|volume|page|continued|ebook|archive|ocr|proofread|"
        r"proofreading|errata|header|footer|rights\s+reserved|"
        r"monday|tuesday|wednesday|thursday|friday|saturday|sunday|"
        r"january|february|march|april|june|july|august|september|october|november|december)\b",
        re.IGNORECASE,
    ),
]

FRAGMENT_START_TERMS = {
    "pressed", "made", "awoke", "said", "called", "went", "came", "took",
    "began", "became", "seemed", "found", "turned", "looked", "asked",
}

FRAGMENT_END_TERMS = {
    "unfamiliar", "little", "another", "however", "therefore", "perhaps",
    "again", "soon", "early", "late", "said", "cried", "asked", "told",
}

GENERIC_BIGRAM_END_TERMS = {
    "said", "cried", "asked", "told", "made", "went", "came", "looked",
    "seemed", "became", "found",
}

ABSTRACT_FRAGMENT_TERMS = {
    "constitute", "constitutes", "constituted", "ingredients", "ingredient",
    "circumstances", "considerations", "objects", "object", "things",
    "means", "ends", "parts", "particulars", "principles",
}

PATTERN_BASED_FAMILY_RULES = [
    (
        re.compile(
            r"\b(?:typhi\s+murium|trypanosoma|trypanozoon|bacillus|plague\s+infection|"
            r"plague\s+rats|indian\s+plague|danysz\s+virus|disease|infected|infection)\b",
            re.IGNORECASE,
        ),
        "Disease / Hazard",
        8,
    ),
    (
        re.compile(
            r"\b(?:guinea\s+pigs|biological\s+survey|white\s+rats|laboratory|"
            r"inoculated|experiment|experiments|specimens?|samples?|observed|culture)\b",
            re.IGNORECASE,
        ),
        "Evidence / Experiment",
        8,
    ),
    (
        re.compile(
            r"\b(?:destruction\s+rats|extermination\s+rats|destroy\s+rats|"
            r"rats\s+killed|carbon\s+bisulphide|poison|traps?|fumigation|"
            r"sanitation|disinfection|prevention|eradication)\b",
            re.IGNORECASE,
        ),
        "Intervention / Control Method",
        8,
    ),
    (
        re.compile(
            r"\b(?:public\s+health|plague\s+commission|marinehospital\s+service|"
            r"marine-hospital\s+service|health\s+service|sanitary\s+board|"
            r"asst\s+surg|assistant\s+surgeon|quarantine|inspection)\b",
            re.IGNORECASE,
        ),
        "Public Health / Institutional Response",
        8,
    ),
    (
        re.compile(
            r"\b(?:federal\s+government|national\s+government|state\s+governments?|"
            r"house\s+representatives|trial\s+jury|state\s+legislatures?|"
            r"executive\s+department|legislative\s+body|supreme\s+court|"
            r"proposed\s+constitution|constitution|union|confederacy|republic)\b",
            re.IGNORECASE,
        ),
        "Institutional Structure / Social Design",
        8,
    ),
    (
        re.compile(
            r"\b(?:chief\s+magistrate|supreme\s+court|judicial\s+power|"
            r"executive\s+power|legislative\s+authority|constitutional\s+authority|"
            r"authority|sovereignty)\b",
            re.IGNORECASE,
        ),
        "Authority / Legitimacy",
        8,
    ),
    (
        re.compile(
            r"\b(?:standing\s+armies|faction|tyranny|invasion|war|security|"
            r"public\s+danger|internal\s+violence|usurpation)\b",
            re.IGNORECASE,
        ),
        "Risk / Concern",
        8,
    ),
    (
        re.compile(
            r"\b(?:imports\s+exports|power\s+taxation|taxation|commerce|"
            r"liberty\s+security|checks?\s+balances?|choice|tradeoff)\b",
            re.IGNORECASE,
        ),
        "Decision / Tradeoff",
        8,
    ),
    (
        re.compile(
            r"\b(?:mending\s+apparatus|power\s+station|communication-system|central\s+power|infrastructure)\b",
            re.IGNORECASE,
        ),
        "Infrastructure / System Dependence",
        7,
    ),
    (
        re.compile(
            r"\b(?:public\s+nurseries|nurseries|parents?|children|duties|assigned|institutional|social\s+design)\b",
            re.IGNORECASE,
        ),
        "Institutional Structure / Social Design",
        7,
    ),
    (
        re.compile(
            r"\b(?:advanced\s+thanks|thanks\s+machine|blessed\s+machine|progress|praise|worship|belief|ideology)\b",
            re.IGNORECASE,
        ),
        "Aspiration / Ideology",
        7,
    ),
    (
        re.compile(
            r"\b(?:machine\s+stops?|stoppage|collapse|failure|outage|broken|defect|dying)\b",
            re.IGNORECASE,
        ),
        "Risk / Failure Mode",
        7,
    ),
    (
        re.compile(
            r"\b(?:exactly\s+)?(?:alike|same|identical|uniform|standardi[sz]ed|interchangeable|homogeneous)\b",
            re.IGNORECASE,
        ),
        "Standardization / Loss of Difference",
        7,
    ),
    (
        re.compile(
            r"\b(?:book|committee|policy|rule|rules|authority|official|governance)\b",
            re.IGNORECASE,
        ),
        "Authority / Legitimacy",
        5,
    ),
    (
        re.compile(
            r"\b(?:advanced?|progress|innovation|future|modern|thanks|blessed|praise|belief|inevitable)\b",
            re.IGNORECASE,
        ),
        "Aspiration / Ideology",
        4,
    ),
    (
        re.compile(
            r"\b(?:northern|southern|hemisphere|colour|color|pearl|cloud|hills|stars|moon|sun|sky|vision|visions)\b",
            re.IGNORECASE,
        ),
        "Motif / Image Pattern",
        3,
    ),
    (
        re.compile(
            r"\b(?:outage|failure|collapse|stops?|broken|defect|alarm|warning|danger|fragile|dying)\b",
            re.IGNORECASE,
        ),
        "Risk / Failure Mode",
        4,
    ),
]

FORCED_SIGNAL_FAMILY_RULES = [
    (
        re.compile(
            r"\b(?:typhi\s+murium|trypanosoma|trypanozoon|bacillus|plague\s+infection|"
            r"plague\s+rats|indian\s+plague|danysz\s+virus|disease|infected|infection)\b",
            re.IGNORECASE,
        ),
        "Disease / Hazard",
    ),
    (
        re.compile(
            r"\b(?:guinea\s+pigs|biological\s+survey|white\s+rats|laboratory|"
            r"inoculated|experiment|experiments|specimens?|samples?|observed|culture)\b",
            re.IGNORECASE,
        ),
        "Evidence / Experiment",
    ),
    (
        re.compile(
            r"\b(?:destruction\s+rats|extermination\s+rats|destroy\s+rats|"
            r"rats\s+killed|carbon\s+bisulphide|poison|traps?|fumigation|"
            r"sanitation|disinfection|prevention|eradication)\b",
            re.IGNORECASE,
        ),
        "Intervention / Control Method",
    ),
    (
        re.compile(
            r"\b(?:public\s+health|plague\s+commission|marinehospital\s+service|"
            r"marine-hospital\s+service|health\s+service|sanitary\s+board|"
            r"asst\s+surg|assistant\s+surgeon|quarantine|inspection)\b",
            re.IGNORECASE,
        ),
        "Public Health / Institutional Response",
    ),
    (
        re.compile(
            r"\b(?:federal\s+government|national\s+government|state\s+governments?|"
            r"house\s+representatives|trial\s+jury|state\s+legislatures?|"
            r"executive\s+department|legislative\s+body|supreme\s+court|"
            r"proposed\s+constitution|constitution|union|confederacy|republic)\b",
            re.IGNORECASE,
        ),
        "Institutional Structure / Social Design",
    ),
    (
        re.compile(
            r"\b(?:chief\s+magistrate|judicial\s+power|executive\s+power|"
            r"legislative\s+authority|constitutional\s+authority|sovereignty)\b",
            re.IGNORECASE,
        ),
        "Authority / Legitimacy",
    ),
    (
        re.compile(
            r"\b(?:standing\s+armies|faction|tyranny|invasion|war|security|"
            r"public\s+danger|internal\s+violence|usurpation)\b",
            re.IGNORECASE,
        ),
        "Risk / Concern",
    ),
    (
        re.compile(
            r"\b(?:imports\s+exports|power\s+taxation|taxation|commerce|"
            r"liberty\s+security|checks?\s+balances?|tradeoff)\b",
            re.IGNORECASE,
        ),
        "Decision / Tradeoff",
    ),
    (
        re.compile(
            r"\b(?:mending\s+apparatus|power\s+station|communication-system|central\s+power)\b",
            re.IGNORECASE,
        ),
        "Infrastructure / System Dependence",
    ),
    (
        re.compile(
            r"\b(?:public\s+nurseries|nurseries)\b",
            re.IGNORECASE,
        ),
        "Institutional Structure / Social Design",
    ),
    (
        re.compile(
            r"\b(?:advanced\s+thanks|thanks\s+machine|blessed\s+machine)\b",
            re.IGNORECASE,
        ),
        "Aspiration / Ideology",
    ),
    (
        re.compile(
            r"\b(?:machine\s+stops?|stoppage|collapse|failure|outage)\b",
            re.IGNORECASE,
        ),
        "Risk / Failure Mode",
    ),
    (
        re.compile(
            r"\b(?:exactly\s+)?(?:alike|same|identical|uniform|standardi[sz]ed|interchangeable|homogeneous)\b",
            re.IGNORECASE,
        ),
        "Standardization / Loss of Difference",
    ),
]


def semantic_family_score(
    signal: str,
    related_terms: str,
    evidence_text: str,
    evidence_tokens: Iterable[str],
) -> Tuple[str, int]:
    combined = f"{signal} {related_terms} {evidence_text}".lower()
    token_set = set(str(token).lower() for token in evidence_tokens)
    token_set.update(re.findall(r"[a-z][a-z-]+", combined))

    scores = {}
    for pattern, family, score in PATTERN_BASED_FAMILY_RULES:
        if pattern.search(signal) or pattern.search(related_terms):
            scores[family] = max(scores.get(family, 0), score)

    for family, config in SEMANTIC_SIGNAL_FAMILIES.items():
        score = scores.get(family, 0)
        required = REQUIRED_SEMANTIC_MARKERS.get(family)
        if required and not any(marker in token_set or marker in combined for marker in required):
            scores[family] = score
            continue

        for marker in config["markers"]:
            marker_l = marker.lower()
            if " " in marker_l or "-" in marker_l:
                if marker_l in combined:
                    score += 2
            elif marker_l in token_set or marker_l in combined:
                score += 1
        scores[family] = score

    best_family, best_score = max(scores.items(), key=lambda item: item[1])
    return best_family, best_score


def is_low_information_signal(signal: str, support: int, distinctiveness: float) -> bool:
    normalized = " ".join(signal.lower().split())
    if normalized in GENERIC_SIGNAL_PHRASES:
        return True
    parts = [part.lower() for part in signal.split()]
    if not parts:
        return True
    if all(part in LOW_INFORMATION_SIGNAL_TERMS for part in parts):
        return True
    if len(parts) == 2 and any(part in LOW_INFORMATION_SIGNAL_TERMS for part in parts) and support <= 3:
        return True
    return support <= 3 and distinctiveness < 0.75 and any(
        part in LOW_INFORMATION_SIGNAL_TERMS for part in parts
    )

def is_boilerplate_signal(signal: str, related_terms: str = "", evidence_text: str = "") -> bool:
    normalized = " ".join(str(signal).lower().split())
    if not normalized:
        return False
    if normalized in BOILERPLATE_SIGNAL_PHRASES:
        return True
    parts = normalized.split()
    if any(part in BOILERPLATE_SIGNAL_TERMS for part in parts):
        return True
    if any(pattern.search(normalized) for pattern in BOILERPLATE_SIGNAL_PATTERNS):
        return True

    combined_context = f"{related_terms} {evidence_text}".lower()
    if "project gutenberg" in combined_context and len(parts) <= 3:
        return True
    return False


def calibrate_signal_card(
    signal: str,
    signal_type: str,
    related_terms: str,
    evidence_text: str,
    evidence_tokens: Iterable[str],
    support: int,
    distinctiveness: float,
) -> Tuple[str, str, str]:
    if is_low_information_signal(signal, support, distinctiveness):
        return (
            "Low-Specificity Signal",
            (
                f"'{signal}' appears statistically, but it may be too generic to interpret "
                "without stronger supporting evidence. Treat this as a weak lead unless the "
                "excerpt clearly shows a substantive pattern."
            ),
            f"Is '{signal}' carrying real meaning here, or is it mostly connective language?",
        )

    if is_boilerplate_signal(signal, related_terms, evidence_text):
        return (
            "Source / Boilerplate",
            (
                f"'{signal}' looks like source, edition, heading, or publication metadata "
                "rather than a substantive theme. Keep it visible for audit purposes, but "
                "do not treat it as the heart of the text."
            ),
            "Is this phrase part of the document's content, or is it repeated source/formatting language?",
        )

    for pattern, forced_family in FORCED_SIGNAL_FAMILY_RULES:
        if pattern.search(signal):
            config = SEMANTIC_SIGNAL_FAMILIES[forced_family]
            related = f" Related language includes: {related_terms}." if related_terms else ""
            interpretation = f"{config['interpretation']}{related}"
            return forced_family, interpretation, config["question"]

    family, family_score = semantic_family_score(
        signal,
        related_terms,
        evidence_text,
        evidence_tokens,
    )

    threshold = 2 if family == "Motif / Image Pattern" else 3
    if family == "Embodiment / Lived Experience":
        threshold = 5
    if family_score >= threshold:
        config = SEMANTIC_SIGNAL_FAMILIES[family]
        related = f" Related language includes: {related_terms}." if related_terms else ""
        interpretation = f"{config['interpretation']}{related}"
        return family, interpretation, config["question"]

    return (
        signal_type,
        build_interpretation(signal, signal_type, related_terms, support),
        build_followup_question(signal, signal_type),
    )


def find_representative_evidence(
    evidence_docs: List[Dict[str, Any]],
    signal_terms: List[str],
    max_items: int = 3,
) -> List[Dict[str, Any]]:
    if not evidence_docs or not signal_terms:
        return []

    wanted = [t.lower() for t in signal_terms if t]
    scored_docs = []
    for doc in evidence_docs:
        doc_tokens = set(str(t).lower() for t in doc.get("tokens", []))
        excerpt = str(doc.get("excerpt", ""))
        excerpt_lower = excerpt.lower()
        score = 0
        for term in wanted:
            parts = term.split()
            if len(parts) > 1:
                if all(part in doc_tokens or part in excerpt_lower for part in parts):
                    score += 3
            elif term in doc_tokens or term in excerpt_lower:
                score += 2
        if any(pattern.search(excerpt) for pattern in TENSION_PATTERNS):
            score += 1
        if score > 0:
            scored_docs.append((score, doc))

    scored_docs.sort(key=lambda item: item[0], reverse=True)
    return [doc for _score, doc in scored_docs[:max_items]]


def summarize_evidence(evidence: List[Dict[str, Any]]) -> str:
    if not evidence:
        return "No representative excerpt captured. Re-scan live source files to enable evidence anchoring."
    snippets = []
    for item in evidence:
        prefix_parts = []
        if item.get("category"):
            prefix_parts.append(str(item["category"]))
        if item.get("date"):
            prefix_parts.append(str(item["date"]))
        prefix = f"[{', '.join(prefix_parts)}] " if prefix_parts else ""
        snippets.append(prefix + str(item.get("excerpt", "")))
    return "\n\n".join(snippets)


def build_interpretation(signal: str, signal_type: str, related_terms: str, support: int) -> str:
    related = f" It travels with: {related_terms}." if related_terms else ""
    if signal_type == "Pain / Friction":
        return f"This may point to lived friction around '{signal}', especially if the excerpts show repeated operational strain.{related}"
    if signal_type == "Need / Request":
        return f"This may represent an unmet need or recurring ask around '{signal}'.{related}"
    if signal_type == "Blocker / Constraint":
        return f"This may indicate something constraining progress, access, capacity, or execution around '{signal}'.{related}"
    if signal_type == "Aspiration / Opportunity":
        return f"This looks like future-facing or improvement-oriented language around '{signal}'.{related}"
    if signal_type == "Risk / Concern":
        return f"This may signal risk, sensitivity, compliance pressure, or concern around '{signal}'.{related}"
    if signal_type == "Decision / Tradeoff":
        return f"This may mark a decision area where ownership, priorities, or tradeoffs need clarification.{related}"
    if signal_type == "Disease / Hazard":
        return f"This may identify a disease, hazard, vector, or biological agent around '{signal}'.{related}"
    if signal_type == "Evidence / Experiment":
        return f"This may point to evidence, observations, test subjects, samples, or experimental support around '{signal}'.{related}"
    if signal_type == "Intervention / Control Method":
        return f"This may identify a practical control method, prevention measure, or intervention around '{signal}'.{related}"
    if signal_type == "Public Health / Institutional Response":
        return f"This may identify an agency, service, commission, inspection process, or organized public-health response around '{signal}'.{related}"
    if signal_type == "Contradiction / Tension":
        return f"This may indicate a tension between intent and reality, or between different stakeholder needs.{related}"
    return f"This is a candidate theme supported by {support} observed signal(s).{related}"


def build_followup_question(signal: str, signal_type: str) -> str:
    if signal_type == "Pain / Friction":
        return f"What specifically is causing friction around '{signal}', and who experiences it most?"
    if signal_type == "Need / Request":
        return f"What support or decision would satisfy the need behind '{signal}'?"
    if signal_type == "Blocker / Constraint":
        return f"What would have to change for '{signal}' to stop acting as a blocker?"
    if signal_type == "Aspiration / Opportunity":
        return f"What would make the opportunity around '{signal}' concrete and measurable?"
    if signal_type == "Risk / Concern":
        return f"What is the real exposure if the concern around '{signal}' is not addressed?"
    if signal_type == "Decision / Tradeoff":
        return f"Who owns the next decision around '{signal}', and what tradeoff is being made?"
    if signal_type == "Disease / Hazard":
        return f"What hazard, disease pathway, or exposed population does '{signal}' reveal?"
    if signal_type == "Evidence / Experiment":
        return f"What evidence or test is '{signal}' contributing to the larger claim?"
    if signal_type == "Intervention / Control Method":
        return f"What control method or prevention strategy does '{signal}' point toward?"
    if signal_type == "Public Health / Institutional Response":
        return f"What institution or coordinated response is organized around '{signal}'?"
    if signal_type == "Contradiction / Tension":
        return f"What contradiction or mismatch is the text revealing around '{signal}'?"
    return f"What is the practical implication of the recurring signal '{signal}'?"


def confidence_label(support: int, evidence_count: int, distinctiveness: float) -> str:
    # Confidence is intentionally conservative: repeated evidence matters more
    # than a single distinctive phrase, but distinctive one-offs can still be Medium.
    if support >= 20 and evidence_count >= 3:
        return "High"
    if support >= 6 and evidence_count >= 2:
        return "Medium"
    if distinctiveness >= 0.5 and evidence_count >= 1:
        return "Medium"
    return "Low"

def signal_phrase_quality(signal: str, signal_type: str, support: int, distinctiveness: float) -> Tuple[int, List[str]]:
    normalized = " ".join(str(signal).lower().split())
    parts = normalized.split()
    reasons = []
    # Start from a neutral midpoint, then reward clear/compact signals and
    # penalize boilerplate, fragments, and generic connective language.
    score = 50

    if not parts:
        return 0, ["empty signal"]

    if 1 <= len(parts) <= 3:
        score += 10
        reasons.append("compact phrase")
    if len(parts) == 2:
        score += 5
        reasons.append("readable two-word signal")
    if support >= 5:
        score += 8
        reasons.append("repeated evidence")
    if distinctiveness >= 0.8:
        score += 8
        reasons.append("distinctive pairing")
    if signal_type not in {"Theme / Topic", "Low-Specificity Signal"}:
        score += 8
        reasons.append("semantic category match")

    if signal_type == "Source / Boilerplate" or is_boilerplate_signal(signal):
        score -= 45
        reasons.append("source or boilerplate language")
    if is_low_information_signal(normalized, support, distinctiveness):
        score -= 35
        reasons.append("generic wording")
    if parts[0] in FRAGMENT_START_TERMS or parts[-1] in FRAGMENT_END_TERMS:
        score -= 18
        reasons.append("phrase may be a fragment")
    if len(parts) == 2 and parts[-1] in GENERIC_BIGRAM_END_TERMS:
        score -= 20
        reasons.append("generic verb pairing")
    if len(parts) == 2 and all(part in ABSTRACT_FRAGMENT_TERMS for part in parts):
        score -= 30
        reasons.append("abstract phrase fragment")
    elif any(part in ABSTRACT_FRAGMENT_TERMS for part in parts) and signal_type == "Contradiction / Tension":
        score -= 18
        reasons.append("abstract tension wording")
    if all(part in CONTEXT_REFERENCE_TERMS for part in parts):
        score -= 25
        reasons.append("mostly context/reference language")
    elif any(part in CONTEXT_REFERENCE_TERMS for part in parts) and support <= 3:
        score -= 12
        reasons.append("context-heavy phrase")

    return max(0, min(100, score)), reasons

def classify_signal_role(
    signal: str,
    signal_type: str,
    support: int,
    distinctiveness: float,
    evidence_count: int,
    phrase_quality: int,
) -> Tuple[str, str]:
    normalized = " ".join(str(signal).lower().split())
    parts = normalized.split()

    if signal_type == "Low-Specificity Signal" or is_low_information_signal(normalized, support, distinctiveness):
        return "Low-Specificity", "Generic or connective wording; useful mainly as a weak lead."

    if signal_type == "Source / Boilerplate" or is_boilerplate_signal(signal):
        return "Context / Reference", "Likely source, edition, heading, footer, or publication metadata."

    if signal_type == "Absence / Weak Signal":
        return "Supporting Signal", "Expected concept check rather than an organically discovered core signal."

    if signal_type == "Motif / Image Pattern":
        if parts and all(part in CONTEXT_REFERENCE_TERMS for part in parts):
            return "Context / Reference", "Mostly location, reference, or descriptive setting language."
        return "Supporting Motif", "Recurring image language that may support interpretation."

    if parts and all(part in CONTEXT_REFERENCE_TERMS for part in parts):
        return "Context / Reference", "Mostly location, reference, or named-context language."
    if len(parts) == 2 and parts[-1] in GENERIC_BIGRAM_END_TERMS:
        return "Low-Specificity", "The phrase is mostly a generic verb pairing rather than a stable concept."
    if len(parts) == 2 and all(part in ABSTRACT_FRAGMENT_TERMS for part in parts):
        return "Low-Specificity", "The phrase is an abstract fragment rather than a clear substantive signal."

    if phrase_quality < 50:
        return "Supporting Signal", "Potentially meaningful, but the phrase is not clean enough to lead the analysis."

    high_priority = SIGNAL_TYPE_PRIORITY.get(signal_type, 1.0) >= 1.15
    enough_evidence = support >= 5 or evidence_count >= 3
    high_lift_candidate = phrase_quality >= 70 and distinctiveness >= 0.65
    if high_priority and enough_evidence and high_lift_candidate:
        return "Core Insight", "Strong enough to help shape the top-level read."

    if support >= 7 and evidence_count >= 3 and phrase_quality >= 75:
        return "Core Insight", "Repeated and clear enough to help shape the top-level read."

    return "Supporting Signal", "Useful supporting evidence, but not necessarily the center of the analysis."

def interpretive_lift_score(
    signal_type: str,
    role: str,
    support: int,
    distinctiveness: float,
    evidence_count: int,
    phrase_quality: int,
) -> int:
    # Lift is directional rather than scientific: log-scaled support prevents
    # very frequent terms from overwhelming distinctiveness, evidence, and role.
    support_score = min(math.log1p(max(support, 0)) * 12, 30)
    distinctiveness_score = min(max(distinctiveness, 0.0) * 20, 20)
    evidence_score = min(evidence_count * 8, 20)
    family_score = SIGNAL_TYPE_PRIORITY.get(signal_type, 1.0) * 8
    quality_score = phrase_quality * 0.18
    role_adjustment = {
        "Core Insight": 14,
        "Supporting Signal": 5,
        "Supporting Motif": 0,
        "Context / Reference": -12,
        "Low-Specificity": -25,
    }.get(role, 0)
    score = support_score + distinctiveness_score + evidence_score + family_score + quality_score + role_adjustment
    return int(max(0, min(100, round(score))))

def build_ranking_rationale(
    role: str,
    role_reason: str,
    phrase_quality_reasons: List[str],
    lift_score: int,
) -> str:
    details = ", ".join(phrase_quality_reasons[:3]) if phrase_quality_reasons else "general signal strength"
    return f"{role}: {role_reason} Interpretive lift {lift_score}/100 based on {details}."

def signal_token_set(signal: str) -> Set[str]:
    return {
        token
        for token in re.findall(r"[a-z][a-z-]+", str(signal).lower())
        if token not in LOW_INFORMATION_SIGNAL_TERMS
    }

def is_near_duplicate_signal(candidate: pd.Series, selected_rows: List[pd.Series]) -> bool:
    candidate_terms = signal_token_set(candidate.get("Signal", ""))
    if not candidate_terms:
        return False
    for selected in selected_rows:
        selected_terms = signal_token_set(selected.get("Signal", ""))
        if not selected_terms:
            continue
        overlap = len(candidate_terms & selected_terms)
        overlap_ratio = overlap / max(1, min(len(candidate_terms), len(selected_terms)))
        same_type = candidate.get("Signal Type") == selected.get("Signal Type")
        same_role = candidate.get("Signal Role") == selected.get("Signal Role")
        # Strong token overlap is treated as duplicate regardless of category.
        # Moderate overlap is only suppressed when the analytical role/type also matches.
        if overlap_ratio >= 0.67:
            return True
        if overlap_ratio >= 0.5 and (same_type or same_role):
            return True
    return False

def select_balanced_insight_cards(df: pd.DataFrame, top_n: int) -> pd.DataFrame:
    if df.empty:
        return df

    sorted_df = df.copy()
    sorted_df["Role Priority"] = sorted_df["Signal Role"].map(SIGNAL_ROLE_PRIORITY).fillna(1)
    sorted_df = sorted_df.sort_values(
        ["Role Priority", "Interpretive Lift", "Evidence Strength", "Distinctiveness"],
        ascending=[False, False, False, False],
    )

    role_caps = {
        "Core Insight": max(4, min(6, top_n - 3)),
        "Supporting Signal": 3,
        "Supporting Motif": 2,
        "Context / Reference": 2,
        "Low-Specificity": 2,
    }
    role_counts = Counter()
    selected_rows = []

    for _, row in sorted_df.iterrows():
        role = row.get("Signal Role", "Supporting Signal")
        if role_counts[role] >= role_caps.get(role, top_n):
            continue
        if is_near_duplicate_signal(row, selected_rows):
            continue
        selected_rows.append(row)
        role_counts[role] += 1
        if len(selected_rows) >= top_n:
            break

    if len(selected_rows) < top_n:
        for _, row in sorted_df.iterrows():
            if any(row.get("Signal") == selected.get("Signal") for selected in selected_rows):
                continue
            if is_near_duplicate_signal(row, selected_rows):
                continue
            selected_rows.append(row)
            if len(selected_rows) >= top_n:
                break

    if len(selected_rows) < top_n:
        for _, row in sorted_df.iterrows():
            if any(row.get("Signal") == selected.get("Signal") for selected in selected_rows):
                continue
            selected_rows.append(row)
            if len(selected_rows) >= top_n:
                break

    if not selected_rows:
        return sorted_df.head(top_n).drop(columns=["Role Priority"], errors="ignore").reset_index(drop=True)

    result = pd.DataFrame(selected_rows).drop(columns=["Role Priority"], errors="ignore")
    return result.reset_index(drop=True)


def build_insight_cards(
    scanner: StreamScanner,
    counts: Counter,
    expected_terms_raw: str = "",
    top_n: int = 8,
) -> pd.DataFrame:
    theme_df = build_theme_evidence_cards(scanner, counts, top_n=max(top_n * 4, 30))
    rows = []
    used = set()

    for _, row in theme_df.iterrows():
        signal = str(row.get("Theme Evidence", "")).strip()
        if not signal or signal in used:
            continue
        used.add(signal)

        signal_terms = signal.split()
        evidence = find_representative_evidence(scanner.evidence_docs, signal_terms, max_items=3)
        evidence_text = " ".join(item.get("excerpt", "") for item in evidence)
        evidence_tokens = []
        for item in evidence:
            evidence_tokens.extend(item.get("tokens", []))

        signal_type, _score = classify_signal_dimension(
            f"{signal} {evidence_text} {row.get('Related Terms', '')}",
            evidence_tokens + signal_terms,
        )
        support = int(row.get("Support", 0) or 0)
        distinctiveness = float(row.get("Distinctiveness", 0.0) or 0.0)
        confidence = confidence_label(support, len(evidence), distinctiveness)
        related_terms = str(row.get("Related Terms", ""))
        calibrated_type, calibrated_interpretation, calibrated_question = calibrate_signal_card(
            signal=signal,
            signal_type=signal_type,
            related_terms=related_terms,
            evidence_text=evidence_text,
            evidence_tokens=evidence_tokens + signal_terms,
            support=support,
            distinctiveness=distinctiveness,
        )
        phrase_quality, phrase_quality_reasons = signal_phrase_quality(
            signal=signal,
            signal_type=calibrated_type,
            support=support,
            distinctiveness=distinctiveness,
        )
        signal_role, role_reason = classify_signal_role(
            signal=signal,
            signal_type=calibrated_type,
            support=support,
            distinctiveness=distinctiveness,
            evidence_count=len(evidence),
            phrase_quality=phrase_quality,
        )
        lift_score = interpretive_lift_score(
            signal_type=calibrated_type,
            role=signal_role,
            support=support,
            distinctiveness=distinctiveness,
            evidence_count=len(evidence),
            phrase_quality=phrase_quality,
        )

        rows.append({
            "Signal": signal,
            "Signal Type": calibrated_type,
            "Signal Role": signal_role,
            "Interpretive Lift": lift_score,
            "Evidence Strength": support,
            "Distinctiveness": round(distinctiveness, 3),
            "Confidence": confidence,
            "Ranking Rationale": build_ranking_rationale(
                signal_role,
                role_reason,
                phrase_quality_reasons,
                lift_score,
            ),
            "Representative Evidence": summarize_evidence(evidence),
            "Interpretation": calibrated_interpretation,
            "Follow-up Question": calibrated_question,
        })

    expected_df = build_expected_terms_df(expected_terms_raw, counts, scanner.global_bigrams)
    if not expected_df.empty:
        missing = expected_df[expected_df["Status"].isin(["Missing", "Weak"])].head(4)
        for _, row in missing.iterrows():
            signal = str(row["Expected Signal"])
            if signal in used:
                continue
            rows.append({
                "Signal": signal,
                "Signal Type": "Absence / Weak Signal",
                "Signal Role": "Supporting Signal",
                "Interpretive Lift": 35,
                "Evidence Strength": int(row["Observed Count"]),
                "Distinctiveness": 0.0,
                "Confidence": "Medium" if int(row["Observed Count"]) == 0 else "Low",
                "Ranking Rationale": "Supporting Signal: expected concept check rather than an organically discovered core signal.",
                "Representative Evidence": "Expected concept was missing or weak in the observed vocabulary.",
                "Interpretation": f"'{signal}' was expected but is {str(row['Status']).lower()}. This may be meaningful absence, source mismatch, or vocabulary mismatch.",
                "Follow-up Question": f"Should '{signal}' be present in this corpus, and if so, why is it not showing up clearly?",
            })

    df = pd.DataFrame(rows)
    if df.empty:
        return df
    return select_balanced_insight_cards(df, top_n)


def build_ai_insight_context(insight_df: pd.DataFrame, max_cards: int = 6) -> str:
    if insight_df.empty:
        return ""
    lines = ["Insight Cards:"]
    for idx, row in insight_df.head(max_cards).iterrows():
        lines.append(
            f"{idx + 1}. {row['Signal']} | {row['Signal Type']} | "
            f"role={row.get('Signal Role', 'Signal')} | "
            f"lift={row.get('Interpretive Lift', 'n/a')} | "
            f"confidence={row['Confidence']} | interpretation={row['Interpretation']}"
        )
    return "\n".join(lines)

def build_ai_context_brief(
    scanner: StreamScanner,
    combined_counts: Counter,
    insight_df: pd.DataFrame,
    text_stats: Dict[str, Any],
    maturity_result: Optional[Dict[str, Any]] = None,
    graph_cluster_info: str = "N/A",
    max_items: int = 20,
) -> str:
    lines = [
        "Signal Foundry AI Context Brief",
        "",
        "Privacy boundary: this context contains summaries, counts, labels, and short insight-card interpretations. It does not include the full raw source documents.",
        "",
        "Corpus Stats:",
        f"- Rows processed: {scanner.total_rows_processed:,}",
        f"- Total tokens after cleaning: {sum(combined_counts.values()):,}",
        f"- Unique vocabulary: {len(combined_counts):,}",
        f"- Evidence snippets retained locally: {len(scanner.evidence_docs):,}",
        f"- Evidence cap reached: {scanner.evidence_limit_reached}",
    ]

    if text_stats:
        lines.extend([
            f"- Avg word length: {text_stats.get('Avg Word Length', 'N/A')}",
            f"- Lexical diversity: {text_stats.get('Lexical Diversity', 'N/A')}",
        ])

    top_terms = ", ".join([f"{w} ({c})" for w, c in combined_counts.most_common(max_items)])
    top_bigrams = ", ".join([
        f"{w1} {w2} ({c})"
        for (w1, w2), c in scanner.global_bigrams.most_common(max_items)
    ])
    top_entities = ", ".join([
        f"{entity} ({count})"
        for entity, count in scanner.entity_counts.most_common(12)
    ])

    lines.extend([
        "",
        "Top Terms:",
        top_terms or "No top terms available.",
        "",
        "Top Bigrams:",
        top_bigrams or "No bigrams available.",
        "",
        "Top Entities:",
        top_entities or "No entities detected.",
    ])

    if graph_cluster_info and graph_cluster_info != "N/A":
        lines.extend(["", "Graph / Community Sketch:", graph_cluster_info])

    if isinstance(insight_df, pd.DataFrame) and not insight_df.empty:
        compass_df = build_signal_compass_df(insight_df)
        if not compass_df.empty:
            compass_lines = [
                f"- {row['Direction']}: {row['Share']:.0%} (example: {row['Example Signal']})"
                for _, row in compass_df.head(5).iterrows()
            ]
            lines.extend(["", "Signal Compass:", *compass_lines])

        resource_shape = build_resource_shape(insight_df)
        lines.extend([
            "",
            "Resource Shape:",
            resource_shape.get("summary", "No Resource Shape summary available."),
        ])

        if resource_shape.get("families"):
            lines.append("Signal Families:")
            for family in resource_shape["families"][:5]:
                lines.append(
                    f"- {family.get('Signal Type')}: {family.get('Cards')} card(s), "
                    f"example {family.get('TopSignal')}"
                )

        role_mix = signal_role_distribution_dataframe(insight_df)
        if not role_mix.empty:
            lines.append("Signal Role Mix:")
            for _, row in role_mix.iterrows():
                lines.append(f"- {row['Signal Role']}: {row['Cards']}")

        lines.extend(["", build_ai_insight_context(insight_df, max_cards=8)])

    if maturity_result:
        lines.extend(["", "Maturity Context:"])
        if maturity_result.get("type") == "domain_based":
            domain_summary_parts = []
            for dk in sorted(maturity_result.get("domain_results", {}).keys()):
                dr = maturity_result["domain_results"][dk]
                if dr.get("signals", 0) > 0:
                    domain_summary_parts.append(
                        f"- {dr['short']}: {dr['score']}/3.0 ({dr['tier_label']}), signals={dr['signals']}"
                    )
            if domain_summary_parts:
                lines.extend(domain_summary_parts)
                lines.append(f"Composite Score: {maturity_result.get('overall_score', 'N/A')}/3.0")
            else:
                lines.append("No maturity domains had enough signal to summarize.")
        else:
            lines.append(f"Maturity Score: {maturity_result.get('overall_score', 'N/A')}/5.0")

    return "\n".join(lines)


def concise_sentence(text: str, max_chars: int = 220) -> str:
    clean = " ".join(str(text).split())
    if len(clean) <= max_chars:
        return clean
    clipped = clean[:max_chars].rsplit(" ", 1)[0].rstrip()
    return f"{clipped}..."


def first_sentence(text: str, max_chars: int = 180) -> str:
    clean = " ".join(str(text).split())
    match = re.search(r"(?<=[.!?])\s+", clean)
    sentence = clean[:match.start()].strip() if match else clean
    return concise_sentence(sentence, max_chars)


def build_resource_shape(insight_df: pd.DataFrame, top_n: int = 5) -> Dict[str, Any]:
    if insight_df.empty:
        return {
            "summary": "Not enough signal yet to describe the shape of this resource.",
            "families": [],
            "key_insights": [],
        }

    usable = insight_df[
        ~insight_df["Signal Type"].isin([
            "Absence / Weak Signal",
            "Low-Specificity Signal",
            "Source / Boilerplate",
        ])
    ].copy()
    if "Signal Role" in usable.columns:
        role_filtered = usable[
            ~usable["Signal Role"].isin(["Context / Reference", "Low-Specificity"])
        ].copy()
        if not role_filtered.empty:
            usable = role_filtered
    if usable.empty:
        usable = insight_df.copy()

    confidence_weight = {"High": 3.0, "Medium": 1.5, "Low": 0.5}
    family_priority = usable["Signal Type"].map(SIGNAL_TYPE_PRIORITY).fillna(1.0)
    role_shape_weight = (
        usable["Signal Role"]
        .map({
            "Core Insight": 1.0,
            "Supporting Signal": 0.8,
            "Supporting Motif": 0.35,
            "Context / Reference": 0.15,
            "Low-Specificity": 0.1,
        })
        .fillna(0.75)
        if "Signal Role" in usable.columns
        else 1.0
    )
    tension_shape_weight = (
        usable["Signal Type"]
        .map({
            "Contradiction / Tension": 0.55,
        })
        .fillna(1.0)
    )
    lift_component = usable["Interpretive Lift"].fillna(50).astype(float) / 10 if "Interpretive Lift" in usable.columns else 0
    # Resource Shape is intentionally directional, not a formal score:
    # base signal strength is adjusted by family priority, signal role, and
    # a tension demotion so one useful-but-common category does not dominate.
    usable["Shape Weight"] = (
        (
            usable["Evidence Strength"].fillna(0).astype(float)
            + usable["Distinctiveness"].fillna(0).astype(float) * 2
            + usable["Confidence"].map(confidence_weight).fillna(0.5)
            + lift_component
        )
        * family_priority
        * role_shape_weight
        * tension_shape_weight
    )

    family_rows = (
        usable.groupby("Signal Type", as_index=False)
        .agg(
            Weight=("Shape Weight", "sum"),
            Cards=("Signal", "count"),
            TopSignal=("Signal", "first"),
        )
        .sort_values("Weight", ascending=False)
        .head(4)
    )
    families = family_rows.to_dict("records")

    if families:
        family_names = [str(item["Signal Type"]) for item in families[:3]]
        if len(family_names) == 1:
            family_phrase = family_names[0]
        elif len(family_names) == 2:
            family_phrase = f"{family_names[0]} and {family_names[1]}"
        else:
            family_phrase = f"{family_names[0]}, {family_names[1]}, and {family_names[2]}"
        summary = (
            f"This resource appears to be shaped mainly by {family_phrase}. "
            "Treat this as a first-pass synthesis: use the supporting evidence below "
            "to confirm, revise, or reject the interpretation."
        )
    else:
        summary = (
            "This resource has candidate signals, but they do not yet cluster into a clear shape. "
            "Review the evidence cards and consider adjusting cleaning settings or stopwords."
        )

    key_rows = usable.sort_values("Shape Weight", ascending=False).head(top_n)
    key_insights = []
    for _, row in key_rows.iterrows():
        key_insights.append({
            "Signal": str(row["Signal"]),
            "Signal Type": str(row["Signal Type"]),
            "Confidence": str(row["Confidence"]),
            "Signal Role": str(row.get("Signal Role", "Supporting Signal")),
            "Interpretive Lift": int(row.get("Interpretive Lift", 0)),
            "Evidence Strength": int(row["Evidence Strength"]),
            "Ranking Rationale": str(row.get("Ranking Rationale", "")),
            "Interpretation": first_sentence(row["Interpretation"], 220),
            "Evidence": concise_sentence(row["Representative Evidence"], 320),
            "Question": str(row["Follow-up Question"]),
        })

    return {
        "summary": summary,
        "families": families,
        "key_insights": key_insights,
    }


def resource_shape_diagnostics_dataframe(insight_df: pd.DataFrame) -> pd.DataFrame:
    """
    Exposes the Resource Shape weighting ingredients for calibration review.
    This mirrors the current formula and does not change rankings.
    """
    columns = [
        "Signal",
        "Signal Type",
        "Signal Role",
        "Evidence Strength",
        "Distinctiveness",
        "Confidence",
        "Confidence Weight",
        "Lift Component",
        "Family Priority",
        "Role Weight",
        "Tension Weight",
        "Base Signal Score",
        "Shape Weight",
        "Ranking Rationale",
    ]
    if insight_df is None or insight_df.empty:
        return pd.DataFrame(columns=columns)

    usable = insight_df[
        ~insight_df["Signal Type"].isin([
            "Absence / Weak Signal",
            "Low-Specificity Signal",
            "Source / Boilerplate",
        ])
    ].copy()
    if "Signal Role" in usable.columns:
        role_filtered = usable[
            ~usable["Signal Role"].isin(["Context / Reference", "Low-Specificity"])
        ].copy()
        if not role_filtered.empty:
            usable = role_filtered
    if usable.empty:
        usable = insight_df.copy()

    confidence_weight = {"High": 3.0, "Medium": 1.5, "Low": 0.5}
    role_weight_map = {
        "Core Insight": 1.0,
        "Supporting Signal": 0.8,
        "Supporting Motif": 0.35,
        "Context / Reference": 0.15,
        "Low-Specificity": 0.1,
    }
    tension_weight_map = {"Contradiction / Tension": 0.55}

    usable["Confidence Weight"] = usable["Confidence"].map(confidence_weight).fillna(0.5)
    usable["Lift Component"] = (
        usable["Interpretive Lift"].fillna(50).astype(float) / 10
        if "Interpretive Lift" in usable.columns
        else 0
    )
    usable["Family Priority"] = usable["Signal Type"].map(SIGNAL_TYPE_PRIORITY).fillna(1.0)
    usable["Role Weight"] = (
        usable["Signal Role"].map(role_weight_map).fillna(0.75)
        if "Signal Role" in usable.columns
        else 1.0
    )
    usable["Tension Weight"] = usable["Signal Type"].map(tension_weight_map).fillna(1.0)
    usable["Base Signal Score"] = (
        usable["Evidence Strength"].fillna(0).astype(float)
        + usable["Distinctiveness"].fillna(0).astype(float) * 2
        + usable["Confidence Weight"]
        + usable["Lift Component"]
    )
    usable["Shape Weight"] = (
        usable["Base Signal Score"]
        * usable["Family Priority"]
        * usable["Role Weight"]
        * usable["Tension Weight"]
    ).round(3)

    if "Ranking Rationale" not in usable.columns:
        usable["Ranking Rationale"] = ""
    if "Signal Role" not in usable.columns:
        usable["Signal Role"] = "Supporting Signal"

    return usable[columns].sort_values("Shape Weight", ascending=False).reset_index(drop=True)

SIGNAL_DIRECTION_LABELS = {
    "Infrastructure / System Dependence": "System-dependence heavy",
    "Risk / Failure Mode": "Failure-mode heavy",
    "Risk / Concern": "Risk-centered",
    "Institutional Structure / Social Design": "Institutional-design heavy",
    "Authority / Legitimacy": "Authority/legitimacy heavy",
    "Decision / Tradeoff": "Tradeoff/decision heavy",
    "Disease / Hazard": "Hazard/disease heavy",
    "Evidence / Experiment": "Evidence/testing heavy",
    "Intervention / Control Method": "Intervention/control heavy",
    "Public Health / Institutional Response": "Public-health response heavy",
    "Standardization / Loss of Difference": "Standardization heavy",
    "Embodiment / Lived Experience": "Lived-experience heavy",
    "Aspiration / Ideology": "Ideology/progress heavy",
    "Motif / Image Pattern": "Motif/image heavy",
    "Isolation / Disconnection": "Isolation/disconnection heavy",
    "Contradiction / Tension": "Tension-heavy",
}

ANALYSIS_HELP_TEXT = {
    "Signal Role": (
        "How this item is being used in the analysis. Core Insights are strongest; "
        "Supporting Signals, Motifs, Context, and Low-Specificity items help explain what was promoted or demoted."
    ),
    "Signal Type": (
        "The app's best-fit analytical category for this signal, such as risk, infrastructure, "
        "institutional structure, evidence, disease, or tradeoff."
    ),
    "Signal Family": (
        "A broader analytical category used to summarize the overall shape of the resource."
    ),
    "Direction": (
        "A plain-language label for the main analytical pull of the text, such as risk-centered or evidence-heavy."
    ),
    "Evidence Strength": (
        "How often this signal appears or is supported in the processed text. Higher means more repeated evidence, not necessarily more importance."
    ),
    "Distinctiveness": (
        "How strongly this phrase or signal stands out compared with ordinary word-pairing patterns. Higher means more unusually associated, not just frequent."
    ),
    "Confidence": (
        "A rough reliability label based on support, available evidence, and distinctiveness. Treat it as a guide, not proof."
    ),
    "Interpretive Lift": (
        "A directional ranking score for how useful this signal is likely to be. It combines evidence strength, distinctiveness, signal role, semantic fit, and phrase quality."
    ),
    "Cards": (
        "How many insight cards contributed to this signal family or role."
    ),
    "Example Signal": (
        "A representative signal that helped place this family or direction in the summary."
    ),
    "Share": (
        "This direction's approximate share of the top Resource Shape weight."
    ),
    "Rows": (
        "How many source units were processed, such as lines, spreadsheet rows, transcript segments, or document chunks."
    ),
    "Tokens": (
        "The total number of cleaned words available for analysis after filtering."
    ),
    "Insight Cards": (
        "The number of ranked evidence cards the app produced from the strongest available signals."
    ),
    "Evidence Snippets": (
        "Short retained excerpts used to connect insight cards back to source evidence."
    ),
    "Groups / Dates": (
        "How many category groups and time points were detected for comparison or trend views."
    ),
    "Total Tokens": (
        "The total count of all words processed after cleaning, such as removing stopwords and numbers."
    ),
    "Unique Vocab": (
        "The number of distinct words found. Higher values often mean broader vocabulary or more varied subject matter."
    ),
    "Docs/Rows": (
        "The number of separate processing units, such as rows, paragraphs, pages, or transcript lines."
    ),
    "Lexical Diversity": (
        "The ratio of unique words to total words. Higher means more varied language; lower means more repetition."
    ),
    "Signal Types": (
        "The number of analytical categories available for classifying signals."
    ),
    "Evidence Cap": (
        "The maximum number of representative snippets retained for evidence review."
    ),
    "Frequency": (
        "How often this term appears in the processed text."
    ),
    "Evidence Score": (
        "A combined directional score based on frequency and distinctiveness."
    ),
    "Quadrant": (
        "A simple label showing whether a term is common backdrop, niche signal, core signal, or low evidence."
    ),
    "Maturity Score": (
        "A directional score based on maturity vocabulary found in the text. It is a guide, not a formal assessment."
    ),
    "Signal Density": (
        "How many maturity-related vocabulary signals were found."
    ),
    "Composite Maturity": (
        "A directional maturity score across assessed domains."
    ),
    "Domains Assessed": (
        "How many maturity domains had enough matching evidence to be scored."
    ),
    "Total Signals": (
        "The total number of maturity-related signals found across assessed domains."
    ),
    "Foundational": (
        "Share of domain evidence that maps to foundational maturity language."
    ),
    "Advanced": (
        "Share of domain evidence that maps to more advanced maturity language."
    ),
    "Leading Edge": (
        "Share of domain evidence that maps to leading-edge maturity language."
    ),
    "Positive Words Observed": (
        "How many processed terms matched the app's positive sentiment vocabulary."
    ),
    "Negative Words Observed": (
        "How many processed terms matched the app's negative sentiment vocabulary."
    ),
}

def build_signal_compass_df(insight_df: pd.DataFrame) -> pd.DataFrame:
    shape = build_resource_shape(insight_df)
    families = shape.get("families", [])
    if not families:
        return pd.DataFrame(columns=["Signal Family", "Direction", "Weight", "Share", "Example Signal"])

    df = pd.DataFrame(families).rename(columns={
        "Signal Type": "Signal Family",
        "TopSignal": "Example Signal",
    })
    df["Direction"] = df["Signal Family"].map(SIGNAL_DIRECTION_LABELS).fillna(
        df["Signal Family"].astype(str) + " heavy"
    )
    total_weight = max(float(df["Weight"].sum()), 1.0)
    df["Share"] = (df["Weight"].astype(float) / total_weight).round(3)
    return df[["Signal Family", "Direction", "Weight", "Share", "Cards", "Example Signal"]]

def render_signal_compass_panel(insight_df: pd.DataFrame):
    compass_df = build_signal_compass_df(insight_df)

    st.subheader("🧭 Signal Compass")
    st.caption(
        "A directional view of the resource: which analytical forces are pulling the text most strongly."
    )

    if compass_df.empty:
        st.info("Not enough signal yet to build the compass.")
        return

    top_rows = compass_df.head(3).reset_index(drop=True)
    badge_cols = st.columns(len(top_rows))
    for idx, row in top_rows.iterrows():
        with badge_cols[idx]:
            st.metric(
                row["Direction"],
                f"{int(row['Share'] * 100)}%",
                help=(
                    f"{ANALYSIS_HELP_TEXT['Direction']} "
                    f"{ANALYSIS_HELP_TEXT['Share']} "
                    f"Example signal: {row['Example Signal']}"
                ),
            )
            st.caption(f"Example: {row['Example Signal']}")

    if alt is not None:
        chart_df = compass_df.copy()
        chart = (
            alt.Chart(chart_df)
            .mark_bar(cornerRadius=4)
            .encode(
                x=alt.X("Share:Q", title="Directional weight", axis=alt.Axis(format="%")),
                y=alt.Y("Direction:N", sort="-x", title=None),
                tooltip=[
                    alt.Tooltip("Signal Family:N", title="Signal Family"),
                    alt.Tooltip("Direction:N", title="Direction"),
                    alt.Tooltip("Cards:Q", title="Cards"),
                    alt.Tooltip("Example Signal:N", title="Example Signal"),
                    alt.Tooltip("Share:Q", title="Share", format=".0%"),
                ],
            )
            .properties(height=max(150, len(chart_df) * 34))
        )
        st.altair_chart(chart, use_container_width=True)
    else:
        display_df = compass_df.copy()
        display_df["Share"] = (display_df["Share"] * 100).round(1).astype(str) + "%"
        st.dataframe(
            display_df,
            use_container_width=True,
            hide_index=True,
            column_config={
                "Signal Family": st.column_config.TextColumn(
                    "Signal Family",
                    help=ANALYSIS_HELP_TEXT["Signal Family"],
                ),
                "Direction": st.column_config.TextColumn(
                    "Direction",
                    help=ANALYSIS_HELP_TEXT["Direction"],
                ),
                "Cards": st.column_config.NumberColumn(
                    "Cards",
                    help=ANALYSIS_HELP_TEXT["Cards"],
                ),
                "Example Signal": st.column_config.TextColumn(
                    "Example Signal",
                    help=ANALYSIS_HELP_TEXT["Example Signal"],
                ),
                "Share": st.column_config.TextColumn(
                    "Share",
                    help=ANALYSIS_HELP_TEXT["Share"],
                ),
            },
        )

    top_direction = str(compass_df.iloc[0]["Direction"])
    secondary = ", ".join(compass_df["Direction"].head(3).tolist()[1:])
    if secondary:
        st.info(f"Directional read: **{top_direction}**, with secondary pull from {secondary}.")
    else:
        st.info(f"Directional read: **{top_direction}**.")


def render_resource_shape_panel(insight_df: pd.DataFrame):
    shape = build_resource_shape(insight_df)

    st.subheader("🧭 Resource Shape")
    st.caption(
        "A plain-language synthesis of the strongest signal families. Start here, then inspect the supporting cards."
    )

    with st.container(border=True):
        st.markdown(f"**First-pass read:** {shape['summary']}")

        if shape["families"]:
            family_df = pd.DataFrame(shape["families"])
            family_df = family_df.rename(columns={
                "Signal Type": "Signal Family",
                "TopSignal": "Example Signal",
            })
            st.dataframe(
                family_df[["Signal Family", "Cards", "Example Signal"]],
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Signal Family": st.column_config.TextColumn(
                        "Signal Family",
                        help=ANALYSIS_HELP_TEXT["Signal Family"],
                    ),
                    "Cards": st.column_config.NumberColumn(
                        "Cards",
                        help=ANALYSIS_HELP_TEXT["Cards"],
                    ),
                    "Example Signal": st.column_config.TextColumn(
                        "Example Signal",
                        help=ANALYSIS_HELP_TEXT["Example Signal"],
                    ),
                },
            )

        if shape["key_insights"]:
            st.markdown("#### Key Insights")
            for idx, item in enumerate(shape["key_insights"], start=1):
                with st.expander(
                    f"{idx}. {item.get('Signal Role', 'Signal')} · {item['Signal Type']} · {item['Signal']}",
                    expanded=(idx <= 3),
                ):
                    c_role, c_lift = st.columns(2)
                    c_role.metric(
                        "Signal Role",
                        item.get("Signal Role", "Supporting Signal"),
                        help=ANALYSIS_HELP_TEXT["Signal Role"],
                    )
                    c_lift.metric(
                        "Interpretive Lift",
                        f"{item.get('Interpretive Lift', 0)}/100",
                        help=ANALYSIS_HELP_TEXT["Interpretive Lift"],
                    )
                    if item.get("Ranking Rationale"):
                        st.caption(item["Ranking Rationale"])
                    st.write(item["Interpretation"])
                    st.markdown("**Evidence preview**")
                    st.write(item["Evidence"])
                    st.markdown("**Question to carry forward**")
                    st.info(item["Question"])
        else:
            st.info("No key insights yet. Scan more content or review cleaning settings.")


def build_dashboard_next_steps(
    scanner: StreamScanner,
    insight_df: pd.DataFrame,
    has_temporal: bool,
    has_categories: bool,
) -> List[str]:
    steps = []
    if not insight_df.empty:
        steps.append("Start with Resource Shape to get the top-level synthesis before opening individual cards.")
        high_conf = insight_df[insight_df["Confidence"].isin(["High", "Medium"])]
        if not high_conf.empty:
            top_signal = str(high_conf.iloc[0]["Signal"])
            steps.append(f"Then inspect the supporting evidence card for '{top_signal}'.")
        pain_or_blocker = insight_df[
            insight_df["Signal Type"].isin(["Pain / Friction", "Blocker / Constraint"])
        ]
        if not pain_or_blocker.empty:
            steps.append("Review pain and blocker cards before interpreting maturity scores.")
    if has_categories:
        steps.append("Use Contrastive Analysis to compare stakeholder groups or source types.")
    if has_temporal:
        steps.append("Use Temporal Drift to see what is rising or fading over time.")
    if scanner.entity_counts:
        steps.append("Check Entities to confirm which people, systems, or organizations dominate the corpus.")
    if not steps:
        steps.append("Confirm the word cloud and top terms look sane before deeper interpretation.")
    return steps[:4]


def render_executive_signal_dashboard(
    scanner: StreamScanner,
    combined_counts: Counter,
    text_stats: Dict,
    insight_df: pd.DataFrame,
):
    st.subheader("🎛️ Executive Signal Dashboard")
    st.caption(
        "A first-pass sensemaking layer: what the corpus seems to be about, where the strongest signals are, "
        "and which deeper views are worth opening next."
    )

    has_temporal = bool(scanner.temporal_counts)
    has_categories = bool(scanner.category_counts)
    high_conf_count = 0
    if not insight_df.empty:
        high_conf_count = int(insight_df["Confidence"].isin(["High", "Medium"]).sum())

    metric_cols = st.columns(5)
    metric_cols[0].metric(
        "Rows",
        f"{text_stats['Total Rows']:,}",
        help=ANALYSIS_HELP_TEXT["Rows"],
    )
    metric_cols[1].metric(
        "Tokens",
        f"{text_stats['Total Tokens']:,}",
        help=ANALYSIS_HELP_TEXT["Tokens"],
    )
    metric_cols[2].metric(
        "Insight Cards",
        f"{len(insight_df):,}",
        help=ANALYSIS_HELP_TEXT["Insight Cards"],
    )
    metric_cols[3].metric(
        "Evidence Snippets",
        f"{len(scanner.evidence_docs):,}",
        help=ANALYSIS_HELP_TEXT["Evidence Snippets"],
    )
    metric_cols[4].metric(
        "Groups / Dates",
        f"{len(scanner.category_counts)} / {len(scanner.temporal_counts)}",
        help=ANALYSIS_HELP_TEXT["Groups / Dates"],
    )

    if scanner.evidence_limit_reached:
        st.warning(
            "Evidence snippets hit the retention cap. Counts still reflect the full scan, "
            "but representative excerpts are based on the retained sample."
        )

    left_col, mid_col, right_col = st.columns([1.05, 1.1, 1.05])

    with left_col:
        st.markdown("#### Strongest Signals")
        if insight_df.empty:
            st.info("Scan more content, or reduce filtering, to generate insight cards.")
        else:
            for _, row in insight_df.head(4).iterrows():
                st.markdown(f"**{row['Signal']}**")
                st.caption(
                    f"{row['Signal Type']} · {row['Confidence']} confidence · "
                    f"support {row['Evidence Strength']}"
                )

    with mid_col:
        st.markdown("#### Signal Mix")
        if insight_df.empty:
            st.caption("No signal mix yet.")
        else:
            mix_df = (
                insight_df["Signal Type"]
                .value_counts()
                .rename_axis("Signal Type")
                .reset_index(name="Cards")
            )
            if alt is not None:
                chart = alt.Chart(mix_df).mark_bar().encode(
                    x=alt.X("Cards:Q", title="Cards"),
                    y=alt.Y("Signal Type:N", sort="-x", title=None),
                    tooltip=[
                        alt.Tooltip("Signal Type:N", title="Signal Type"),
                        alt.Tooltip("Cards:Q", title="Cards"),
                    ],
                ).properties(height=max(160, 28 * len(mix_df)))
                st.altair_chart(chart, use_container_width=True)
            else:
                st.dataframe(
                    mix_df,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Signal Type": st.column_config.TextColumn(
                            "Signal Type",
                            help=ANALYSIS_HELP_TEXT["Signal Type"],
                        ),
                        "Cards": st.column_config.NumberColumn(
                            "Cards",
                            help=ANALYSIS_HELP_TEXT["Cards"],
                        ),
                    },
                )

    with right_col:
        st.markdown("#### What To Do Next")
        for step in build_dashboard_next_steps(scanner, insight_df, has_temporal, has_categories):
            st.write(f"- {step}")
        if high_conf_count:
            st.success(f"{high_conf_count} insight card(s) have medium/high confidence.")
        else:
            st.info("Treat current outputs as early leads until more evidence accumulates.")

    with st.expander("At-a-glance corpus fingerprint", expanded=False):
        top_words = ", ".join([f"{w} ({c})" for w, c in combined_counts.most_common(15)])
        top_entities = ", ".join([f"{e} ({c})" for e, c in scanner.entity_counts.most_common(10)])
        top_phrases = ", ".join([
            f"{w1} {w2} ({c})"
            for (w1, w2), c in scanner.global_bigrams.most_common(10)
        ])
        st.markdown("**Top terms**")
        st.write(top_words or "No terms available.")
        st.markdown("**Top phrases**")
        st.write(top_phrases or "No bigrams available.")
        st.markdown("**Top entities**")
        st.write(top_entities or "No entities detected.")

        if scanner.dashboard_summary:
            st.markdown("**Offline harvester summary**")
            st.json(scanner.dashboard_summary, expanded=False)


def render_workflow_guide():
    with st.expander("📘 Comprehensive App Guide: How to use this Tool", expanded=False):
        st.markdown("""
        ### 🌟 What is Signal Foundry?

        Signal Foundry turns messy text into structured signal. Feed it PDFs, CSVs, transcripts, PowerPoints, pasted notes, URLs, or pre-computed harvester sketches and it will surface recurring language, relationships, themes, entities, maturity signals, a top-level resource shape, key insights, and supporting evidence cards.

        It is designed for fast exploratory sensemaking:

        - What is this corpus about?
        - What terms and phrases actually matter?
        - What seems distinctive, not just frequent?
        - What tensions, needs, blockers, risks, or opportunities appear?
        - What concepts are missing or weak?
        - What should a human analyst inspect next?

        ---

        ### 🚦 Recommended Reading Order

        After scanning, use this order:

        1. **Executive Signal Dashboard**  
           Start here. It summarizes corpus size, evidence coverage, strongest signals, signal-type mix, and suggested next steps.

        2. **Signal Compass**  
           Use this first inside the Insight Engine. It shows the main directional pull of the text, such as risk-centered, institution-heavy, evidence-heavy, or system-dependence-heavy.

        3. **Resource Shape**  
           Read this next inside the Insight Engine. It groups the strongest signals into a short synthesis of what the uploaded resource appears to be about.

        4. **Supporting Insight Cards**  
           Use these to inspect the evidence behind the synthesis. Cards connect statistical signals to interpretation, confidence, representative excerpts, and follow-up questions.

        5. **Word Cloud & Stats**  
           Confirm the scan looks sane. If boilerplate or junk dominates, adjust stopwords and rescan.

        6. **Themes**  
           Inspect theme evidence cards, frequency-vs-distinctiveness, missing expected signals, category contrasts, and temporal drift.

        7. **Keyphrases**  
           Use TF-IDF to find words that are unusually specific to this corpus.

        8. **Entities**  
           Check which people, organizations, systems, programs, places, or named concepts keep appearing.

        9. **Network Graph**  
           Explore relationships between terms after the core signals make sense. The graph controls are capped for safer browser rendering on Streamlit Community Cloud.

        10. **Maturity**  
           Use this when the source material fits one of the maturity lenses.

        11. **AI Analyst**  
           Use last. It works best when the visible dashboard and evidence cards already look reasonable.

        12. **Calibration Export**  
           Admin-only. Use this when testing or comparing runs. The ZIP includes insight cards, Resource Shape files, and weighting diagnostics that help explain why signals ranked where they did.

        ---

        ### 🚀 How to Use It

        #### Path A - Standard Scan

        Best for most users.

        1. Open **Workspace**.
        2. Upload one or more files, paste text/URLs, or load an offline sketch.
        3. Keep **Clear previous data** enabled unless intentionally combining scans.
        4. Leave default cleaning settings on for the first pass.
        5. Click the scan button.
        6. Start with the **Executive Signal Dashboard**, then read **Signal Compass** and **Resource Shape** in the Insight Engine.

        #### Path B - Structured Files

        Best for CSV or Excel files with meaningful columns.

        1. Expand the file configuration panel.
        2. Choose the text column or columns.
        3. Choose a date column if you want trends.
        4. Choose a category column if you want group comparison.
        5. Scan the file.

        Category columns can power contrastive analysis, such as comparing teams, departments, clients, speakers, document types, or phases.

        Date columns can power temporal drift and trend analysis.

        #### Path C - Offline Harvester

        Best for very large files or secure environments.

        Run the harvester outside the app:

        ```bash
        python harvester.py --input data.csv --col text --output sketch.json
        ```

        With richer dashboard support:

        ```bash
        python harvester.py --input data.csv --col text --date-col date --category-col team --output sketch.json
        ```

        Then upload the generated sketch in **Load Offline Analysis (Harvester)**.

        Use `--no-evidence` if you do not want representative excerpts stored in the sketch.

        ---

        ### 💡 How to Read the Insight Engine

        The Insight Engine now has three layers:

        1. **Signal Compass** gives a quick directional read of the uploaded material.
        2. **Resource Shape** gives a short synthesis of the main signal families and strongest key insights.
        3. **Supporting Insight Cards** show the evidence behind that synthesis.

        Each supporting card includes:

        - **Signal:** the phrase or concept being surfaced
        - **Signal Type:** likely analytical category
        - **Evidence Strength:** how much support exists
        - **Distinctiveness:** how much the signal stands out
        - **Confidence:** low, medium, or high
        - **Representative Evidence:** short source excerpts when available
        - **Interpretation:** what the signal may indicate
        - **Follow-up Question:** what a human should ask next

        Signal types include:

        - Pain / Friction
        - Need / Request
        - Blocker / Constraint
        - Aspiration / Opportunity
        - Risk / Concern
        - Decision / Tradeoff
        - Contradiction / Tension
        - Infrastructure / System Dependence
        - Disease / Hazard
        - Evidence / Experiment
        - Intervention / Control Method
        - Public Health / Institutional Response
        - Embodiment / Lived Experience
        - Institutional Structure / Social Design
        - Standardization / Loss of Difference
        - Motif / Image Pattern
        - Low-Specificity Signal
        - Absence / Weak Signal

        The Resource Shape is a first-pass synthesis, not a final conclusion. The supporting cards are analytical leads that help you verify or challenge it.

        ---

        ### 🤖 How to Use the AI Analyst

        The AI Analyst is a second-pass helper. It does **not** read the full raw documents. It receives a privacy-preserving context brief built from:

        - corpus stats
        - top terms and bigrams
        - top detected entities
        - Signal Compass
        - Resource Shape
        - Insight Cards and signal roles
        - graph/community summary when available
        - maturity scores and domains when a maturity assessment has been run

        Use **What the AI can see** to preview the exact briefing before asking the AI anything.

        Good uses:

        - Ask for a plain-language synthesis of the current scan.
        - Ask what the Signal Compass and Resource Shape imply.
        - Ask which evidence cards deserve human review.
        - Ask what maturity domains look strongest or weakest after running Maturity.
        - Ask what questions a human analyst should investigate next.

        Limits:

        - It cannot quote or inspect the full source document unless that text appears in the summarized context.
        - It may miss details that are visible in other tabs but not included in the AI context brief.
        - It should be treated as an analyst aide, not a final judgment.

        ---

        ### 🧪 Calibration Export

        The Calibration Export is password-gated because it is mainly an admin/testing tool.

        Use it when you want to:

        - compare repeated runs of the same document
        - test stopword or concept-check changes
        - preserve a snapshot of an analysis
        - review why a Resource Shape or insight-card ranking changed

        The export includes a `resource_shape_weight_diagnostics.csv` file. This file breaks rankings into readable ingredients such as evidence strength, distinctiveness, confidence weight, interpretive lift, family priority, role weight, and final shape weight.

        Safe exports omit representative evidence from the insight-card CSV. Full diagnostic exports include evidence excerpts and should be shared carefully.

        ---

        ### 🧭 Reading Between the Lines

        Signal Foundry does not magically know intent. It surfaces patterns that help a person interpret the text.

        Useful moves:

        - Read Signal Compass and Resource Shape first.
        - Use Key Insights to decide what deserves attention.
        - Open supporting evidence cards before making claims.
        - Look for high-confidence insight cards after the top-level synthesis.
        - Compare frequent terms against distinctive terms.
        - Enter expected concepts to see what is missing.
        - Use category comparison when different stakeholder groups are present.
        - Use temporal drift when dates are available.
        - Read representative evidence before making claims.
        - Treat surprising absences as questions, not proof.

        Good follow-up questions often sound like:

        - What is causing this friction?
        - Who experiences this blocker most?
        - What need is being repeated but not resolved?
        - What risk is implied but not directly named?
        - What is missing from the language?
        - What does the corpus avoid saying?

        ---

        ### 🕸️ Network Graph Guidance

        The graph is useful, but it is also the most browser-sensitive visualization.

        Start with conservative settings:

        - **Min Link Frequency:** raise this to remove weak connections.
        - **Max Nodes:** lower this if the graph is too dense.
        - **Max Links:** lower this if rendering becomes unstable.
        - **Physics:** turn off if the graph jumps around too much.
        - **GEXF Export:** use this for heavier graph analysis in external tools.

        Suggested presets:

        - **Safe:** Min Link 5+, Max Nodes 40, Max Links 60
        - **Balanced:** Min Link 3-5, Max Nodes 55, Max Links 90
        - **Dense:** Min Link 2-3, Max Nodes 80+, Max Links 150+

        If the graph fails to render, lower nodes and links first.

        ---

        ### 🏆 Maturity Model Guidance

        Use maturity scoring when the source material actually contains maturity-relevant language.

        Strong sources:

        - meeting transcripts
        - coaching notes
        - strategic reviews
        - policy discussions
        - operational retrospectives
        - implementation planning notes

        Weaker sources:

        - technical logs
        - boilerplate documents
        - short excerpts
        - OCR-heavy PDFs
        - documents with little organizational language

        Important: maturity scoring measures language signals, not verified organizational reality. Treat it as a conversation starter.

        ---

        ### 🧼 Data Hygiene Tips

        If results look odd:

        - Check whether **Clear previous data** was off accidentally.
        - Add recurring junk terms to stopwords.
        - Lower **Min Word Len** if acronyms matter.
        - Raise **Min Word Len** if short junk words dominate.
        - Keep **Bigrams** enabled for phrase detection.
        - Use structured columns when possible.
        - For OCR-heavy PDFs, expect some noisy terms.
        - Re-scan after changing important cleaning settings.

        ---

        ### 🔐 Privacy Notes

        The app can analyze raw uploaded text inside the Streamlit session.

        The AI Analyst is designed to use statistical sketches and insight-card context rather than full raw documents.

        Offline harvester sketches may include representative excerpts unless generated with:

        ```bash
        --no-evidence
        ```

        Use `--no-evidence` for highly sensitive corpora.

        ---

        ### ✅ Practical Testing Checklist

        For a new deployment or code update:

        1. Load the app.
        2. Upload one small TXT or PDF.
        3. Scan with **Clear previous data** enabled.
        4. Confirm the Executive Signal Dashboard appears.
        5. Confirm Insight Engine cards appear.
        6. Confirm Word Cloud and Frequency Tables render.
        7. Test the Network Graph with conservative settings.
        8. Test offline sketch loading from the harvester.
        9. Test AI Analyst last.

        Bottom line: start with the dashboard, inspect evidence, then drill into the deeper tools.
        """)

def render_lit_case_study():
    # We use Unicode "Math Sans" characters to simulate bold/italics in the title
    # Italic 'another': 𝘢𝘯𝘰𝘵𝘩𝘦𝘳
    # Bold 'specific': 𝘀𝗽𝗲𝗰𝗶𝗳𝗶𝗰
    title = "🔦 Spotlight: Digital Humanities & Literary Forensics (𝘢𝘯𝘰𝘵𝘩𝘦𝘳 𝘀𝗽𝗲𝗰𝗶𝗳𝗶𝗰 Case Study)"
    
    with st.expander(title, expanded=False):
        st.markdown("""
        ### The Scenario
        **The Artifact:** The full text of Ovid's **<a href="https://www.gutenberg.org/files/21765/21765-h/21765-h.htm" target="_blank">"Metamorphoses"</a>** (via Project Gutenberg URL).
        **The User:** A Digital Humanities Researcher or Student.
        **The Goal:** To rapidly map the "Pantheon" of characters and distinguish the original narrative from the translator's artifacts.

        ---

        ### 1. The "Pantheon Map" (Entities Tab)
        *   **The Question:** "Who are the dominant power players in this 15-book epic?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** The engine immediately surfaces **"Jupiter," "Apollo," "Ceres,"** and **"Minerva"** as the top nodes.
        *   **The Value / Insight:** Without reading a single line, you have a hierarchical map of the Roman deities driving the plot.

        ### 2. The "Translator's Fingerprint" (NPMI & Bigrams)
        *   **The Question:** "Is this pure text, or is there structural noise?"
        *   **The Signal:** Sticky Concepts (Bigrams).
        *   **The Result:** The engine identifies **"Clarke translates"** and **"-ver Clarke"** as top phrases.
        *   **The Value / Insight:** **Forensic Separation.** The engine detected that *John Clarke* (the translator) is statistically inseparable from the text. It highlights "Data Hygiene" issues—showing you exactly what "boilerplate" needs to be cleaned (e.g., "Project Gutenberg" headers) before deep analysis.

        ### 3. The "Narrative Arcs" (Topic Modeling)
        *   **The Question:** "What are the distinct recurring themes?"
        *   **The Signal:** NMF/LDA Mathematical Bucketing.
        *   **The Result:**
            *   **Topic A:** [Daughter, Jupiter, Cadmus, Wife] -> *The Genealogy & Origin Myths.*
            *   **Topic B:** [Thou, Thee, Thus, Said] -> *The Dialogue & Poetic Structure.*
        *   **The Value / Insight:** The engine successfully separates the *Style* (Archaic English) from the *Substance* (Mythological Events).

        ### 4. The "Semantic Network" (Graph Tab)
        *   **The Question:** "How do the main characters interact?"
        *   **The Signal:** Proximity-based linking.
        *   **The Result:** "jupiter" is the central "hub" node, with spokes connecting to various "nymphs" and "daughters."
        *   **The Value / Insight:** Visualizes the centralized power structure of the mythology, confirming Jupiter as the primary driver of the transformations.
        """, unsafe_allow_html=True)

def render_auto_insights(scanner, proc_conf):
    # Only run if we have data
    if not scanner.global_counts: return

    # --- 1. PREPARE DATA ---
    # Entities
    top_ents = scanner.entity_counts.most_common(3)
    ent_str = ", ".join([f"**{e[0]}**" for e in top_ents]) if top_ents else "(No entities detected)"
    
    # Sticky Concepts (NPMI)
    df_npmi = calculate_npmi(scanner.global_bigrams, scanner.global_counts, scanner.total_rows_processed)
    top_npmi = df_npmi.head(3)["Bigram"].tolist() if not df_npmi.empty else []
    npmi_str = ", ".join([f"**{b}**" for b in top_npmi]) if top_npmi else "(No strong phrases found)"
    
    # Tech Signal (TF-IDF)
    df_tfidf = calculate_tfidf(scanner, 20)
    top_idf = df_tfidf.head(3)["Term"].tolist() if not df_tfidf.empty else []
    idf_str = ", ".join([f"**{t}**" for t in top_idf]) if top_idf else "(Not enough documents for TF-IDF)"

    # -render reporting
    with st.expander("⚡ High-Level Signal Report (Auto-Generated, ymmv)", expanded=True):
        st.markdown(f"""
        ### 1. The "Stakeholder Map" (Entities)
        *   **The Question:** "Who are the dominant actors or organizations?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** {ent_str}
        *   **The Insight:** These nodes appear most frequently, suggesting they are the primary drivers of the narrative or the key subjects of the file.

        ### 2. The "Sticky Concepts" (Phrase Significance)
        *   **The Question:** "What is the specific 'Term of Art' or jargon here?"
        *   **The Signal:** NPMI (Normalized Pointwise Mutual Information).
        *   **The Result:** {npmi_str}
        *   **The Insight:** These words appear together mathematically more often than random chance, indicating they represent specific concepts (e.g. "Credit Card" vs "Red Card") rather than generic language.

        ### 3. The "Technical Signal" (Keyphrases)
        *   **The Question:** "What makes this specific document unique?"
        *   **The Signal:** TF-IDF (Inverse Document Frequency).
        *   **The Result:** {idf_str}
        *   **The Insight:** While words like "the" or "report" might be frequent, *these* specific words are statistically unique to this dataset, representing its core technical signature.
        """)

def render_neurotech_case_study():
    with st.expander("🔦 Spotlight: Analyzing Mi|itary Neurotechno|ogy (a very *specific* Case Study)", expanded=False):
        st.markdown("""
        ### The Scenario
        **The Artifact:** A dense, 50-page UNIDIR report titled <a href="https://unidir.org/wp-content/uploads/2025/11/UNIDIR_Neurotechnology_Military-Domain_A-Primer.pdf" target="_blank"><b>"Neurotechnology in the Military Domain"</b></a>.
        **The User:** A Defense Ana|yst with 5 minutes to extract actionable inte||igence.
        **The Goal:** Move beyond "what is this paper about?" to "what are the threats and opportunities?"

        ---

        ### 1. The "Sticky Concepts" (NPMI Tab)
        *   **The Question:** "What specific types of risks are discussed?"
        *   **The Signal:** The engine finds words that mathematically *stick together* more than random chance.
        *   **The Result:** It surfaces **"Dua| Use"** and **"Cognitive Liberty."**
        *   **The Insight:** The strategic risk isn't just new weap0ns; it is civi|ian medica| techno|ogy being repurposed for mi|itary app|ications (Dua| Use), necessitating a legal/ethical framework (Liberty).

        ### 2. The "Technical Signal" (Keyphrases Tab)
        *   **The Question:** "Do I need to worry about brain implants yet?"
        *   **The Signal:** TF-IDF filters out generic words to find unique technical terms.
        *   **The Result:** High scores for **"Non-invasive,"** **"Transcranial,"** and **"Wearable."**
        *   **The Insight:** The immediate operationa| reality is external headsets/helmets, not surgical implants.

        ### 3. The "Semantic Network" (Graph Tab)
        *   **The Question:** "How is the techno|ogy being applied?"
        *   **The Signal:** The Graph links words based on proximity in the text.
        *   **The Result:** 
            *   Cluster A links **"Stimulation"** to **"Performance"** (Enhancement/Super-Soldiers).
            *   Cluster B links **"Stimulation"** to **"Interrogation"** (Weap0nization/T0rture).
        *   **The Insight:** The paper treats "Enhancement" and "Weaponization" as distinct operational clusters.

        ### 4. The "Stakeholder Map" (Entities Tab)
        *   **The Question:** "Who is involved?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** **"D@RPA," "Neura|ink," "Geneva Convention," "Human Rights Council."**
        *   **The Insight:** Identifies the funding sources (DARP@) vs. the regulatory blockers (Geneva).
        """, unsafe_allow_html=True)

def render_maturity_guide():
    with st.expander("🏆 Guide: Understanding the Maturity Models", expanded=False):
        st.markdown("""
        ### What is the Maturity Engine?

        While standard analytics count *what* words appear, the Maturity Engine measures the **intent and capability** behind those words. It compares your text against known frameworks of organizational development.

        ---

        ### 🎭 The Personas (Context Matters!)

        Words change meaning based on context. You must select the model that matches your document's intent.

        #### 1. 🏫 EdTech & LMS Ops (5-Level)
        *Best for: course designs, syllabi, LMS usage reviews, instructional strategy documents.*
        *   **L1 (Repository):** Static content dumps and basic file hosting.
        *   **L3 (Integrated):** Connected systems, tools, APIs, and interoperability.
        *   **L5 (Strategic):** Institutional thinking around equity, accessibility, and transformation.

        #### 2. 🏢 General Business Ops (5-Level, CMMI)
        *Best for: project updates, business operations notes, internal reports, emails.*
        *   **L1 (Reactive):** Urgent, manual, firefighting language.
        *   **L4 (Measured):** KPI, ROI, forecasting, and quantitative management language.

        #### 3. ⚖️ Policy & Governance (5-Level)
        *Best for: policy documents, compliance reviews, NGO/government publications.*
        *   **L1 (Enforcement):** Violations, bans, restrictions, sanctions.
        *   **L5 (Systemic):** Sustainable, resilient, holistic, ecosystem-level language.

        #### 4. 🎓 TAM Maturity Model — 12-Domain (3-Tier)
        *Best for: meeting transcripts, coaching notes, strategic reviews, TAM/client discussions.*

        This persona evaluates the client across **12 separate domains**, each scored on 3 tiers:

        | Tier | Label | Score Range | Meaning |
        |------|-------|-------------|---------|
        | 🔴 | **Foundational** | 1.0–1.49 | Basic, manual, reactive |
        | 🟠 | **Advanced** | 1.50–2.49 | Structured, proactive, repeatable |
        | 🟢 | **Leading Edge** | 2.50–3.0 | Strategic, scalable, innovation-oriented |

        **The 12 Domains**
        1. Platform & Technical Administration
        2. Curriculum Development & Delivery
        3. Student Engagement & Success
        4. Data & Learning Analytics
        5. Assessment & Evaluation
        6. Instructor Efficiency
        7. Change Management
        8. Knowledge & Resource Management
        9. Accessibility & Compliance
        10. User Support & Training
        11. Innovation & Emerging Technologies
        12. Collaboration & Communication

        ---

        ### 🧮 How Scoring Works

        **For the 5-level personas**
        *   Weighted average across the detected maturity-level signals.
        *   Score range: **1.0 to 5.0**

        **For the 12-domain TAM persona**
        *   Each domain is scored independently using tiers 1-3.
        *   The composite score is the average of all domains with at least one signal.
        *   Score range: **1.0 to 3.0**
        *   Domains with zero signal are shown as **No Data** and excluded from the composite.

        **Phrase weighting**
        *   Multi-word phrases receive extra weight relative to single words because they are usually more specific signals.

        ---

        ### ⚙️ Important Processing Note

        Maturity scoring uses the same cleaned token stream as the rest of the app. That means settings like:
        *   **Min Word Len**
        *   **Stopwords**
        *   **Remove Generic Filler Words and Prepositions**
        *   **Use Lemmatization**
        *   **Keep Hyphens**
        *   **Chat / HTML / URL cleanup**

        can all affect the maturity result.

        ---

        ### 🕸️ How to Read the Charts

        **5-Level Radar**
        *   Round shape = balanced maturity
        *   Sharp spikes = uneven capability

        **12-Domain Radar**
        *   One spoke per domain
        *   Outer edge = stronger maturity signal
        *   Uneven shape = domain-specific strengths and gaps

        **12-Domain Breakdown Bars**
        *   Red = Foundational share
        *   Orange = Advanced share
        *   Green = Leading Edge share

        ---

        ### ⚠️ Caveats

        > **Signal vs. reality:** The maturity engine measures language, not actual implementation. Treat it as a conversation starter and evidence aid, not a final audit.

        > **Persona mismatch:** If you scan a generic business strategy memo with the TAM model, the score may look plausible but be contextually wrong.

        > **Data quality matters:** Transcripts, coaching notes, and strategy discussions usually work better than exported system logs or boilerplate-heavy files.
        """)

def render_use_cases():
    with st.expander("📖 Playbook: High-Value Use Cases", expanded=False):
        st.markdown("""
        ### 🎓 Specialized: Education & EdTech
        
        #### 1. The "Syllabus Audit" (Maturity Scan)
        *   **Goal:** Determine if a department is truly "Modernizing" or just digitizing old habits.
        *   **Action:** Bulk scan 50 syllabi/course descriptions using the **EdTech Persona**.
        *   **Signal:**
            *   High **L1 (Repository)** scores = The LMS is just a file dump.
            *   High **L3 (Integrated)** scores = They are actually using the ecosystem tools.
        
        #### 2. Vendor RFP Analysis
        *   **Goal:** Cut through sales fluff.
        *   **Action:** Scan a vendor's whitepaper or proposal.
        *   **Signal:** Do they speak **L2 (Features/Tools)** or **L5 (Partnership/Success)**? A strategic partner should score >3.5.
        
        ---

        ### 🏢 General: Corporate & Intelligence
        
        #### 3. Strategic Alignment Check
        *   **Goal:** See if the IT Department aligns with the C-Suite.
        *   **Action:** Scan the CTO's emails vs. the CEO's annual letter.
        *   **Signal:** If the CEO is L5 (Vision) and IT is L1 (Firefighting), you have an execution gap.

        #### 4. "Chain of Custody" Verification
        *   **Goal:** Prove that a visualization came from *this* specific file.
        *   **Action:** Generate the **Hybrid Signature** (QR Code) in the Graph tab.
        *   **Value:** Encrypts the document hash into the image. If the screenshot leaks, you can prove exactly which source file generated it.

        #### 5. Crisis Timeline Reconstruction
        *   **Goal:** Pinpoint when a project went off the rails.
        *   **Action:** Scan weekly status reports with the **Time-Travel Slider**.
        *   **Signal:** Watch for the crossover point where **"Plan" (L2)** words drop and **"Fix/Urgent" (L1)** words spike.

        ---

        ### 🎓 Specialized: TAM & Client Maturity Assessment

        #### 6. The "Client Maturity Snapshot" (12-Domain Scan)
        *   **Goal:** Quickly gauge where a client stands across all 12 TAM admin maturity domains.
        *   **Action:** Upload 3–5 meeting transcripts or coaching session recordings (VTT/PDF) from a single client. Select the **🎓 TAM Maturity Model** persona.
        *   **Signal:**
            *   The **Radar Chart** instantly shows which domains the client discusses at a strategic level vs. which are still basic.
            *   The **Breakdown Chart** reveals the distribution — is the client *mostly* Foundational with a few Advanced pockets, or broadly Advanced with Leading Edge potential?
        *   **Action:** Use the domain cards to identify the 2–3 domains with the lowest scores and build your coaching plan around them.
        *   **Processing note:** Keep **Bigrams** on, and if acronym-heavy terminology is common, lower **Min Word Len** before scanning.

        #### 7. Client Progress Over Time (Longitudinal Tracking)
        *   **Goal:** Measure whether a client's maturity is improving after coaching.
        *   **Action:** Export a snapshot after each major engagement (quarterly reviews, coaching milestones). Later, upload all snapshots to the Progress Tracking section.
        *   **Signal:** Compare the composite scores and per-domain scores across dates. Are the reds turning orange? Are the oranges turning green?
        *   **Value:** Concrete, data-backed evidence of TAM impact for quarterly reviews and stakeholder reporting.

        #### 8. Cross-Client Benchmarking
        *   **Goal:** Identify which clients need the most attention.
        *   **Action:** Scan transcripts from 3–4 different clients separately. Note each client's composite score and lowest-scoring domains.
        *   **Signal:** Client A scores 2.4 composite, Client B scores 1.3. Client B needs more coaching. Client A might be ready for Leading Edge challenges.
        *   **Value:** Helps TAMs and OS leadership prioritize resource allocation across the client portfolio.
        """)

def render_analyst_help():
    with st.expander("🎓 Analyst's Guide & Troubleshooting", expanded=False):
        st.markdown("""
        ### Quick Diagnostic Rule

        If the **Word Cloud** and **Frequency Tables** look wrong, the downstream outputs will also be wrong. Fix cleaning, stopwords, file columns, or scan mode first.

        ---

        **Symptom: The app output feels wrong right away**
        * **Fix:** Confirm you scanned the correct file, sheet, and text column.
        * **Fix:** Check whether **Clear previous data** was off and old data was accidentally merged into the current scan.
        * **Fix:** Look at the top frequency table. If the top words are boilerplate, add them to **Stopwords** and rescan.

        **Symptom: Important acronyms disappeared**
        * **Likely cause:** **Min Word Len** is too high.
        * **Fix:** Lower **Min Word Len** to 2 or 3 for acronym-heavy corpora such as API, LTI, SSO, AI, SIS, or WCAG.

        **Symptom: Maturity results feel too low or missing domains**
        * **Fix:** Confirm the selected maturity persona matches the source material.
        * **Fix:** Lower **Min Word Len** if important short terms are being filtered out.
        * **Fix:** Keep **Bigrams** enabled so phrase signals like "student success" and "change management" can be detected.
        * **Fix:** Feed the model more relevant meeting notes, transcripts, coaching notes, or strategic review material.
        * **Interpretation:** A **No Data** domain means no mapped language was detected for that domain. It does not prove the client lacks that capability.

        **Symptom: Maturity result feels too generic**
        * **Fix:** Use the domain detail cards. The top linguistic drivers tell you exactly why the score was produced.
        * **Fix:** If drivers look like boilerplate, add that boilerplate to **Stopwords** and rescan.
        * **Fix:** If a meaningful client-specific phrase is missing, the maturity vocabulary may need expansion.

        **Symptom: Topics look like gibberish or random words**
        * **Fix:** Check **Rows per Doc**.
            * For chats, tickets, comments, and transcripts, use 1–5.
            * For reports, PDFs, books, and long documents, use 100+.
        * **Fix:** Try **NMF** for shorter, cleaner records.
        * **Fix:** Try **LDA** for longer mixed-content documents.
        * **Fix:** Add recurring boilerplate terms to **Stopwords**.

        **Symptom: The Network Graph is a giant blob**
        * **Fix:** Increase **Min Link Frequency** to cut weak connections.
        * **Fix:** Increase **Repulsion** to push nodes apart.
        * **Fix:** Lower **Max Nodes** to focus on the most important terms.

        **Symptom: The Graph has disconnected islands**
        * **Fix:** Decrease **Min Link Frequency** to reveal subtler connections.
        * **Fix:** Increase **Edge Length** to give clusters room to breathe.

        **Symptom: Seeing duplicates such as "run" and "running"**
        * **Fix:** Enable **Use Lemmatization**. This can merge word variations into a shared root form.

        **Symptom: High-ranking words are boring, such as "page", "copyright", or "transcript"**
        * **Fix:** These are corpus artifacts. Add them to **Stopwords** and rescan.
        * **Fix:** For transcripts, keep **Remove Chat Artifacts** enabled.

        **Symptom: Trend charts are empty**
        * **Fix:** Re-scan after selecting the correct date column in the file config panel.
        * **Fix:** Check whether source dates are parseable, such as 2025-01-31, 01/31/2025, or Jan 31 2025.

        **Symptom: AI Analyst answer seems too vague**
        * **Fix:** The AI Analyst only sees the statistical sketch, not raw documents.
        * **Fix:** Ask narrower questions about visible outputs, such as "Which maturity domains look weakest?" or "Which graph clusters appear most central?"
        """)
# visualization helpers
@st.cache_data(show_spinner="Analyzing term sentiment...")
def get_sentiments(_analyzer, terms: Tuple[str, ...]) -> Dict[str, float]:
    if not _analyzer or not terms: return {}
    return {term: _analyzer.polarity_scores(term)['compound'] for term in terms}

def create_sentiment_color_func(sentiments: Dict[str, float], pos_color, neg_color, neu_color, pos_thresh, neg_thresh):
    def color_func(word, font_size, position, orientation, random_state=None, **kwargs):
        score = sentiments.get(word, 0.0)
        if score >= pos_thresh: return pos_color
        elif score <= neg_thresh: return neg_color
        else: return neu_color
    return color_func

def get_sentiment_category(score: float, pos_threshold: float, neg_threshold: float) -> str:
    if score >= pos_threshold: return "Positive"
    if score <= neg_threshold: return "Negative"
    return "Neutral"

def build_wordcloud_figure_from_counts(counts, max_words, width, height, bg_color, colormap, font_path, random_state, color_func):
    limited = dict(counts.most_common(max_words))
    if not limited: return plt.figure(), None
    wc = WordCloud(width=width, height=height, background_color=bg_color, colormap=colormap, font_path=font_path, random_state=random_state, color_func=color_func, collocations=False, normalize_plurals=False).generate_from_frequencies(limited)
    fig, ax = plt.subplots(figsize=(max(6, width/100), max(3, height/100)), dpi=100)
    ax.imshow(wc, interpolation="bilinear"); ax.axis("off"); plt.tight_layout()
    return fig, wc

def fig_to_png_bytes(fig):
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0.1)
    buf.seek(0)
    return buf

def dataframe_to_csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8-sig")

def make_json_safe(value: Any) -> Any:
    if isinstance(value, dict):
        return {str(k): make_json_safe(v) for k, v in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [make_json_safe(v) for v in value]
    if isinstance(value, Counter):
        return {str(k): make_json_safe(v) for k, v in value.items()}
    if isinstance(value, (np.integer,)):
        return int(value)
    if isinstance(value, (np.floating,)):
        return float(value)
    if pd.isna(value) if not isinstance(value, (dict, list, tuple, set)) else False:
        return None
    return value

def json_to_bytes(payload: Dict[str, Any]) -> bytes:
    return json.dumps(make_json_safe(payload), indent=2, ensure_ascii=False).encode("utf-8")

def top_terms_dataframe(counts: Counter, top_n: int = 100) -> pd.DataFrame:
    return pd.DataFrame(
        [{"Term": term, "Count": count} for term, count in counts.most_common(top_n)]
    )

def top_bigrams_dataframe(bigrams: Counter, top_n: int = 100) -> pd.DataFrame:
    return pd.DataFrame(
        [
            {"Bigram": f"{w1} {w2}", "Count": count}
            for (w1, w2), count in bigrams.most_common(top_n)
        ]
    )

def counter_dataframe(counter: Counter, key_name: str, value_name: str = "Count", top_n: int = 100) -> pd.DataFrame:
    return pd.DataFrame(
        [{key_name: key, value_name: value} for key, value in counter.most_common(top_n)]
    )

def signal_type_distribution_dataframe(insight_df: pd.DataFrame) -> pd.DataFrame:
    if insight_df.empty or "Signal Type" not in insight_df.columns:
        return pd.DataFrame(columns=["Signal Type", "Cards"])
    return (
        insight_df["Signal Type"]
        .value_counts()
        .rename_axis("Signal Type")
        .reset_index(name="Cards")
    )

def signal_role_distribution_dataframe(insight_df: pd.DataFrame) -> pd.DataFrame:
    if insight_df.empty or "Signal Role" not in insight_df.columns:
        return pd.DataFrame(columns=["Signal Role", "Cards"])
    return (
        insight_df["Signal Role"]
        .value_counts()
        .rename_axis("Signal Role")
        .reset_index(name="Cards")
    )

def resource_shape_to_dataframes(resource_shape: Dict[str, Any]) -> Tuple[pd.DataFrame, pd.DataFrame]:
    families_df = pd.DataFrame(resource_shape.get("families", []))
    key_insights_df = pd.DataFrame(resource_shape.get("key_insights", []))
    return families_df, key_insights_df

def evidence_dataframe(scanner: StreamScanner, include_excerpts: bool) -> pd.DataFrame:
    if not include_excerpts:
        return pd.DataFrame(columns=["Evidence Snippets", "Evidence Limit Reached"])
    rows = []
    for item in scanner.evidence_docs:
        rows.append({
            "ID": item.get("id"),
            "Date": item.get("date"),
            "Category": item.get("category"),
            "Excerpt": item.get("excerpt"),
            "Unique Token Preview": ", ".join(item.get("tokens", [])[:40]),
        })
    return pd.DataFrame(rows)

def calibration_notes_markdown(
    scanner: StreamScanner,
    insight_df: pd.DataFrame,
    resource_shape: Dict[str, Any],
    export_mode: str,
    run_label: str,
) -> str:
    signal_mix = signal_type_distribution_dataframe(insight_df)
    role_mix = signal_role_distribution_dataframe(insight_df)
    top_signal_types = []
    if not signal_mix.empty:
        top_signal_types = [
            f"{row['Signal Type']} ({row['Cards']})"
            for _, row in signal_mix.head(5).iterrows()
        ]
    top_roles = []
    if not role_mix.empty:
        top_roles = [
            f"{row['Signal Role']} ({row['Cards']})"
            for _, row in role_mix.head(5).iterrows()
        ]

    low_specificity_count = 0
    if not insight_df.empty and "Signal Type" in insight_df.columns:
        low_specificity_count = int(
            (insight_df["Signal Type"] == "Low-Specificity Signal").sum()
        )

    lines = [
        "# Signal Foundry Calibration Export",
        "",
        f"- Run label: {run_label or 'Not provided'}",
        f"- Export mode: {export_mode}",
        f"- Created: {datetime.now().isoformat(timespec='seconds')}",
        f"- Rows processed: {scanner.total_rows_processed:,}",
        f"- Evidence snippets retained: {len(scanner.evidence_docs):,}",
        f"- Evidence cap reached: {scanner.evidence_limit_reached}",
        f"- Insight cards: {len(insight_df):,}",
        f"- Low-specificity cards: {low_specificity_count:,}",
        "",
        "## Resource Shape",
        "",
        resource_shape.get("summary", "No Resource Shape summary available."),
        "",
        "## Signal Mix",
        "",
        ", ".join(top_signal_types) if top_signal_types else "No signal mix available.",
        "",
        "## Signal Roles",
        "",
        ", ".join(top_roles) if top_roles else "No signal role mix available.",
        "",
        "## Calibration Review Prompts",
        "",
        "1. Do the top signal families match a human reading of the source?",
        "2. Are generic connector phrases being demoted into Low-Specificity Signal?",
        "3. Are concrete systems, institutions, bodies, experiences, and ideas being separated cleanly?",
        "4. Are the strongest cards useful as a starting point for interpretation?",
        "5. Did changing sidebar settings materially improve or degrade the output?",
    ]
    return "\n".join(lines)

def build_calibration_export_zip(
    scanner: StreamScanner,
    combined_counts: Counter,
    insight_df: pd.DataFrame,
    text_stats: Dict[str, Any],
    proc_conf: ProcessingConfig,
    export_mode: str,
    run_label: str = "",
    expected_terms_raw: str = "",
) -> bytes:
    include_excerpts = export_mode == "Full diagnostic export"
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    resource_shape = build_resource_shape(insight_df)
    families_df, key_insights_df = resource_shape_to_dataframes(resource_shape)

    insight_export_df = insight_df.copy()
    if not include_excerpts and "Representative Evidence" in insight_export_df.columns:
        insight_export_df = insight_export_df.drop(columns=["Representative Evidence"])

    shape_diagnostics_df = resource_shape_diagnostics_dataframe(insight_df)
    total_words = max(sum(combined_counts.values()), 1)
    npmi_df = calculate_npmi(scanner.global_bigrams, combined_counts, total_words)
    tfidf_df = calculate_tfidf(scanner, top_n=100)
    evidence_df = evidence_dataframe(scanner, include_excerpts)

    settings_payload = {
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "run_label": run_label,
        "export_mode": export_mode,
        "includes_evidence_excerpts": include_excerpts,
        "app_limits": {
            "max_file_size_mb": MAX_FILE_SIZE_MB,
            "max_topic_docs": MAX_TOPIC_DOCS,
            "max_evidence_docs": MAX_EVIDENCE_DOCS,
            "max_evidence_chars": MAX_EVIDENCE_CHARS,
        },
        "scan_totals": {
            "rows_processed": scanner.total_rows_processed,
            "total_tokens": sum(combined_counts.values()),
            "unique_terms": len(combined_counts),
            "topic_docs_sampled": len(scanner.topic_docs),
            "evidence_snippets": len(scanner.evidence_docs),
            "evidence_limit_reached": scanner.evidence_limit_reached,
            "temporal_buckets": len(scanner.temporal_counts),
            "category_buckets": len(scanner.category_counts),
            "entities": len(scanner.entity_counts),
        },
        "processing_settings": {
            "min_word_len": proc_conf.min_word_len,
            "drop_integers": proc_conf.drop_integers,
            "compute_bigrams": proc_conf.compute_bigrams,
            "use_lemmatization": proc_conf.use_lemmatization,
            "stopword_count": len(proc_conf.stopwords),
            "excluded_speaker_count": len(proc_conf.excluded_speakers),
            "partial_speaker_match": proc_conf.partial_speaker_match,
        },
        "hypothesis_concept_check": expected_terms_raw,
        "text_stats": text_stats,
        "resource_shape_summary": resource_shape.get("summary"),
        "resource_shape_weighting": {
            "base_signal_score": "evidence_strength + distinctiveness*2 + confidence_weight + interpretive_lift/10",
            "shape_weight": "base_signal_score * family_priority * role_weight * tension_weight",
            "purpose": "Directional ranking transparency for calibration review; not a formal statistical score.",
        },
    }

    notes_md = calibration_notes_markdown(
        scanner=scanner,
        insight_df=insight_df,
        resource_shape=resource_shape,
        export_mode=export_mode,
        run_label=run_label,
    )

    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("calibration_run_summary.json", json_to_bytes(settings_payload))
        zf.writestr("calibration_notes.md", notes_md.encode("utf-8"))
        zf.writestr("insight_cards.csv", dataframe_to_csv_bytes(insight_export_df))
        zf.writestr("resource_shape.json", json_to_bytes(resource_shape))
        zf.writestr("resource_shape_families.csv", dataframe_to_csv_bytes(families_df))
        zf.writestr("resource_shape_key_insights.csv", dataframe_to_csv_bytes(key_insights_df))
        zf.writestr("resource_shape_weight_diagnostics.csv", dataframe_to_csv_bytes(shape_diagnostics_df))
        zf.writestr("signal_type_distribution.csv", dataframe_to_csv_bytes(signal_type_distribution_dataframe(insight_df)))
        zf.writestr("signal_role_distribution.csv", dataframe_to_csv_bytes(signal_role_distribution_dataframe(insight_df)))
        zf.writestr("text_stats.csv", dataframe_to_csv_bytes(pd.DataFrame([text_stats])))
        zf.writestr("top_terms.csv", dataframe_to_csv_bytes(top_terms_dataframe(combined_counts, 150)))
        zf.writestr("top_bigrams.csv", dataframe_to_csv_bytes(top_bigrams_dataframe(scanner.global_bigrams, 150)))
        zf.writestr("top_entities.csv", dataframe_to_csv_bytes(counter_dataframe(scanner.entity_counts, "Entity", top_n=100)))
        zf.writestr("npmi_phrases.csv", dataframe_to_csv_bytes(npmi_df.head(150)))
        zf.writestr("tfidf_keyphrases.csv", dataframe_to_csv_bytes(tfidf_df.head(150)))
        zf.writestr("evidence_snippets.csv", dataframe_to_csv_bytes(evidence_df))
        zf.writestr(
            "README.txt",
            (
                "Signal Foundry calibration export.\n\n"
                "Use this package to compare repeated runs while tuning stopwords, "
                "concept checks, signal calibration, and scan settings.\n\n"
                "resource_shape_weight_diagnostics.csv shows the component scores "
                "behind Resource Shape ranking so unexpected rankings are easier to debug.\n\n"
                "Safe exports omit representative evidence from insight_cards.csv. "
                "Full diagnostic exports include evidence excerpts and should be shared carefully.\n"
            ).encode("utf-8"),
        )
    zip_buffer.seek(0)
    return zip_buffer.getvalue()

def render_calibration_export_panel(
    scanner: StreamScanner,
    combined_counts: Counter,
    insight_df: pd.DataFrame,
    text_stats: Dict[str, Any],
    proc_conf: ProcessingConfig,
    expected_terms_raw: str = "",
    key_prefix: str = "calibration_export",
):
    st.divider()
    with st.expander("🧪 Admin Calibration Export", expanded=False):
        st.caption(
            "Password-gated diagnostic export for calibration testing. "
            "Use this when comparing repeated runs or sharing results for review."
        )
        export_mode = st.radio(
            "Export mode",
            ["Safe calibration export", "Full diagnostic export"],
            index=0,
            horizontal=False,
            key=f"{key_prefix}_mode",
            help=(
                "Safe mode omits representative evidence from insight cards. "
                "Full mode includes evidence snippets and should be treated as sensitive."
            ),
        )
        run_label = st.text_input(
            "Optional run label",
            placeholder="Example: machine-stops-clean-txt-v4-default-settings",
            key=f"{key_prefix}_run_label",
            help="Adds a human-readable label to the export summary.",
        )
        if export_mode == "Full diagnostic export":
            st.warning(
                "Full diagnostic export includes representative evidence excerpts. "
                "Use only with anonymized or shareable material."
            )

        zip_bytes = build_calibration_export_zip(
            scanner=scanner,
            combined_counts=combined_counts,
            insight_df=insight_df,
            text_stats=text_stats,
            proc_conf=proc_conf,
            export_mode=export_mode,
            run_label=run_label,
            expected_terms_raw=expected_terms_raw,
        )
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_label = re.sub(r"[^a-zA-Z0-9_-]+", "-", run_label.strip()).strip("-").lower()
        label_part = f"_{safe_label}" if safe_label else ""
        st.download_button(
            "📦 Download calibration ZIP",
            zip_bytes,
            f"signal_foundry_calibration{label_part}_{timestamp}.zip",
            "application/zip",
            key=f"{key_prefix}_download",
            help="Downloads a ZIP containing CSV, JSON, and Markdown files for calibration review.",
        )

# 🤖 AI logic
def call_llm_and_track_cost(system_prompt: str, user_prompt: str, config: dict):
    try:
        client = openai.OpenAI(api_key=config['api_key'], base_url=config['base_url'])
        response = client.chat.completions.create(
            model=config['model_name'],
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
        )
        in_tok = 0
        out_tok = 0
        if hasattr(response, 'usage') and response.usage:
            in_tok = response.usage.prompt_tokens
            out_tok = response.usage.completion_tokens
        
        cost = (in_tok * config['price_in'] / 1_000_000) + (out_tok * config['price_out'] / 1_000_000)
        st.session_state['total_tokens'] += (in_tok + out_tok)
        st.session_state['total_cost'] += cost
        return response.choices[0].message.content
    except Exception as e:
        return f"AI Error: {str(e)}"


# 🚀 main app ui
# ==========================================

st.set_page_config(page_title="Signal Foundry", layout="wide")
st.toast("app loaded/updated successfully", icon="🚀") # cache buster
st.title("🧠 Signal Foundry: Unstructured Data Analytics")
st.markdown("### *(or: dirty-data geiger counter~)*")

if "privacy_notice_dismissed" not in st.session_state:
    st.session_state["privacy_notice_dismissed"] = False

if not st.session_state["privacy_notice_dismissed"]:
    with st.container(border=True):
        st.warning(
            "Privacy reminder: this app is deployed on Streamlit Community Cloud. "
            "Before uploading or pasting content, anonymize or remove sensitive, confidential, personal, "
            "regulated, or client-identifiable information unless you are authorized to process it here."
        )
        if st.button("I understand - dismiss this notice", key="dismiss_privacy_notice"):
            st.session_state["privacy_notice_dismissed"] = True
            st.rerun()

# Initialize NLP globally
analyzer, lemmatizer = setup_nlp_resources()

# --- SIDEBAR (Global Inputs) ---
with st.sidebar:
    with st.expander("🧭 Start Here", expanded=True):
        st.markdown(
            """
            **Before scanning:**

            1. Choose whether to **clear previous data**.
            2. Set cleaning options and stopwords below.
            3. Add speaker exclusions if using transcripts.
            4. Open **Workspace** and upload or paste your source material.
            5. Scan, then start with the dashboard.

            For sensitive material, anonymize before upload.
            """
        )

    st.header("🧭 Scan Settings")

    # checking if there's currently data
    has_data = st.session_state['sketch'].total_rows_processed > 0
    
    # showing the checkbox
    clear_on_scan = st.checkbox(
        "Clear previous data",
        value=False,
        help="Turn this on when you want a fresh analysis. Leave it off only when you intentionally want to add new files into the current corpus.",
    )
    
    # 'banner' logic: only showing if there's ambiguity
    if has_data and not clear_on_scan:
        st.info("⚠️ **Additive Mode Active:** New scans will be ADDED to current results. Check the 'Clear previous data' check-box above to start fresh", icon="ℹ️")
    elif has_data and clear_on_scan:
        st.caption("✅ Next scan will overwrite current data.")
        
    if st.button("🗑️ Reset All"):
        reset_sketch()
        st.rerun()
        
    st.divider()
   
    st.header("🔐 AI Setup", help="This AI feature reads the *metadata only*; it does not upload raw document text to the cloud provider, nor does it 'read' it.")
    
    if st.session_state['authenticated']:
        st.success("Unlocked")
        with st.expander("🤖 Provider Settings", expanded=True):
            ai_provider = st.radio(
                "Provider",
                ["DeepSeek AI", "xAI (Grok)", "OpenAI (GPT-4o)"],
                index=0,
            )

            if ai_provider == "DeepSeek AI":
                api_key_name = "deepseek_api_key"
                base_url = "https://api.deepseek.com"
                model_name = st.selectbox(
                    "Model",
                    ["deepseek-v4-flash", "deepseek-v4-pro"],
                    index=0,
                )
                # Pricing can change; these placeholders keep the existing
                # cost-estimator plumbing working without affecting API calls.
                price_in, price_out = 0.00, 0.00

            elif "OpenAI" in ai_provider:
                api_key_name = "openai_api_key"
                base_url = None
                model_name = st.selectbox("Model", ["gpt-4o", "gpt-4o-mini"])
                price_in, price_out = (0.15, 0.60) if "mini" in model_name else (2.50, 10.00)

            else:
                api_key_name = "xai_api_key"
                base_url = "https://api.x.ai/v1"
                model_name = "grok-4-0709"
                price_in, price_out = 3.00, 15.00

            api_key = st.secrets.get(api_key_name)
            if not api_key:
                api_key = st.text_input(f"Enter {api_key_name}", type="password")

            ai_config = {
                "api_key": api_key,
                "base_url": base_url,
                "model_name": model_name,
                "price_in": price_in,
                "price_out": price_out,
            }

        with st.expander("💰 Cost Estimator", expanded=False):
            c1, c2 = st.columns(2)
            c1.markdown(f"**Tokens:**\n{st.session_state['total_tokens']:,}")
            c2.markdown(f"**Cost:**\n`${st.session_state['total_cost']:.5f}`")
            if st.button("Reset Cost"):
                st.session_state['total_cost'] = 0.0
                st.session_state['total_tokens'] = 0
                st.rerun()
        if st.button("Logout"): logout(); st.rerun()
    else:
        st.text_input("Password", type="password", key="password_input", on_change=perform_login)
        if st.session_state['auth_error']: st.error("Incorrect password")

    st.divider()
    
    # [NEW] Health Check Section
    with st.expander("❤️ System Health", expanded=False):
        st.markdown(f"{'✅' if requests else '❌'} **Web Scraper** (requests)")
        st.markdown(f"{'✅' if pypdf else '❌'} **PDF Support** (pypdf)")
        st.markdown(f"{'✅' if openpyxl else '❌'} **Excel Support** (openpyxl)")
        st.markdown(f"{'✅' if pptx else '❌'} **PowerPoint** (python-pptx)")
        st.markdown(f"{'✅' if nltk else '❌'} **Sentiment** (nltk)")
        st.markdown(f"{'✅' if LatentDirichletAllocation else '❌'} **Topic Modeling** (sklearn)")
        st.markdown(f"{'✅' if qrcode else '❌'} **QR Generator** (qrcode)")
    
    st.header("⚙️ Configuration")
    
    st.markdown("**Cleaning**")
    clean_conf = CleaningConfig(
        remove_chat=st.checkbox("Remove Chat Artifacts", True, help="Strips metadata like timestamps, usernames (e.g., <@U1234>), and system messages from logs/transcripts to focus purely on the conversation content."),
        remove_html=st.checkbox("Remove HTML", True, help="Removes HTML tags such as <div>, <br>, and other markup from exported web or LMS content."),
        remove_urls=st.checkbox("Remove URLs", True, help="Removes links and email addresses so they do not pollute the vocabulary and graph."),
        unescape=st.checkbox("Unescape HTML", True, help="Converts coded entities (e.g., &amp ; amp;, &amp ; quot;) back into readable symbols (&, \").")
    )
    
    st.markdown("**Processing**")
    use_lemma = st.checkbox("Use Lemmatization", False, help="Merges 'running' -> 'run'. Slower but cleaner.")
    if use_lemma and lemmatizer is None: st.warning("NLTK Lemmatizer not found.")
    
    proc_conf = ProcessingConfig(
        min_word_len=st.slider(
            "Min Word Len",
            1,
            10,
            4,
            help=(
                "Filters out short tokens before analysis. This affects Word Cloud, "
                "Keyphrases, Graphs, Topic Modeling, and Maturity scoring. Lower this "
                "for acronyms and short jargon such as LTI, API, SSO, or AI. Raise it "
                "only when short words are mostly noise."
            ),
        ),
        drop_integers=st.checkbox(
            "Drop Integers",
            True,
            help=(
                "Removes standalone numbers from the analysis. Turn this off if years, "
                "course codes, section numbers, ticket IDs, or numeric labels carry meaning."
            ),
        ),
        compute_bigrams=st.checkbox(
            "Bigrams",
            True,
            help=(
                "Tracks two-word phrases such as 'student success' or 'change management'. "
                "Required for NPMI, graph strength, and phrase-style maturity signals."
            ),
        ),
        use_lemmatization=use_lemma,
        translate_map=build_punct_translation(
            st.checkbox(
                "Keep Hyphens",
                help=(
                    "Useful when hyphenated terms matter, such as 'cross-functional' "
                    "or 'competency-based'. Hyphen handling can affect phrase matching."
                ),
            ),
            st.checkbox(
                "Keep Apostrophes",
                help=(
                    "Useful if apostrophes are meaningful in your text, but most corpora "
                    "are cleaner with them removed."
                ),
            ),
        )
    )

    st.markdown("**Transcript Speaker Exclusion**")
    detected_speaker_counts = Counter()
    workspace_files_for_detection = st.session_state.get("workspace_scan_upload") or []
    manual_input_for_detection = st.session_state.get("manual_text_input") or ""
    if workspace_files_for_detection:
        for uploaded_file in workspace_files_for_detection:
            try:
                detected_speaker_counts.update(
                    collect_speaker_labels_from_file(
                        uploaded_file.getvalue(),
                        uploaded_file.name,
                    )
                )
            except Exception:
                pass
    if manual_input_for_detection:
        detected_speaker_counts.update(
            collect_speaker_labels_from_text(manual_input_for_detection)
        )

    selected_detected_speakers: List[str] = []
    if detected_speaker_counts:
        speaker_options = [
            speaker
            for speaker, _count in detected_speaker_counts.most_common()
        ]
        with st.expander("Detected Speaker Labels", expanded=False):
            st.caption(
                "These are labels found before a colon in transcript-style text. "
                "Select labels to exclude their full utterances from the next scan."
            )
            selected_detected_speakers = st.multiselect(
                "Select detected speakers to exclude",
                speaker_options,
                format_func=lambda speaker: (
                    f"{speaker} ({detected_speaker_counts[speaker]:,} lines)"
                ),
            )
    else:
        st.caption(
            "No speaker labels detected yet. Upload or paste transcript-style text, "
            "or type speaker labels manually below."
        )

    excluded_speaker_input = st.text_area(
        "Exclude Speakers (comma-separated)",
        "",
        help=(
            "Optional. For transcripts with speaker labels, enter speakers whose "
            "entire utterances should be excluded before analysis. Exact matching "
            "is used by default, for example: Omar Akhtar, Speaker 18 (Blaze)."
        ),
    )
    partial_speaker_match = st.checkbox(
        "Use partial speaker matching",
        False,
        help=(
            "When enabled, an entry like 'Omar' also matches labels such as "
            "'Omar Akhtar' or 'Speaker 4 (Omar)'. Leave off for safer exact matching."
        ),
    )
    excluded_speakers = parse_speaker_exclusions(excluded_speaker_input)
    excluded_speakers.update(
        normalize_speaker_name(speaker)
        for speaker in selected_detected_speakers
        if normalize_speaker_name(speaker)
    )
    proc_conf.excluded_speakers = excluded_speakers
    proc_conf.partial_speaker_match = partial_speaker_match
    if excluded_speakers:
        st.caption(
            f"Speaker exclusion active for {len(excluded_speakers)} entr"
            f"{'y' if len(excluded_speakers) == 1 else 'ies'}. "
            "Re-scan files after changing this list."
        )
    
    # stopwords
    user_sw = st.text_area(
        "Stopwords (comma-separated)",
        "firstname.lastname, jane doe, okay, ok, really",
        help="Add domain-specific junk words, recurring names, or boilerplate terms you want removed before analysis.",
    )
    phrases, singles = parse_user_stopwords(user_sw)
    clean_conf.phrase_pattern = build_phrase_pattern(phrases)
    stopwords = set(STOPWORDS).union(singles)
    if st.checkbox(
        "Remove Generic Filler Words and Prepositions",
        True,
        help="Removes common low-value linking words such as 'of', 'to', 'with', and similar filler terms that usually add noise rather than meaning. Maturity-model vocabulary is protected separately.",
    ):
        stopwords.update(default_prepositions())

    stopwords = protect_maturity_vocabulary(stopwords)
    st.caption(
        "Safety note: maturity-model vocabulary is protected from generic stopword removal, so terms that drive maturity scoring are less likely to disappear accidentally."
    )
    proc_conf.stopwords = stopwords
    
    st.markdown("### 🎨 Appearance")
    bg_color = st.color_picker("Background Color", "#000000")
    colormap = st.selectbox("Colormap", ["viridis", "plasma", "inferno", "magma", "cividis", "tab10", "Blues", "Reds", "Greys"], 0)
    top_n = st.number_input("Top Terms to Display", min_value=5, max_value=1000, value=51, help="Controls table depth and how many top terms or bigrams are shown in ranked outputs.")
    max_words = st.slider("Max Words (Cloud)", 50, 3000, 1000, 50, help="Controls how many words are eligible for the word cloud. Higher values create denser clouds, but can make them harder to read.")
    
    # font selection
    font_map, font_names = list_system_fonts(), list(list_system_fonts().keys())
    default_font_idx = 0
    desired_font = "DejaVu Sans Mono"
    if desired_font in font_names:
        default_font_idx = font_names.index(desired_font)
    combined_font_name = st.selectbox("Font", font_names or ["(default)"], 0, help="Applies to generated word clouds and some rendered visual outputs.")
    combined_font_path = font_map.get(combined_font_name) if font_names else None

    st.markdown("### 🔬 Sentiment")
    enable_sentiment = st.checkbox("Enable Sentiment", False, help="Adds sentiment scoring to top terms, bigrams, and the Bayesian sentiment view. Best used on opinion-rich text, not neutral technical exports.")
    if enable_sentiment and analyzer is None:
        st.error("NLTK not found.")
        enable_sentiment = False
    
    pos_threshold, neg_threshold, pos_color, neu_color, neg_color = 0.05, -0.05, '#2ca02c', '#808080', '#d62728'
    if enable_sentiment:
        c1, c2 = st.columns(2)
        with c1: pos_threshold = st.slider("pos threshold", 0.0, 1.0, 0.05, 0.01)
        with c2: neg_threshold = st.slider("neg threshold", -1.0, 0.0, -0.05, 0.01)
        c1, c2, c3 = st.columns(3)
        with c1: pos_color = st.color_picker("pos color", value=pos_color)
        with c2: neu_color = st.color_picker("neu color", value=neu_color)
        with c3: neg_color = st.color_picker("neg color", value=neg_color)

    # Updated with help text
    doc_granularity = st.select_slider(
        "Rows per Doc", 
        options=[1, 5, 10, 100, 500], 
        value=5,
        help="Defines how much text gets grouped into one synthetic document for topic modeling and TF-IDF-like behavior. Use 1-5 for chats, tickets, and transcripts. Use 100+ for books, reports, and long PDFs."
    )
    st.session_state['sketch'].set_batch_size(doc_granularity)
    
    speaker_settings_hash = "|".join(sorted(proc_conf.excluded_speakers))
    current_settings_hash = (
        f"{doc_granularity}_{proc_conf.min_word_len}_"
        f"{speaker_settings_hash}_{int(proc_conf.partial_speaker_match)}"
    )
    if 'last_settings_hash' not in st.session_state: 
        st.session_state['last_settings_hash'] = current_settings_hash

    if st.session_state['last_settings_hash'] != current_settings_hash:
        if st.session_state['sketch'].total_rows_processed > 0:
            reset_sketch()
            st.warning("⚙️ Settings changed. Data reset for consistency. Please Scan again.")
        st.session_state['last_settings_hash'] = current_settings_hash

    # Updated with help text
    topic_model_type = st.selectbox(
        "Topic Model", 
        ["LDA", "NMF"],
        help="LDA is better for longer documents where each chunk may mix several themes. NMF is better for shorter, cleaner records where topics should separate more sharply."
    )
    st.caption("Topic model rule of thumb: use **NMF** for shorter rows and transcripts; use **LDA** for longer reports and mixed-content documents.")
    
    # Updated with help text
    n_topics = st.slider(
        "Topics", 
        2, 10, 
        4,
        help="How many topic buckets to generate. Start with 3-5. If topics feel too broad, increase it slightly. If they feel repetitive or noisy, decrease it."
    )

# --- TABS LAYOUT ---
# --- TABS LAYOUT ---
st.markdown("### Start Here")
st.caption(
    "Use Analyze Documents to scan files and review signals. "
    "Use Guide & Use Cases for workflow guidance, examples, and interpretation help."
)
tab_work, tab_learn = st.tabs(["🚀 Analyze Documents", "📚 Guide & Use Cases"])

# 1. THE LEARNING TAB (Guides, Examples)
with tab_learn:
    render_workflow_guide()
    render_maturity_guide()
    render_use_cases()
    render_neurotech_case_study()
    render_lit_case_study()
    render_analyst_help()

# 2. THE WORKSPACE TAB (Main Engine)
with tab_work:
    st.subheader("🚀 Scan Workspace")
    st.caption(
        "Upload files, paste text or URLs, or load an offline sketch here. "
        "Use the sidebar for cleaning, stopwords, scan mode, and analysis settings."
    )

    workspace_files = st.file_uploader(
        "Upload files to scan",
        type=["csv", "xlsx", "vtt", "txt", "json", "pdf", "pptx"],
        accept_multiple_files=True,
        key="workspace_scan_upload",
        help="Upload one or more source files. CSV/XLSX are best for structured data, VTT/TXT for transcripts, PDF/PPTX for document decks, and JSON for line-delimited text or sketches.",
    )

    with st.expander("🌐 Quick Web / Text Paste", expanded=False):
        url_input = st.text_area(
            "URLs (1 per line)",
            placeholder="https://example.com/article",
            key="workspace_url_input",
            help="Paste public URLs to scrape page text into the current scan. Best for articles, public reports, and web pages with readable body text.",
        )
        manual_input = st.text_area(
            "Manual Text Paste",
            placeholder="Paste raw text content here...",
            key="manual_text_input",
            help="Use this for notes, excerpts, transcripts, or ad hoc text that you do not want to upload as a file.",
        )

    with st.expander("📡 Load Offline Analysis (Harvester)", expanded=False):
        st.markdown("Upload pre-computed analysis from your secure server.")
        st.caption("⚠️ **Note:** Offline sketches provide Graphs, Counts & NPMI. Time-series and Entities are disabled unless included in the sketch.")

        sketch_upload = st.file_uploader(
            "Upload Sketch File (.json)",
            type=["json"],
            key="workspace_sketch_upload",
            help="Use this for very large datasets. Run the harvester.py script on your data server, then upload the resulting JSON here.",
        )

        if sketch_upload:
            file_hash = hash(sketch_upload.getvalue())
            if st.session_state.get('last_sketch_hash') != file_hash:
                try:
                    reset_sketch()
                    json_str = sketch_upload.getvalue().decode('utf-8')
                    if st.session_state['sketch'].load_from_json(json_str):
                        st.session_state['last_sketch_hash'] = file_hash
                        st.success(f"✅ Loaded Sketch: {sketch_upload.name}")
                    else:
                        st.error("❌ Invalid Schema: JSON does not match Signal Foundry structure.")
                except Exception as e:
                    st.error(f"❌ Load Error: {e}")

    with st.expander("🛠️ Data Refinery (only if you need to split very large data files; **NOTE: sanitize first**)"):
        ref_file = st.file_uploader("CSV to Refine/split", type=['csv'], key="refinery_csv_upload")
        if ref_file and st.button("🚀 Run Refinery"):
            zip_data = perform_refinery_job(ref_file, 50000, clean_conf)
            if zip_data: st.download_button("Download ZIP", zip_data, "refined.zip", "application/zip")

    # --scanning phase
    all_inputs = list(workspace_files) if workspace_files else []
    if url_input:
        for u in url_input.split('\n'):
            if u.strip(): 
                txt = fetch_url_content(u.strip())
                if txt: 
                    all_inputs.append(VirtualFile(f"url_{hash(u)}.txt", txt))
                    time.sleep(URL_SCRAPE_RATE_LIMIT_SECONDS) # RATE LIMITING
    if manual_input: all_inputs.append(VirtualFile("manual.txt", manual_input))

    if not all_inputs and st.session_state['sketch'].total_rows_processed == 0:
        st.info("Upload a file, paste text/URLs, or load an offline sketch above to begin.")

    if all_inputs:
        st.subheader("🚀 Scanning Phase")
        
        # --batch scanner
        # only shows this button if we have more than 1 file/url
        if len(all_inputs) > 1:
            # dynamic button labeling
            item_count = len(all_inputs)
            has_data = st.session_state['sketch'].total_rows_processed > 0
            
            if has_data:
                if clear_on_scan:
                    # data exists + clear is TRUE -> "re-scan / overwrite"
                    batch_btn_label = f"♻️ Re-Scan ALL {item_count} Items (Overwrite)"
                    batch_btn_help = "This will WIPE the current analysis and process the file list from scratch."
                else:
                    # if data exists + clear is FALSE -> "add" (risking duplicates)
                    batch_btn_label = f"➕ Scan & Add ALL {item_count} Items (Additive)"
                    batch_btn_help = "⚠️ CAUTION: This will process these files and ADD them to the existing results. If you have already scanned them, this will create DUPLICATES. Use 'Clear previous data' to start fresh."
            else:
                # no data -> standard start
                batch_btn_label = f"⚡ Start Batch Scan ({item_count} Items)"
                batch_btn_help = "Process all items in the list above."

            if st.button(batch_btn_label, type="primary", help=batch_btn_help):
                # handles reset logic based on user checkbox
                if clear_on_scan: 
                    reset_sketch()
                
                # 2. Setup Progress
                prog_bar = st.progress(0)
                status_box = st.empty()
                
                # 3. Iterate through all files
                for i, item in enumerate(all_inputs):
                    status_box.markdown(f"**Processing {i+1}/{len(all_inputs)}:** *{item.name}*...")
                    
                    # Detect format
                    f_bytes = item.getvalue()
                    fname = item.name.lower()
                    
                    # Logic to pick the reader (Simplified detection for batch mode)
                    batch_iter = iter([])
                    if fname.endswith(".csv"):
                        headers = detect_csv_headers(f_bytes)
                        if headers:
                            # Header exists → use first column (safe default)
                            batch_iter = read_rows_csv_structured(
                                f_bytes, "auto", ",", True, [headers[0]], None, None, " "
                            )
                        else:
                            # No header detected → fall back to treating as raw text lines (current safe behavior)
                            # Alternative (more aggressive): treat as CSV with no header and join all columns
                            # Uncomment the line below if you want to capture everything in headerless CSVs
                            # batch_iter = read_rows_csv_structured(f_bytes, "auto", ",", False,
                            #                                       [f"col_{i}" for i in range(20)], None, None, " ")
                            batch_iter = read_rows_raw_lines(f_bytes)

                    elif fname.endswith(".xlsx"):
                        sheets = get_excel_sheetnames(f_bytes)
                        if sheets:
                            batch_iter = iter_excel_structured(f_bytes, sheets[0], True, ["col_0"], None, None, " ")

                    elif fname.endswith(".pdf"):
                        batch_iter = read_rows_pdf(f_bytes)

                    elif fname.endswith(".pptx"):
                        batch_iter = read_rows_pptx(f_bytes)

                    elif fname.endswith(".vtt") or is_probably_vtt(f_bytes):
                        batch_iter = read_rows_vtt(
                            f_bytes,
                            excluded_speakers=proc_conf.excluded_speakers,
                            partial_speaker_match=proc_conf.partial_speaker_match,
                        )

                    elif fname.endswith(".json"):
                        batch_iter = read_rows_json(f_bytes)

                    else:
                        batch_iter = read_rows_raw_lines(f_bytes)
                        
                    # Process
                    process_chunk_iter(batch_iter, clean_conf, proc_conf, st.session_state['sketch'], lemmatizer)
                    
                    # Update Progress
                    prog_bar.progress((i + 1) / len(all_inputs))
                
                status_box.success(f"✅ Batch Complete! Processed {len(all_inputs)} files.")
                st.rerun()
        # -------------------------

        for idx, f in enumerate(all_inputs):
            try:
                # resource limit check
                if f.getbuffer().nbytes > MAX_FILE_SIZE_MB * 1024 * 1024:
                    st.error(f"❌ File **{f.name}** exceeds {MAX_FILE_SIZE_MB}MB limit.")
                    continue

                file_bytes, fname, lower = f.getvalue(), f.name, f.name.lower()
                is_csv = lower.endswith(".csv")
                is_xlsx = lower.endswith((".xlsx", ".xlsm"))
                is_json = lower.endswith(".json")
                is_vtt = lower.endswith(".vtt") or is_probably_vtt(file_bytes)
                is_pdf = lower.endswith(".pdf")
                is_pptx = lower.endswith(".pptx")
                
                # Default Scan Settings
                scan_settings = {
                    "date_col": None,
                    "cat_col": None,
                    "text_cols": [],
                    "has_header": False,
                    "sheet_name": None,
                    "json_key": None
                }
                
                with st.expander(f"🧩 Config: {fname}", expanded=False):
                    if is_csv:
                        headers = detect_csv_headers(file_bytes)
                        if headers:
                            scan_settings["has_header"] = True
                            st.info(f"Detected {len(headers)} columns.")
                            scan_settings["text_cols"] = st.multiselect("Text Columns", headers, default=[headers[0]], key=f"txt_{idx}", help="Choose the column or columns that contain the main text you want analyzed.")
                            scan_settings["date_col"] = st.selectbox("Date Column (Optional)", ["(None)"] + headers, key=f"date_{idx}", help="Pick a date-like column if you want the Trends tab and time slider to work.")
                            scan_settings["cat_col"] = st.selectbox("Category Column (Optional)", ["(None)"] + headers, key=f"cat_{idx}", help="Optional grouping field. Useful when the file contains segments such as department, source, or content type.")
                            
                            if scan_settings["date_col"] == "(None)": scan_settings["date_col"] = None
                            if scan_settings["cat_col"] == "(None)": scan_settings["cat_col"] = None
                        else:
                            st.warning("No headers detected. Scanning as raw text.")
                    elif is_xlsx:
                        sheets = get_excel_sheetnames(file_bytes)
                        scan_settings["sheet_name"] = st.selectbox("Sheet", sheets, key=f"sheet_{idx}", help="Choose the Excel sheet that contains the text you want to analyze.")
                        if scan_settings["sheet_name"]:
                             scan_settings["has_header"] = st.checkbox("Has Header Row", True, key=f"xls_head_{idx}", help="Leave this on if the first row contains column names rather than actual text data.")
                    elif is_json:
                        scan_settings["json_key"] = st.text_input("JSON Key (Optional)", "", key=f"json_{idx}", help="For line-delimited JSON objects, enter the field that contains the text to analyze.")

                if st.button(f"Start Scan: {fname}", key=f"btn_{idx}"):
                    if clear_on_scan: reset_sketch()
                    bar = st.progress(0)
                    status = st.empty()
                    
                    # select iterator
                    rows_iter = iter([])
                    approx = estimate_row_count_from_bytes(file_bytes)

                    if is_csv and scan_settings["has_header"] and scan_settings["text_cols"]:
                        rows_iter = read_rows_csv_structured(
                            file_bytes, "auto", ",", True, 
                            scan_settings["text_cols"], scan_settings["date_col"], scan_settings["cat_col"], " "
                        )
                    elif is_xlsx and scan_settings["sheet_name"]:
                        rows_iter = iter_excel_structured(
                            file_bytes, scan_settings["sheet_name"], scan_settings["has_header"], 
                            ["col_0"], None, None, " " 
                        )
                    elif is_pdf:
                        rows_iter = read_rows_pdf(file_bytes)
                    elif is_pptx:
                        rows_iter = read_rows_pptx(file_bytes)
                    elif is_vtt:
                        rows_iter = read_rows_vtt(
                            file_bytes,
                            excluded_speakers=proc_conf.excluded_speakers,
                            partial_speaker_match=proc_conf.partial_speaker_match,
                        )
                    elif is_json:
                        rows_iter = read_rows_json(file_bytes, scan_settings["json_key"])
                    else:
                        # fallback raw line reader
                        rows_iter = read_rows_raw_lines(file_bytes)
                    
                    # run
                    process_chunk_iter(rows_iter, clean_conf, proc_conf, st.session_state['sketch'], lemmatizer, lambda n: status.text(f"Rows: {n:,}"))
                    bar.progress(100)
                    status.success("Done!")
                    if not clear_on_scan: st.rerun()

            except Exception as e:
                st.error(f"Error: {e}")

    # --- analysis phase

    scanner = st.session_state['sketch']

    # [NEW] Time-Travel Logic
    counts_source = scanner.global_counts
    
    if scanner.temporal_counts:
        all_dates = sorted(scanner.temporal_counts.keys())
        if len(all_dates) > 1:
            st.divider()
            st.subheader("⏳ Time-Travel Filter")
            
            # Create the slider
            start_date, end_date = st.select_slider(
                "Filter Analysis by Timeframe",
                options=all_dates,
                value=(all_dates[0], all_dates[-1])
            )
            
            # if user moves the slider, reconstruct counts from the specific dates
            if (start_date != all_dates[0]) or (end_date != all_dates[-1]):
                subset_counts = Counter()
                # to iterate through dates in range
                in_range = False
                for d in all_dates:
                    if d == start_date: in_range = True
                    if in_range:
                        subset_counts.update(scanner.temporal_counts[d])
                    if d == end_date: break
                
                counts_source = subset_counts
                st.caption(f"Showing data from **{start_date}** to **{end_date}**")

    # now: dynamic filtering with visualization fix
    combined_counts = Counter({
        k: v for k, v in counts_source.items() 
        if len(str(k)) >= proc_conf.min_word_len
        and k not in proc_conf.stopwords # dynamic filtering
    })

    if combined_counts:
        st.divider()
        st.header("📊 Analysis Dashboard")
        
        # calculate stats upfront
        text_stats = calculate_text_stats(combined_counts, scanner.total_rows_processed)
        insight_df = build_insight_cards(
            scanner,
            combined_counts,
            expected_terms_raw=st.session_state.get("insight_expected_terms", ""),
            top_n=10,
        )
        render_executive_signal_dashboard(scanner, combined_counts, text_stats, insight_df)
        render_auto_insights(scanner, proc_conf)
        if st.session_state.get("authenticated"):
            render_calibration_export_panel(
                scanner=scanner,
                combined_counts=combined_counts,
                insight_df=insight_df,
                text_stats=text_stats,
                proc_conf=proc_conf,
                expected_terms_raw=st.session_state.get("insight_expected_terms", ""),
                key_prefix="calibration_export_quick_access",
            )
        # main tabs
        analysis_tab_labels = [
            "💡 Insight Engine",
            "☁️ Word Cloud & Stats",
            "🧭 Themes",
            "📈 Trends",
            "👥 Entities",
            "🔑 Keyphrases",
            "🏆 Maturity",
        ]
        if st.session_state.get("authenticated"):
            analysis_tab_labels.append("🧪 Calibration Export")

        analysis_tabs = st.tabs(analysis_tab_labels)
        tab_insight, tab_main, tab_theme, tab_trend, tab_ent, tab_key, tab_mat = analysis_tabs[:7]
        tab_calibration = analysis_tabs[7] if st.session_state.get("authenticated") else None
        
        with tab_insight:
            st.subheader("💡 Insight Engine")
            st.caption(
                "This layer connects statistical signals to representative excerpts, likely signal types, "
                "interpretive notes, and follow-up questions. Treat it as an analyst's first-pass map, not a final verdict."
            )

            c_insight_1, c_insight_2, c_insight_3 = st.columns(3)
            c_insight_1.metric(
                "Evidence Snippets",
                f"{len(scanner.evidence_docs):,}",
                help=ANALYSIS_HELP_TEXT["Evidence Snippets"],
            )
            c_insight_2.metric(
                "Signal Types",
                f"{len(SIGNAL_TAXONOMY):,}",
                help=ANALYSIS_HELP_TEXT["Signal Types"],
            )
            c_insight_3.metric(
                "Evidence Cap",
                f"{MAX_EVIDENCE_DOCS:,}",
                help=ANALYSIS_HELP_TEXT["Evidence Cap"],
            )

            if scanner.evidence_limit_reached:
                st.warning(
                    "Evidence capture hit its safety cap. Counts and charts still use the full scan, "
                    "but insight evidence is based on the retained sample."
                )

            insight_expected_terms = st.text_area(
                "Hypothesis / Concept Check",
                key="insight_expected_terms",
                placeholder="Optional: governance, workload, trust, escalation, training",
                help=(
                    "Optional. Use this if you already have concepts, risks, themes, or required topics "
                    "you want to check against the text. Leave blank for open-ended discovery."
                ),
            )

            insight_df = build_insight_cards(
                scanner,
                combined_counts,
                expected_terms_raw=insight_expected_terms,
                top_n=10,
            )

            if insight_df.empty:
                st.info(
                    "Not enough evidence yet to build insight cards. Scan more text or reduce filtering."
                )
            else:
                render_signal_compass_panel(insight_df)
                st.divider()
                render_resource_shape_panel(insight_df)
                st.divider()
                st.subheader("🔎 Supporting Insight Cards")
                st.caption(
                    "Use these cards to inspect the evidence behind the Resource Shape. "
                    "They are supporting leads, not the final answer."
                )
                role_filter = st.multiselect(
                    "Filter by signal role",
                    sorted(insight_df["Signal Role"].unique()) if "Signal Role" in insight_df.columns else [],
                    default=[],
                    help="Core Insight items are strongest. Supporting, motif, context, and low-specificity items help explain what was promoted or demoted.",
                )
                signal_type_filter = st.multiselect(
                    "Filter by signal type",
                    sorted(insight_df["Signal Type"].unique()),
                    default=[],
                    help="Leave empty to show all insight cards.",
                )
                visible_insights = insight_df
                if role_filter and "Signal Role" in visible_insights.columns:
                    visible_insights = visible_insights[
                        visible_insights["Signal Role"].isin(role_filter)
                    ]
                if signal_type_filter:
                    visible_insights = visible_insights[
                        visible_insights["Signal Type"].isin(signal_type_filter)
                    ]

                for idx, row in visible_insights.iterrows():
                    title = (
                        f"{row.get('Signal Role', 'Signal')} · {row['Confidence']} confidence · {row['Signal Type']} · "
                        f"{row['Signal']}"
                    )
                    with st.expander(title, expanded=(idx < 3)):
                        m1, m2, m3, m4 = st.columns(4)
                        m1.metric(
                            "Evidence Strength",
                            row["Evidence Strength"],
                            help=ANALYSIS_HELP_TEXT["Evidence Strength"],
                        )
                        m2.metric(
                            "Distinctiveness",
                            row["Distinctiveness"],
                            help=ANALYSIS_HELP_TEXT["Distinctiveness"],
                        )
                        m3.metric(
                            "Confidence",
                            row["Confidence"],
                            help=ANALYSIS_HELP_TEXT["Confidence"],
                        )
                        m4.metric(
                            "Interpretive Lift",
                            f"{row.get('Interpretive Lift', 0)}/100",
                            help=ANALYSIS_HELP_TEXT["Interpretive Lift"],
                        )
                        if row.get("Ranking Rationale"):
                            st.caption(row["Ranking Rationale"])
                        st.markdown("**Interpretation**")
                        st.write(row["Interpretation"])
                        st.markdown("**Representative Evidence**")
                        st.write(row["Representative Evidence"])
                        st.markdown("**Follow-up Question**")
                        st.info(row["Follow-up Question"])

                st.download_button(
                    "📥 Download insight cards CSV",
                    dataframe_to_csv_bytes(insight_df),
                    "insight_cards.csv",
                    "text/csv",
                )

        if tab_calibration is not None:
            with tab_calibration:
                st.subheader("🧪 Calibration Export")
                st.caption(
                    "Admin-only export for comparing repeated test runs. "
                    "Use this to review calibration quality without relying on screenshots."
                )
                if insight_df.empty:
                    st.info(
                        "Not enough insight-card data yet. Scan a resource first, then return here."
                    )
                else:
                    render_calibration_export_panel(
                        scanner=scanner,
                        combined_counts=combined_counts,
                        insight_df=insight_df,
                        text_stats=text_stats,
                        proc_conf=proc_conf,
                        expected_terms_raw=insight_expected_terms,
                        key_prefix="calibration_export_tab",
                    )

        with tab_main:
            if enable_sentiment:
                top_keys = [k for k,v in combined_counts.most_common(1000)]
                term_sentiments = get_sentiments(analyzer, tuple(top_keys))
                if proc_conf.compute_bigrams:
                     top_bg_keys = [" ".join(k) for k,v in scanner.global_bigrams.most_common(2000)]
                     term_sentiments.update(get_sentiments(analyzer, tuple(top_bg_keys)))
                c_color_func = create_sentiment_color_func(term_sentiments, pos_color, neg_color, neu_color, pos_threshold, neg_threshold)
                fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, 800, 400, bg_color, colormap, combined_font_path, 42, c_color_func)
            else:
                term_sentiments = {}
                fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, 800, 400, bg_color, colormap, combined_font_path, 42, None)
                
            st.pyplot(fig, use_container_width=True)
            st.download_button("📥 download combined png", fig_to_png_bytes(fig), "combined_wc.png", "image/png")
            
            c1, c2, c3, c4 = st.columns(4)
            
            c1.metric(
                "Total Tokens", 
                f"{text_stats['Total Tokens']:,}",
                help=ANALYSIS_HELP_TEXT["Total Tokens"],
            )
            
            c2.metric(
                "Unique Vocab", 
                f"{text_stats['Unique Vocabulary']:,}",
                help=ANALYSIS_HELP_TEXT["Unique Vocab"],
            )
            
            c3.metric(
                "Docs/Rows", 
                f"{text_stats['Total Rows']:,}",
                help=ANALYSIS_HELP_TEXT["Docs/Rows"],
            )
            
            c4.metric(
                "Lexical Diversity", 
                f"{text_stats['Lexical Diversity']}",
                help=ANALYSIS_HELP_TEXT["Lexical Diversity"],
            )

        with tab_theme:
            st.subheader("🧭 Interpretive Lens")
            st.caption(
                "Use these views to read between the lines: recurring phrases, distinctive terms, missing expected signals, "
                "category differences, and time-based shifts. These are evidence prompts, not automated conclusions."
            )

            st.markdown("#### Theme Evidence Cards")
            theme_df = build_theme_evidence_cards(scanner, combined_counts, top_n=8)
            if not theme_df.empty:
                st.dataframe(
                    theme_df,
                    use_container_width=True,
                    column_config={
                        "Theme Evidence": st.column_config.TextColumn("Theme Evidence", help="A phrase or term that may point to an underlying theme."),
                        "Evidence Type": st.column_config.TextColumn("Evidence Type", help="Whether the signal comes from phrase association or corpus distinctiveness."),
                        "Support": st.column_config.NumberColumn("Support", help="How often this signal appears."),
                        "Distinctiveness": st.column_config.NumberColumn("Distinctiveness", help=ANALYSIS_HELP_TEXT["Distinctiveness"]),
                        "Related Terms": st.column_config.TextColumn("Related Terms", help="Nearby terms that often travel with this signal."),
                        "Read As": st.column_config.TextColumn("Read As", help="Plain-language interpretation guidance."),
                    }
                )
                st.download_button(
                    "📥 Download theme evidence CSV",
                    dataframe_to_csv_bytes(theme_df),
                    "theme_evidence.csv",
                    "text/csv",
                )
            else:
                st.info("Not enough phrase or TF-IDF signal yet to build theme evidence cards.")

            st.markdown("#### Signal Quadrant: Frequency vs. Distinctiveness")
            st.caption(
                "High-frequency terms show the backdrop; high-distinctiveness terms reveal stronger signal. "
                "Hover over points for term names. The table below keeps the full ranked detail without cluttering the chart."
            )
            quadrant_df = build_signal_quadrant_df(scanner, combined_counts, top_n=150)
            if not quadrant_df.empty:
                q_order = ["Core signal", "Niche signal", "Common backdrop", "Low evidence"]
                q_colors = {
                    "Core signal": "#2ca02c",
                    "Common backdrop": "#1f77b4",
                    "Niche signal": "#ff7f0e",
                    "Low evidence": "#7f7f7f",
                }
                chart_df = quadrant_df.head(120).copy()
                chart_df["Evidence Score"] = chart_df["Frequency"] * chart_df["Distinctiveness"]

                x_split = float(chart_df["Frequency"].median())
                y_split = float(chart_df["Distinctiveness"].median())

                if alt is not None:
                    base = alt.Chart(chart_df).encode(
                        x=alt.X("Frequency:Q", title="Frequency"),
                        y=alt.Y("Distinctiveness:Q", title="Distinctiveness"),
                    )
                    points = base.mark_circle(size=78, opacity=0.72).encode(
                        color=alt.Color(
                            "Quadrant:N",
                            sort=q_order,
                            scale=alt.Scale(
                                domain=q_order,
                                range=[q_colors.get(q, "#7f7f7f") for q in q_order],
                            ),
                            legend=alt.Legend(title="Quadrant", orient="bottom-right"),
                        ),
                        tooltip=[
                            alt.Tooltip("Term:N", title="Term"),
                            alt.Tooltip("Frequency:Q", title="Frequency"),
                            alt.Tooltip("Distinctiveness:Q", title="Distinctiveness", format=".2f"),
                            alt.Tooltip("Evidence Score:Q", title="Evidence Score", format=".1f"),
                            alt.Tooltip("Quadrant:N", title="Quadrant"),
                        ],
                    )
                    v_rule = alt.Chart(pd.DataFrame({"x": [x_split]})).mark_rule(strokeDash=[4, 4], opacity=0.45).encode(x="x:Q")
                    h_rule = alt.Chart(pd.DataFrame({"y": [y_split]})).mark_rule(strokeDash=[4, 4], opacity=0.45).encode(y="y:Q")
                    chart = (points + v_rule + h_rule).properties(
                        height=430,
                        title="Signal Quadrant",
                    ).interactive()
                    st.altair_chart(chart, use_container_width=True)
                else:
                    # Fallback for minimal environments where Altair is unavailable.
                    fig_q, ax_q = plt.subplots(figsize=(8, 4.5))
                    for q_name, q_group in chart_df.groupby("Quadrant"):
                        ax_q.scatter(
                            q_group["Frequency"],
                            q_group["Distinctiveness"],
                            label=q_name,
                            alpha=0.72,
                            s=42,
                            color=q_colors.get(q_name, "#7f7f7f"),
                        )
                    ax_q.axvline(x_split, linestyle="--", alpha=0.35)
                    ax_q.axhline(y_split, linestyle="--", alpha=0.35)
                    ax_q.set_xlabel("Frequency")
                    ax_q.set_ylabel("Distinctiveness")
                    ax_q.set_title("Signal Quadrant")
                    ax_q.legend(loc="best", fontsize=8)
                    ax_q.grid(alpha=0.2)
                    st.pyplot(fig_q, use_container_width=True)
                    plt.close(fig_q)

                top_signal_df = chart_df.sort_values("Evidence Score", ascending=False).head(15)
                with st.expander("Top high-signal points shown in the quadrant", expanded=False):
                    st.dataframe(
                        top_signal_df[["Term", "Frequency", "Distinctiveness", "Evidence Score", "Quadrant"]],
                        use_container_width=True,
                        column_config={
                            "Term": st.column_config.TextColumn("Term", help="A word or phrase extracted from the processed text."),
                            "Frequency": st.column_config.NumberColumn("Frequency", help=ANALYSIS_HELP_TEXT["Frequency"]),
                            "Distinctiveness": st.column_config.NumberColumn("Distinctiveness", format="%.2f", help=ANALYSIS_HELP_TEXT["Distinctiveness"]),
                            "Evidence Score": st.column_config.NumberColumn("Evidence Score", format="%.1f", help=ANALYSIS_HELP_TEXT["Evidence Score"]),
                            "Quadrant": st.column_config.TextColumn("Quadrant", help=ANALYSIS_HELP_TEXT["Quadrant"]),
                        },
                    )

                st.dataframe(
                    quadrant_df.head(80),
                    use_container_width=True,
                    column_config={
                        "Term": st.column_config.TextColumn("Term", help="A word or phrase extracted from the processed text."),
                        "Frequency": st.column_config.NumberColumn("Frequency", help=ANALYSIS_HELP_TEXT["Frequency"]),
                        "Distinctiveness": st.column_config.NumberColumn("Distinctiveness", help=ANALYSIS_HELP_TEXT["Distinctiveness"]),
                        "Evidence Score": st.column_config.NumberColumn("Evidence Score", help=ANALYSIS_HELP_TEXT["Evidence Score"]),
                        "Quadrant": st.column_config.TextColumn("Quadrant", help=ANALYSIS_HELP_TEXT["Quadrant"]),
                    },
                )
                st.download_button(
                    "📥 Download signal quadrant CSV",
                    dataframe_to_csv_bytes(quadrant_df),
                    "signal_quadrant.csv",
                    "text/csv",
                )
            else:
                st.info("Not enough term data yet to build the signal quadrant.")

            st.markdown("#### Hypothesis / Concept Check")
            expected_raw = st.text_area(
                "Optional concepts, risks, themes, or required topics to check",
                key="expected_signal_terms",
                placeholder="Optional: procurement delay, escalation, training, governance",
                help=(
                    "Use this when you already have concepts, risks, themes, or required topics "
                    "you want to check against the text. Leave blank for open-ended discovery."
                ),
            )
            expected_df = build_expected_terms_df(expected_raw, combined_counts, scanner.global_bigrams)
            if not expected_df.empty:
                st.dataframe(expected_df, use_container_width=True)
                st.download_button(
                    "📥 Download expected signal CSV",
                    dataframe_to_csv_bytes(expected_df),
                    "expected_signal_check.csv",
                    "text/csv",
                )
            else:
                st.caption("Add expected terms to check whether important themes are present, weak, or absent in the scanned text.")

            st.markdown("#### Contrastive Analysis")
            if scanner.category_counts and len(scanner.category_counts) >= 2:
                categories = sorted(scanner.category_counts.keys())
                col_left, col_right = st.columns(2)
                left_category = col_left.selectbox("Compare category A", categories, index=0, key="contrast_left_category")
                right_default = 1 if len(categories) > 1 else 0
                right_category = col_right.selectbox("Compare category B", categories, index=right_default, key="contrast_right_category")

                if left_category != right_category:
                    contrast_df = compare_counter_terms(
                        scanner.category_counts[left_category],
                        scanner.category_counts[right_category],
                        proc_conf.stopwords,
                        proc_conf.min_word_len,
                        top_n=50,
                    )
                    if not contrast_df.empty:
                        contrast_df = contrast_df.rename(columns={
                            "Left Count": f"{left_category} Count",
                            "Right Count": f"{right_category} Count",
                            "Left Rate": f"{left_category} Rate",
                            "Right Rate": f"{right_category} Rate",
                        })
                        contrast_df["Leans Toward"] = contrast_df["Leans Toward"].replace({
                            "Left": left_category,
                            "Right": right_category,
                        })
                        st.dataframe(contrast_df, use_container_width=True)
                        st.download_button(
                            "📥 Download contrastive analysis CSV",
                            dataframe_to_csv_bytes(contrast_df),
                            "contrastive_analysis.csv",
                            "text/csv",
                        )
                    else:
                        st.info("No strong category differences were detected with the current settings.")
                else:
                    st.caption("Choose two different categories to compare.")
            else:
                st.info("Contrastive analysis appears when scanned data includes a category column.")

            st.markdown("#### Temporal Drift")
            if scanner.temporal_counts and len(scanner.temporal_counts) >= 2:
                drift_df = build_temporal_drift_df(
                    scanner.temporal_counts,
                    proc_conf.stopwords,
                    proc_conf.min_word_len,
                    top_n=50,
                )
                if not drift_df.empty:
                    st.dataframe(drift_df, use_container_width=True)
                    st.download_button(
                        "📥 Download temporal drift CSV",
                        dataframe_to_csv_bytes(drift_df),
                        "temporal_drift.csv",
                        "text/csv",
                    )
                else:
                    st.info("No strong early-vs-late shifts were detected with the current settings.")
            else:
                st.info("Temporal drift appears when scanned data includes at least two valid dates.")

        with tab_trend:
            if scanner.temporal_counts:
                st.markdown("#### Word Volume Over Time")
                trend_data = []
                for d_str, counts in scanner.temporal_counts.items():
                    trend_data.append({"Date": d_str, "Volume": sum(counts.values())})
                
                df_trend = pd.DataFrame(trend_data).sort_values("Date")
                st.line_chart(df_trend.set_index("Date"))
                
                st.markdown("#### Specific Term Trends")
                terms_to_plot = st.multiselect("Select terms to plot", [t for t, c in combined_counts.most_common(50)])
                if terms_to_plot:
                    term_trend_data = []
                    for d_str, counts in scanner.temporal_counts.items():
                        row = {"Date": d_str}
                        for t in terms_to_plot: row[t] = counts[t]
                        term_trend_data.append(row)
                    df_term_trend = pd.DataFrame(term_trend_data).sort_values("Date").set_index("Date")
                    st.line_chart(df_term_trend)
            else:
                st.info("No Date column was selected during scan (or no valid dates found).")

        with tab_ent:
            st.markdown("#### Top Entities (Polymorphic NER)")
            st.caption(
                "**What this finds:** Stakeholders (Who) and Systems (What). "
                "This engine adapts to capture **Standard Names** (e.g., 'John Doe'), "
                "**Acronyms** (e.g., 'DARPA'), and **Technical IDs** (e.g., 'COVID-19' or 'F-35')."
            )
            
            # smarter ratio filter
            # if "apple" (fruit) appears 100 times, and "Apple" (entity) appears 2 times, drop it
            # if "Apple" (Company) appears 100 times, keeps it
            refined_entities = Counter()
            for ent, count in scanner.entity_counts.items():
                lower_k = ent.lower()
                total_occurrences = scanner.global_counts.get(lower_k, 0)
                
                # if word appears in the text, check the ratio
                if total_occurrences > 0:
                    capitalization_ratio = count / total_occurrences
                    # keeps if it's capitalized >30% of the time OR it's a complex ID (digits/hyphens)
                    if capitalization_ratio > 0.3 or not ent.isalpha():
                        refined_entities[ent] = count
                else:
                    # if not in global_counts (e.g. stopped out), keeps it to be safe
                    refined_entities[ent] = count

            if refined_entities:
                ent_df = pd.DataFrame(refined_entities.most_common(50), columns=["Entity", "Count"])
                st.dataframe(ent_df, use_container_width=True)
                
                # simple entity cloud (safety wrapped)
                try:
                    fig_e, _ = build_wordcloud_figure_from_counts(refined_entities, 100, 800, 400, "#111111", "Pastel1", combined_font_path, 42, None)
                    st.pyplot(fig_e)
                except Exception as e:
                    st.warning(f"Could not generate Entity Cloud: {e}")
            else:
                st.info("No capitalized entities detected.")

        with tab_key:
            st.subheader("🔑 TF-IDF Keyphrases (The 'Technical DNA')")
            
            # 1 plain language explanation banner
            st.info(
                "**How to read this:** Unlike simple word counts, **TF-IDF** penalizes words that appear everywhere (like 'report' or 'email') and boosts words that are unique to specific incidents or documents. \n\n"
                "👉 **High Score** = Rare, specific, high-signal (e.g. 'Oximetry').\n"
                "👉 **Low Score** = Common, generic, low-signal.",
                icon="ℹ️"
            )
            
            df_tfidf = calculate_tfidf(scanner, 50)
            
            # 2 DataFrame with hover-over tooltips
            st.dataframe(
                df_tfidf, 
                use_container_width=True,
                column_config={
                    "Term": st.column_config.TextColumn("Term", help="The extracted vocabulary word."),
                    "TF (Count)": st.column_config.NumberColumn("TF (Count)", help="Term Frequency: Total number of times this word appears."),
                    "DF (Docs)": st.column_config.NumberColumn("DF (Docs)", help="Document Frequency: Number of distinct documents (or chunks) containing this word. Low DF = Specific."),
                    "Keyphrase Score": st.column_config.NumberColumn("Keyphrase Score", help="Mathematical Uniqueness. Higher = More 'Technical' and less 'Generic'.")
                }
            )

        with tab_mat:
            st.subheader("🏆 Maturity Assessment Engine")
            # 1. Initialize
            assessor = MaturityAssessor()
            model_options = assessor.get_model_names()
            # 2. Persona Selector
            c_sel1, c_sel2 = st.columns([1, 3])
            with c_sel1:
                selected_model = st.selectbox(
                    "Select Organization Persona:", model_options, index=0
                )
            with c_sel2:
                st.info(f"ℹ️ {assessor.get_model_desc(selected_model)}")
                st.caption(
                    "Maturity scoring uses the current cleaned token stream. "
                    "Settings such as Min Word Len, stopwords, lemmatization, "
                    "hyphen handling, and chat/HTML cleanup can change the result."
                )
            # 3. Run Assessment
            maturity_result = assessor.assess(
                combined_counts, scanner.global_bigrams, selected_model
            )
            if maturity_result:
                result_type = maturity_result.get("type", "flat")
                if result_type == "flat":
                    # === ORIGINAL FLAT RENDERING (unchanged) ===
                    st.divider()
                    m_col1, m_col2 = st.columns([1, 2])
                    with m_col1:
                        score = maturity_result['overall_score']
                        dom_stage = maturity_result['dominant_stage']
                        gauge_color = dom_stage['color']
                        st.metric(
                            "Maturity Score (1.0 - 5.0)",
                            f"{score} / 5.0",
                            help=ANALYSIS_HELP_TEXT["Maturity Score"],
                        )
                        st.markdown(
                            f"#### Phase: <span style='color:{gauge_color}'>"
                            f"{dom_stage['name']}</span>",
                            unsafe_allow_html=True
                        )
                        st.metric(
                            "Signal Density",
                            f"{maturity_result['total_signals_found']} words",
                            help=ANALYSIS_HELP_TEXT["Signal Density"],
                        )
                    with m_col2:
                        fig_mat = assessor.render_radar_chart(maturity_result)
                        if fig_mat:
                            st.pyplot(fig_mat, use_container_width=True)
                        st.caption(
                            "**Analyst Tip:** A balanced polygon (round) indicates "
                            "alignment between Strategy and Execution. A sharp "
                            "'spike' indicates a capability gap in adjacent areas."
                        )
                    st.subheader("🔍 Linguistic Drivers")
                    cols = st.columns(5)
                    levels_ref = maturity_result['levels_ref']
                    for i in range(1, 6):
                        with cols[i - 1]:
                            lvl_data = levels_ref[i]
                            st.markdown(f"**L{i}: {lvl_data['name']}**")
                            terms = lvl_data['terms']
                            found = {
                                t: combined_counts[t]
                                for t in terms if t in combined_counts
                            }
                            if found:
                                top_words = sorted(
                                    found.items(), key=lambda x: x[1], reverse=True
                                )[:4]
                                for w, c in top_words:
                                    st.caption(f"{w} ({c})")
                            else:
                                st.caption("(No signals)")
                elif result_type == "domain_based":
                    # === NEW DOMAIN-BASED RENDERING ===
                    st.divider()
                    # --- Top-level metrics ---
                    m_col1, m_col2, m_col3 = st.columns(3)
                    with m_col1:
                        st.metric(
                            "Composite Maturity (1.0 - 3.0)",
                            f"{maturity_result['overall_score']} / 3.0",
                            help=ANALYSIS_HELP_TEXT["Composite Maturity"],
                        )
                    with m_col2:
                        st.metric(
                            "Domains Assessed",
                            f"{maturity_result['domains_assessed']} / "
                            f"{maturity_result['domains_total']}",
                            help=ANALYSIS_HELP_TEXT["Domains Assessed"],
                        )
                    with m_col3:
                        st.metric(
                            "Total Signals",
                            f"{maturity_result['total_signals_found']}",
                            help=ANALYSIS_HELP_TEXT["Total Signals"],
                        )
                    # --- Radar Chart ---
                    # --- Radar Chart ---
                    st.markdown("#### 🕸️ Domain Maturity Radar")

                    with st.expander("🎨 Radar font controls", expanded=False):
                        font_col1, font_col2, font_col3 = st.columns(3)

                        with font_col1:
                            radar_font_family = st.selectbox(
                                "Font family",
                                ["sans-serif", "serif", "monospace", "DejaVu Sans", "Arial"],
                                index=0,
                                key="radar_font_family",
                            )
                            radar_label_size = st.slider(
                                "Domain label size",
                                min_value=6,
                                max_value=16,
                                value=9,
                                step=1,
                                key="radar_label_size",
                            )

                        with font_col2:
                            radar_label_color = st.color_picker(
                                "Label color",
                                value="#222222",
                                key="radar_label_color",
                            )
                            radar_wrap_width = st.slider(
                                "Domain label wrap width",
                                min_value=6,
                                max_value=24,
                                value=14,
                                step=1,
                                key="radar_wrap_width",
                            )

                        with font_col3:
                            radar_show_tier_labels = st.checkbox(
                                "Show tier names",
                                value=True,
                                key="radar_show_tier_labels",
                            )
                            radar_tier_label_angle = st.slider(
                                "Tier label angle",
                                min_value=0,
                                max_value=359,
                                value=35,
                                step=5,
                                key="radar_tier_label_angle",
                            )

                    fig_radar = assessor.render_domain_radar_chart(
                        maturity_result,
                        font_family=radar_font_family,
                        label_size=radar_label_size,
                        label_color=radar_label_color,
                        wrap_width=radar_wrap_width,
                        show_tier_labels=radar_show_tier_labels,
                        tier_label_angle=radar_tier_label_angle,
                    )

                    if fig_radar:
                        st.pyplot(fig_radar, use_container_width=True)
                    st.caption(
                        "🔴 Foundational (<1.5) | 🟠 Advanced (1.5–2.5) | "
                        "🟢 Leading Edge (>2.5)"
                    )

                    render_tam_domain_glossary(maturity_result["domain_results"])

                    # --- Stacked Bar Breakdown ---
                    st.markdown("#### 📊 Tier Distribution by Domain")
                    fig_bar = assessor.render_domain_breakdown_chart(maturity_result)
                    if fig_bar:
                        st.pyplot(fig_bar, use_container_width=True)
                    # --- Per-Domain Detail Cards ---
                    st.markdown("#### 🔍 Domain Details")
                    domain_results = maturity_result["domain_results"]
                    for domain_key in sorted(domain_results.keys()):
                        dr = domain_results[domain_key]

                        if dr["signals"] == 0:
                            indicator = "⚪"
                        elif dr["score"] >= 2.5:
                            indicator = "🟢"
                        elif dr["score"] >= 1.5:
                            indicator = "🟠"
                        else:
                            indicator = "🔴"

                        with st.expander(
                            f"{indicator} {dr['name']} — "
                            f"{dr['tier_label']} ({dr['score']}/3.0, "
                            f"{dr['signals']} signals)"
                        ):
                            domain_help = get_tam_domain_help(domain_key)
                            st.markdown(f"**What this domain means:** {domain_help['plain']}")
                            st.caption(f"**What it covers:** {domain_help['covers']}")
                            st.caption(f"**Leading-edge signals:** {domain_help['leading']}")

                            if dr["signals"] == 0:
                                st.info(
                                    "No matching vocabulary from this dataset was mapped to this domain yet. "
                                    "This usually means either the source material did not discuss this area, "
                                    "or the current domain vocabulary is too narrow for the way it was discussed."
                                )
                                continue

                            dist = dr["distribution"]
                            dc1, dc2, dc3 = st.columns(3)
                            dc1.metric(
                                "Foundational",
                                f"{dist.get(1, 0):.0%}",
                                help=ANALYSIS_HELP_TEXT["Foundational"],
                            )
                            dc2.metric(
                                "Advanced",
                                f"{dist.get(2, 0):.0%}",
                                help=ANALYSIS_HELP_TEXT["Advanced"],
                            )
                            dc3.metric(
                                "Leading Edge",
                                f"{dist.get(3, 0):.0%}",
                                help=ANALYSIS_HELP_TEXT["Leading Edge"],
                            )
                            st.markdown("**Top Linguistic Drivers:**")
                            tier_names = {
                                1: "Foundational", 2: "Advanced", 3: "Leading Edge"
                            }
                            for tier_num in [1, 2, 3]:
                                drivers = dr["drivers"].get(tier_num, [])
                                if drivers:
                                    driver_str = ", ".join(
                                        [f"`{t}` ({c})" for t, c in drivers]
                                    )
                                    st.caption(
                                        f"**{tier_names[tier_num]}:** {driver_str}"
                                    )
                    # ============================================
                    # PHASE 10: DOWNLOAD ASSESSMENT (inline here)
                    # ============================================
                    st.divider()
                    st.markdown("#### 💾 Export Assessment")
                    client_name_input = st.text_input(
                        "Client Name (for export filename):",
                        value="Client",
                        key="mat_client_name"
                    )
                    # Sanitize client name for filename
                    safe_client = re.sub(r'[^\w\s-]', '', client_name_input).strip()
                    safe_client = re.sub(r'[\s]+', '_', safe_client)
                    export_date = datetime.now().strftime("%Y-%m-%d")
                    export_filename = f"{safe_client}_Maturity_Assessment_{export_date}.json"
                    # Build export payload
                    export_payload = {
                        "schema_version": "1.0",
                        "client_name": client_name_input.strip(),
                        "assessment_date": export_date,
                        "model_name": selected_model,
                        "composite_score": maturity_result["overall_score"],
                        "max_score": maturity_result["max_score"],
                        "domains_assessed": maturity_result["domains_assessed"],
                        "domains_total": maturity_result["domains_total"],
                        "total_signals": maturity_result["total_signals_found"],
                        "domain_results": {}
                    }
                    for dk, dr in maturity_result["domain_results"].items():
                        export_payload["domain_results"][dk] = {
                            "name": dr["name"],
                            "short": dr["short"],
                            "score": dr["score"],
                            "tier_label": dr["tier_label"],
                            "distribution": {
                                "foundational": round(dr["distribution"].get(1, 0), 4),
                                "advanced": round(dr["distribution"].get(2, 0), 4),
                                "leading_edge": round(dr["distribution"].get(3, 0), 4)
                            },
                            "signals": dr["signals"],
                            "top_drivers": {
                                "foundational": [
                                    f"{t} ({c})" for t, c in dr["drivers"].get(1, [])
                                ],
                                "advanced": [
                                    f"{t} ({c})" for t, c in dr["drivers"].get(2, [])
                                ],
                                "leading_edge": [
                                    f"{t} ({c})" for t, c in dr["drivers"].get(3, [])
                                ]
                            }
                        }
                    export_json_str = json.dumps(export_payload, indent=2)
                    st.download_button(
                        label=f"⬇️ Download: {export_filename}",
                        data=export_json_str,
                        file_name=export_filename,
                        mime="application/json"
                    )
                    # =====================================================
                    # PHASE 11: UPLOAD SNAPSHOTS + LONGITUDINAL TRACKING
                    # =====================================================
                    st.divider()
                    st.markdown("#### 📈 Longitudinal Progress Tracking")
                    st.caption(
                        "Upload previous assessment JSON files to track maturity "
                        "progress over time. The current live assessment (if running) "
                        "will be included automatically."
                    )
                    uploaded_snapshots = st.file_uploader(
                        "Upload previous assessment snapshots (.json):",
                        type=["json"],
                        accept_multiple_files=True,
                        key="mat_snapshot_upload"
                    )
                    # Parse uploaded snapshots
                    all_snapshots = []
                    parse_errors = []
                    if uploaded_snapshots:
                        for uf in uploaded_snapshots:
                            try:
                                content = json.loads(uf.read())
                                uf.seek(0)  # reset for potential re-read
                                # Validate schema
                                if "schema_version" not in content:
                                    parse_errors.append(
                                        f"⚠️ {uf.name}: Missing schema_version. Skipping."
                                    )
                                    continue
                                if "assessment_date" not in content:
                                    parse_errors.append(
                                        f"⚠️ {uf.name}: Missing assessment_date. Skipping."
                                    )
                                    continue
                                all_snapshots.append(content)
                            except json.JSONDecodeError:
                                parse_errors.append(
                                    f"⚠️ {uf.name}: Invalid JSON. Skipping."
                                )
                    # Append current live assessment as latest snapshot
                    if maturity_result and maturity_result.get("type") == "domain_based":
                        live_snapshot = {
                            "schema_version": "1.0",
                            "client_name": client_name_input.strip(),
                            "assessment_date": export_date,
                            "composite_score": maturity_result["overall_score"],
                            "domain_results": {}
                        }
                        for dk, dr in maturity_result["domain_results"].items():
                            live_snapshot["domain_results"][dk] = {
                                "name": dr["name"],
                                "short": dr["short"],
                                "score": dr["score"],
                                "tier_label": dr["tier_label"],
                                "signals": dr["signals"]
                            }
                        all_snapshots.append(live_snapshot)
                    # Sort by date
                    all_snapshots.sort(
                        key=lambda s: s.get("assessment_date", "0000-00-00")
                    )
                    # Show parse errors
                    for err in parse_errors:
                        st.warning(err)
                    if len(all_snapshots) >= 2:
                        st.success(
                            f"📊 {len(all_snapshots)} snapshots loaded "
                            f"(including current assessment)."
                        )
                        # --- Composite Trend ---
                        st.markdown("##### Composite Score Trend")
                        fig_trend = assessor.render_trend_chart(all_snapshots)
                        if fig_trend:
                            st.pyplot(fig_trend, use_container_width=True)
                        # --- Per-Domain Trend ---
                        st.markdown("##### Per-Domain Score Trend")
                        fig_domain_trend = assessor.render_domain_trend_chart(
                            all_snapshots
                        )
                        if fig_domain_trend:
                            st.pyplot(fig_domain_trend, use_container_width=True)
                        # --- Domain Movement Table ---
                        st.markdown("##### Domain Movement (First → Latest)")
                        first = all_snapshots[0]
                        latest = all_snapshots[-1]
                        movement_data = []
                        all_domain_keys = set(
                            list(first.get("domain_results", {}).keys()) +
                            list(latest.get("domain_results", {}).keys())
                        )
                        for dk in sorted(all_domain_keys):
                            first_dr = first.get("domain_results", {}).get(dk, {})
                            latest_dr = latest.get("domain_results", {}).get(dk, {})
                            first_score = first_dr.get("score", 0.0)
                            latest_score = latest_dr.get("score", 0.0)
                            delta = round(latest_score - first_score, 2)
                            if delta > 0:
                                direction = "⬆️ Improved"
                            elif delta < 0:
                                direction = "⬇️ Declined"
                            else:
                                direction = "➡️ No Change"
                            short_name = (
                                latest_dr.get("short") or
                                first_dr.get("short") or dk
                            )
                            movement_data.append({
                                "Domain": short_name,
                                "First Score": first_score,
                                "Latest Score": latest_score,
                                "Change": delta,
                                "Direction": direction
                            })
                        if movement_data:
                            st.dataframe(
                                movement_data,
                                use_container_width=True,
                                hide_index=True
                            )
                    elif len(all_snapshots) == 1:
                        st.info(
                            "Upload at least one previous snapshot to see "
                            "progress trends. Currently showing only the "
                            "live assessment."
                        )
                    else:
                        st.info(
                            "No snapshots loaded. Run an assessment and export "
                            "it, then upload past exports here to track progress."
                        )
            else:
                st.warning(
                    "No sufficient vocabulary found for this specific Maturity "
                    "Model. Try switching Personas or adding more text data."
                )

        st.divider()
        
        # advanced sections
        
        if enable_sentiment and beta_dist:
            st.subheader("⚖️ Bayesian Sentiment Inference")
            with st.expander("🧠 How to read this chart (and why it matters)", expanded=False):
                st.markdown("""
                **The Problem:** Standard sentiment analysis gives you a single number (e.g., "52% Positive"). But is that 52% based on 5 tweets or 5 million? A single number hides that uncertainty.
                
                **The Solution:** This chart calculates the **Probability** of the true sentiment.
                *   **The Curve (PDF):** Represents likelihood. The higher the peak, the more likely that specific sentiment score is the "truth."
                *   **The Shape:** 
                    *   **Narrow & Tall:** We have lots of data. We are highly confident the sentiment is exactly here.
                    *   **Wide & Flat:** We don't have enough data. The true sentiment could be almost anything.
                *   **The Green Zone (95% CI):** There is a 95% probability the "True" sentiment falls within this range. 
                
                **Decision Tip:** If the green zone is very wide (e.g., spanning 30% to 70%), **do not** make business decisions based on sentiment yet; you need more data.
                """)

            bayes_result = perform_bayesian_sentiment_analysis(combined_counts, term_sentiments, pos_threshold, neg_threshold)
            if bayes_result:
                b_col1, b_col2 = st.columns([1, 2])
                with b_col1:
                    st.metric(
                        "Positive Words Observed",
                        f"{bayes_result['pos_count']:,}",
                        help=ANALYSIS_HELP_TEXT["Positive Words Observed"],
                    )
                    st.metric(
                        "Negative Words Observed",
                        f"{bayes_result['neg_count']:,}",
                        help=ANALYSIS_HELP_TEXT["Negative Words Observed"],
                    )
                    st.info(f"Mean Expected Positive Rate: **{bayes_result['mean_prob']:.1%}**")
                    st.success(f"95% Credible Interval:\n**{bayes_result['ci_low']:.1%} — {bayes_result['ci_high']:.1%}**")
                with b_col2:
                    fig_bayes, ax_bayes = plt.subplots(figsize=(8, 4))
                    ax_bayes.plot(bayes_result['x_axis'], bayes_result['pdf_y'], lw=2, color='blue', label='Posterior PDF')
                    ax_bayes.fill_between(bayes_result['x_axis'], 0, bayes_result['pdf_y'], 
                                        where=(bayes_result['x_axis'] > bayes_result['ci_low']) & (bayes_result['x_axis'] < bayes_result['ci_high']),
                                        color='green', alpha=0.3, label='95% Credible Interval')
                    ax_bayes.set_title("Bayesian Update of Sentiment Confidence", fontsize=10)
                    ax_bayes.legend()
                    ax_bayes.grid(True, alpha=0.2)
                    st.pyplot(fig_bayes)
                    plt.close(fig_bayes)

        GRAPH_RENDER_NODE_LIMIT = 90
        GRAPH_RENDER_EDGE_LIMIT = 180

        show_graph = (
            proc_conf.compute_bigrams
            and scanner.global_bigrams
            and st.checkbox(
                "🕸️ Show Network Graph & Advanced Analytics",
                value=False,
                help=(
                    "Graph rendering can be heavy for dense corpora. Turn this on after "
                    "reviewing the dashboard, themes, and keyphrases."
                ),
            )
        )

        if show_graph:
            st.subheader("🔗 Network Graph")
            with st.expander("🛠️ Graph Settings & Physics", expanded=False):
                st.caption(
                    f"Rendering is capped at {GRAPH_RENDER_NODE_LIMIT} nodes and "
                    f"{GRAPH_RENDER_EDGE_LIMIT} links for browser stability on Streamlit Community Cloud."
                )
                c1, c2, c3 = st.columns(3)
                min_edge_weight = c1.slider(
                    "Min Link Frequency", 2, 100, 3,
                    help="Minimum shared occurrences required to draw a line. Increase this to remove weak connections and de-clutter the graph."
                )
                max_nodes_graph = c1.slider(
                    "Max Nodes", 10, GRAPH_RENDER_NODE_LIMIT, 55,
                    help="Hard limit on the number of words displayed. Lower this if the graph fails to render or looks too dense."
                )
                max_edges_graph = c1.slider(
                    "Max Links", 10, GRAPH_RENDER_EDGE_LIMIT, 90,
                    help="Hard limit on graph links. Lower this for safer rendering on Streamlit Community Cloud."
                )
                repulsion_val = c2.slider(
                    "Repulsion", 100, 3000, 1000, 
                    help="Physics Force: How strongly nodes push away from each other. Increase this if the graph looks like a tight ball."
                )
                edge_len_val = c2.slider(
                    "Edge Length", 50, 500, 250, 
                    help="Target length for the connecting lines. Increase this to space out distinct clusters."
                )
                physics_enabled = c3.checkbox(
                    "Enable Physics", True, 
                    help="If checked, the graph simulates gravity to organize itself. Uncheck to freeze the nodes in place."
                )
                directed_graph = c3.checkbox(
                    "Directed Arrows", False, 
                    help="Draws arrows (->) to show word order/flow, rather than just simple connections."
                )
                color_options = ["Community (Topic)", "Sentiment"]
                # Add maturity domain coloring if available
                if ('maturity_result' in locals() and maturity_result
                    and maturity_result.get("type") == "domain_based"):
                    color_options.append("Maturity Domain")
                color_mode = c3.radio(
                    "Color By:", color_options, index=0,
                    help="Community: Colors by structural cluster.\n"
                    "Sentiment: Colors by positive/negative.\n"
                    "Maturity Domain: Colors by which maturity domain the word belongs to (if 12-domain model is active)."
                )
                if (
                    max_nodes_graph >= GRAPH_RENDER_NODE_LIMIT
                    or max_edges_graph >= GRAPH_RENDER_EDGE_LIMIT
                ):
                    st.info(
                        "You are at the graph safety cap. If rendering is slow or blank, "
                        "lower Max Nodes or Max Links before rescanning the graph."
                    )

            G = nx.DiGraph() if directed_graph else nx.Graph()
            sorted_connections = []
            selected_nodes = set()

            for (src, tgt), weight in sorted(
                scanner.global_bigrams.items(),
                key=lambda item: item[1],
                reverse=True,
            ):
                if weight < min_edge_weight:
                    continue

                src = str(src).strip()
                tgt = str(tgt).strip()
                if not src or not tgt or src == tgt:
                    continue

                next_nodes = selected_nodes | {src, tgt}
                if len(next_nodes) > max_nodes_graph:
                    continue

                sorted_connections.append(((src, tgt), weight))
                selected_nodes = next_nodes

                if len(sorted_connections) >= max_edges_graph:
                    break

            if sorted_connections:
                G.add_edges_from(
                    (src, tgt, {"weight": weight})
                    for (src, tgt), weight in sorted_connections
                )
                try:
                    deg_centrality = nx.degree_centrality(G)
                except Exception:
                    deg_centrality = {node: 1 for node in G.nodes()}
                community_map = {}
                ai_cluster_info = ""
                
                if color_mode == "Community (Topic)":
                    G_undir = G.to_undirected() if directed_graph else G
                    try:
                        communities = nx_comm.greedy_modularity_communities(G_undir)
                        cluster_descriptions = []
                        for group_id, community in enumerate(communities):
                            top_in_cluster = sorted(list(community), key=lambda x: combined_counts[x], reverse=True)[:5]
                            cluster_descriptions.append(f"- Cluster {group_id+1}: {', '.join(top_in_cluster)}")
                            for node in community: community_map[node] = group_id
                        ai_cluster_info = "\n".join(cluster_descriptions)
                    except: pass

                community_colors = ["#FF4B4B", "#4589ff", "#ffa421", "#3cdb82", "#8b46ff", "#ff4b9f", "#00c0f2"]
                nodes, edges = [], []
                for node_id in G.nodes():
                    size = 15 + (deg_centrality.get(node_id, 0) * 80)
                    node_color = "#808080"
                    if color_mode == "Sentiment" and enable_sentiment:
                        s = term_sentiments.get(node_id, 0)
                        if s >= pos_threshold: node_color = pos_color
                        elif s <= neg_threshold: node_color = neg_color
                    elif color_mode == "Community (Topic)":
                        gid = community_map.get(node_id, 0)
                        node_color = community_colors[gid % len(community_colors)]
                    elif color_mode == "Maturity Domain" and maturity_result:
                        # Approximation: a word may appear in multiple domains, so the first mapped domain wins.
                        domain_colors_map = {
                            "01": "#e6194b", "02": "#3cb44b", "03": "#ffe119",
                            "04": "#4363d8", "05": "#f58231", "06": "#911eb4",
                            "07": "#42d4f4", "08": "#f032e6", "09": "#bfef45",
                            "10": "#fabed4", "11": "#469990", "12": "#dcbeff"
                        }
                        term_domain_map = {}
                        for dk, dd in maturity_result.get("domains_ref", {}).items():
                            prefix = dk[:2]
                            for tier_num, tier_data in dd["tiers"].items():
                                for term in tier_data.get("terms", set()):
                                    if term not in term_domain_map:
                                        term_domain_map[term] = prefix
                                for phrase in tier_data.get("phrases", []):
                                    for word in phrase.split():
                                        if word not in term_domain_map:
                                            term_domain_map[word] = prefix
                        matched_domain = term_domain_map.get(node_id)
                        if matched_domain:
                            node_color = domain_colors_map.get(matched_domain, "#808080")
                        else:
                            node_color = "#808080"

                    # re-added font config for white, legible text
                    nodes.append(Node(
                        id=node_id, 
                        label=node_id, 
                        size=size, 
                        color=node_color,
                        font={'color': 'white', 'size': 20, 'strokeWidth': 2, 'strokeColor': '#000000'}
                    ))

                for (source, target), weight in sorted_connections:
                    width = 1 + math.log(weight) * 0.8
                    edges.append(Edge(source=source, target=target, width=width, color="#e0e0e0"))
                
                # re-added interaction dict for zoom/pan buttons
                config = Config(
                    width=1000, 
                    height=700, 
                    directed=directed_graph, 
                    physics=physics_enabled, 
                    hierarchy=False, 
                    interaction={"navigationButtons": True, "zoomView": True}, 
                    physicsSettings={"solver": "forceAtlas2Based", "forceAtlas2Based": {"gravitationalConstant": -abs(repulsion_val), "springLength": edge_len_val, "springConstant": 0.05, "damping": 0.4}}
                )
                
                if len(nodes) > GRAPH_RENDER_NODE_LIMIT or len(edges) > GRAPH_RENDER_EDGE_LIMIT:
                    st.warning(
                        "Graph is too dense to render safely in-browser. "
                        "Lower Max Nodes / Max Links, or download the GEXF file for external graph tools."
                    )
                elif not nodes or not edges:
                    st.info("Not enough graph structure to render with the current settings.")
                else:
                    st.info("💡 **Navigation Tip:** Use the buttons in the **bottom-right** of the graph to Zoom & Pan.")
                    try:
                        agraph(nodes=nodes, edges=edges, config=config)
                    except Exception as e:
                        st.warning(
                            "The interactive graph could not render safely. "
                            "Try lowering Max Nodes, Max Links, or disabling Physics."
                        )

                # [NEW] Gephi Export
                if st.button("📥 Download Graph File (.gexf)"):
                    try:
                        # Write to memory buffer
                        gexf_buffer = io.BytesIO()
                        nx.write_gexf(G, gexf_buffer)
                        gexf_buffer.seek(0)
                        st.download_button("Click to Save GEXF", gexf_buffer, "network.gexf", "application/xml")
                    except Exception as e:
                        st.error(f"Export failed: {e}")
                
                # graph analytics tabs
                tab_g1, tab_g2, tab_g3, tab_g4 = st.tabs(["Basic Stats", "Top Nodes", "Text Stats", "🔥 Heatmap"])
                with tab_g1:
                    col_b1, col_b2, col_b3 = st.columns(3)
                    col_b1.metric(
                        "Nodes", 
                        G.number_of_nodes(),
                        help="The number of unique concepts (dots) currently displayed in the network."
                    )
                    col_b2.metric(
                        "Edges", 
                        G.number_of_edges(),
                        help="The number of connections (lines) between concepts. A connection indicates these words appear together frequently."
                    )
                    try: 
                        col_b3.metric(
                            "Density", 
                            f"{nx.density(G):.4f}",
                            help="A score from 0 to 1 indicating how interconnected the network is.\n\n• Low (<0.1): Distinct, separate topics.\n• High (>0.5): Everything is related to everything (a 'hairball')."
                        )
                    except: pass
                with tab_g2:
                    node_weights = {n: 0 for n in G.nodes()}
                    for u, v, data in G.edges(data=True):
                        w = data.get('weight', 1)
                        node_weights[u] += w
                        node_weights[v] += w
                    st.dataframe(pd.DataFrame(list(node_weights.items()), columns=["Node", "Weighted Degree"]).sort_values("Weighted Degree", ascending=False).head(50), use_container_width=True)
                with tab_g3:
                     c1, c2, c3 = st.columns(3)
                     c1.metric(
                        "Total Tokens", 
                        f"{text_stats['Total Tokens']:,}",
                        help="The total count of all words processed after cleaning (removing stopwords, numbers, etc.). Represents the sheer volume of signal."
                     )
                     c2.metric(
                        "Unique Vocab", 
                        f"{text_stats['Unique Vocabulary']:,}",
                        help="The count of distinct, unique words found. A higher number indicates a broader range of topics or more complex language."
                     )
                     c3.metric(
                        "Lexical Diversity", 
                        f"{text_stats['Lexical Diversity']}",
                        help="The Ratio of Unique Words to Total Words (Unique / Total). \n\n• High (>0.5): Dense information, varied vocabulary (e.g., Poetry, Abstracts).\n• Low (<0.1): Highly repetitive, consistent language (e.g., Logs, Legal Boilerplate)."
                     )
                with tab_g4:
                    # hybrid heatmap-QR generator
                    viz_mode = st.radio("Visualization Mode", ["Standard Heatmap", "Hybrid Signature (Scanable)"], horizontal=True, label_visibility="collapsed")
                    
                    top_20 = [w for w, c in combined_counts.most_common(20)]
                    
                    if len(top_20) > 1:
                        # generating matrix data
                        mat = np.zeros((len(top_20), len(top_20)))
                        for i, w1 in enumerate(top_20):
                            for j, w2 in enumerate(top_20):
                                if i != j: mat[i][j] = scanner.global_bigrams.get((w1, w2), 0) + scanner.global_bigrams.get((w2, w1), 0)
                        
                        # to plot heatmap to a PIL image (memory buffer)
                        fig_h, ax_h = plt.subplots(figsize=(10, 10)) 
                        ax_h.imshow(mat, cmap=colormap, interpolation='nearest') 
                        
                        if viz_mode == "Hybrid Signature (Scanable)":
                            ax_h.axis('off')
                            plt.tight_layout(pad=0)
                        else:
                            ax_h.set_xticks(np.arange(len(top_20)))
                            ax_h.set_yticks(np.arange(len(top_20)))
                            ax_h.set_xticklabels(top_20, rotation=45, ha="right")
                            ax_h.set_yticklabels(top_20)
                        
                        buf = BytesIO()
                        fig_h.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
                        buf.seek(0)
                        
                        if viz_mode == "Standard Heatmap":
                            st.pyplot(fig_h)
                        
                        elif viz_mode == "Hybrid Signature (Scanable)":
                            if qrcode is None:
                                st.error("🚨 Please install: `pip install qrcode[pil]`")
                            else:
                                from PIL import Image, ImageEnhance
                                
                                # to allow for custom "steganographic" payload (with char limit)
                                custom_msg = st.text_input(
                                    "🔒 Encode Custom Payload (Optional)", 
                                    placeholder="Leave empty for standard metadata...",
                                    max_chars=1000,
                                    help="Limit: ~1000 characters to ensure QR readability with high error correction."
                                )
                                
                                # prepping heatmap
                                heatmap_img = Image.open(buf).convert("RGBA")
                                enhancer = ImageEnhance.Brightness(heatmap_img)
                                heatmap_img = enhancer.enhance(1.5) 
                                
                                # defining payload
                                if custom_msg.strip():
                                    signature_payload = custom_msg
                                else:
                                    signature_payload = (
                                        f"SIGNAL FOUNDRY\nRef: {st.session_state.get('last_sketch_hash', 'SESSION')}\n"
                                        f"Top: {', '.join(top_20[:3])}"
                                    )
                                
                                # to generate QR with safety check
                                try:
                                    qr = qrcode.QRCode(error_correction=qrcode.constants.ERROR_CORRECT_H, border=1)
                                    qr.add_data(signature_payload)
                                    qr.make(fit=True)
                                    qr_img = qr.make_image(fill_color="black", back_color="transparent").convert("RGBA")
                                    
                                    # composite
                                    heatmap_resized = heatmap_img.resize(qr_img.size)
                                    hybrid_img = Image.alpha_composite(heatmap_resized, qr_img)
                                    
                                    c1, c2 = st.columns([2, 1])
                                    with c1:
                                        st.image(hybrid_img, caption="Scan to verify data (or read payload).", use_container_width=True)
                                    with c2:
                                        st.markdown("### 🧬 Hybrid Signature")
                                        if custom_msg.strip():
                                            st.warning("⚠️ **Mode:** Custom Payload Active")
                                        else:
                                            st.info("ℹ️ **Mode:** Standard Metadata")
                                        
                                        st.caption("The colors represent the data relationships (Heatmap). The dark overlay pattern encodes the message (QR).")
                                        
                                        final_buf = BytesIO()
                                        hybrid_img.save(final_buf, format="PNG")
                                        st.download_button("📥 Download Signature", final_buf.getvalue(), "heatmap_signature.png", "image/png")
                                        
                                except Exception as e:
                                    st.error("❌ **Payload Too Large:** The text is too long to fit into a High-Security QR code. Please shorten your message.")
                    else:
                        st.info("Not enough data to generate signature.")
                    #

        st.subheader("🔍 Bayesian Theme Discovery")
        if len(scanner.topic_docs) > 5 and DictVectorizer:
            with st.spinner(f"Running {topic_model_type} Topic Modeling..."):
                topics = perform_topic_modeling(scanner.topic_docs, n_topics, topic_model_type)
            if topics:
                cols = st.columns(len(topics))
                for idx, topic in enumerate(topics):
                    with cols[idx]:
                        st.markdown(f"**Topic {topic['id']}**")
                        for w in topic['words']: st.markdown(f"`{w}`")
        else:
            st.info("Needs more data/docs to model topics.")

        # detailed frequency tables
        st.divider()
        st.subheader(f"📊 Frequency Tables (Top {top_n})")
        most_common = combined_counts.most_common(top_n)
        data = []
        if enable_sentiment:
            for w, f in most_common:
                score = term_sentiments.get(w, 0.0)
                category = get_sentiment_category(score, pos_threshold, neg_threshold)
                data.append([w, f, score, category])
        else:
            data = [[w, f] for w, f in most_common]

        cols = ["word", "count"] + (["sentiment", "category"] if enable_sentiment else [])
        st.dataframe(pd.DataFrame(data, columns=cols), use_container_width=True)
        
        if proc_conf.compute_bigrams and scanner.global_bigrams:
            st.write("Bigrams (By Frequency)")
            top_bg = scanner.global_bigrams.most_common(top_n)
            bg_data = []
            if enable_sentiment:
                for bg_tuple, f in top_bg:
                    bg_str = " ".join(bg_tuple)
                    score = term_sentiments.get(bg_str, 0.0)
                    category = get_sentiment_category(score, pos_threshold, neg_threshold)
                    bg_data.append([bg_str, f, score, category])
            else:
                bg_data = [[" ".join(bg), f] for bg, f in top_bg]
            bg_cols = ["bigram", "count"] + (["sentiment", "category"] if enable_sentiment else [])
            st.dataframe(pd.DataFrame(bg_data, columns=bg_cols), use_container_width=True)

            # NPMI in expander (original style)
            with st.expander("🔬 Phrase Significance (NPMI Score)", expanded=False):
                st.markdown("""
                **NPMI (Normalized Pointwise Mutual Information)** finds words that *belong* together, rather than just words that appear often.
                *   High Score (> 0.5): Strong association (e.g., "Artificial Intelligence").
                *   Low Score (< 0.1): Random association (e.g., "of the").
                """)
                df_npmi = calculate_npmi(scanner.global_bigrams, combined_counts, scanner.total_rows_processed)
                st.dataframe(df_npmi.head(top_n), use_container_width=True)

    # --- AI analyst (restored full mode)
    if combined_counts and st.session_state['authenticated']:
        st.divider()
        st.subheader("🤖 AI Analyst")
        st.caption(
            "**Privacy Note:** This AI does *not* read your raw documents. "
            "It interprets the statistical 'Sketch' (frequencies, connections, and clusters) generated by this engine. "
            "This allows for high-level pattern recognition without exposing the full text content."
        )
        
        ai_ctx_str = build_ai_context_brief(
            scanner=scanner,
            combined_counts=combined_counts,
            insight_df=insight_df if 'insight_df' in locals() and isinstance(insight_df, pd.DataFrame) else pd.DataFrame(),
            text_stats=text_stats if 'text_stats' in locals() else {},
            maturity_result=maturity_result if 'maturity_result' in locals() else None,
            graph_cluster_info=locals().get('ai_cluster_info', 'N/A'),
            max_items=25,
        )

        with st.expander("🔎 What the AI can see", expanded=False):
            st.caption(
                "This is the privacy-preserving briefing sent to the AI provider. "
                "It includes summaries, counts, signal labels, and maturity summaries when available, not full raw documents."
            )
            st.text_area(
                "AI context preview",
                ai_ctx_str,
                height=320,
                disabled=True,
                help="Use this to understand what the AI Analyst can and cannot use when answering.",
            )
        
        col_ai_1, col_ai_2 = st.columns(2)
        
        with col_ai_1:
            st.markdown("**1. One-Click Theme Detection**")
            if st.button("✨ Identify Key Themes", type="primary"):
                with st.status("Analyzing..."):
                    system_prompt = (
                        "You are a qualitative data analyst using a privacy-preserving analytical sketch. "
                        "Analyze the Signal Compass, Resource Shape, signal roles, top terms, entities, "
                        "Insight Cards, and maturity summaries if provided. Identify the main shape of the "
                        "resource, 3-5 key themes, notable risks/tensions/anomalies, and the most useful "
                        "human follow-up questions. Do not claim to have read raw source documents. "
                        "Distinguish evidence from interpretation."
                    )
                    user_prompt = f"Data Context:\n{ai_ctx_str}"
                    response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                    st.session_state["ai_response"] = response
                    st.rerun()

        with col_ai_2:
            st.markdown("**2. Ask the Data**")
            user_question = st.text_area("Ask a specific question:", height=100, placeholder="e.g., 'What are the main complaints about pricing?'")
            if st.button("Ask Question"):
                if user_question.strip():
                    with st.status("Thinking..."):
                        system_prompt = (
                            "You are an expert analyst. Answer the user's question based ONLY on the provided "
                            "Signal Foundry context brief: summary statistics, Signal Compass, Resource Shape, "
                            "associations, Insight Cards, entities, and maturity summaries when present. "
                            "If you cannot answer from that context, say so. Do not imply access to raw source "
                            "documents. Distinguish evidence from interpretation."
                        )
                        user_prompt = f"Data Context:\n{ai_ctx_str}\n\nUser Question: {user_question}"
                        response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                        st.session_state["ai_response"] = f"**Q: {user_question}**\n\n{response}"
                        st.rerun()
                else:
                    st.warning("Please enter a question.")

        if st.session_state["ai_response"]:
            st.divider()
            st.markdown("### 📋 AI Output")
            st.markdown(st.session_state["ai_response"])

st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #808080; font-size: 12px;'>"
    "Open Source software licensed under the MIT License."
    "</div>", 
    unsafe_allow_html=True
)
